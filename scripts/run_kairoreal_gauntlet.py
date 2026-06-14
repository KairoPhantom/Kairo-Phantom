#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
run_kairoreal_gauntlet.py — KairoReal 200-scenario headless gauntlet.

Loads all 200 scenarios from scenarios.json and runs the 50 'active' ones
through category-specific headless executors (no UI automation, no network).
Produces task_completion_rate.json and exits 0 if pass_rate_active >= 80%.

Usage:
    python scripts/run_kairoreal_gauntlet.py [--workers N] [--output PATH]
"""
from __future__ import annotations

import argparse
import concurrent.futures
import json
import logging
import os
import shutil
import sys
import tempfile
import time
import traceback
from datetime import datetime, timezone
from typing import Any, Dict, List, Optional

# ── Path setup ────────────────────────────────────────────────────────────────
_SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
_REPO_ROOT = os.path.dirname(_SCRIPT_DIR)
_SIDECAR_ROOT = os.path.join(_REPO_ROOT, "kairo-sidecar")
for _p in (_SIDECAR_ROOT, os.path.join(_SIDECAR_ROOT, "sidecar")):
    if _p not in sys.path:
        sys.path.insert(0, _p)

# ── Logging ───────────────────────────────────────────────────────────────────
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[logging.StreamHandler(stream=sys.stdout)],
)
log = logging.getLogger("KairoRealGauntlet")

# ── Constants ─────────────────────────────────────────────────────────────────
SCENARIOS_JSON = os.path.join(_REPO_ROOT, "scenarios.json")
GATE_THRESHOLD = 80.0
GAUNTLET_VERSION = "kairoreal-headless-v1"


# ── Verdict helpers ───────────────────────────────────────────────────────────

def _pass(reason: str) -> Dict[str, str]:
    return {"oracle_verdict": "PASS", "reason": reason}


def _fail(reason: str) -> Dict[str, str]:
    return {"oracle_verdict": "FAIL", "reason": reason}


def _skip(reason: str) -> Dict[str, str]:
    return {"oracle_verdict": "SKIP", "reason": reason}


# ── Category executors ────────────────────────────────────────────────────────

def _exec_word(sandbox_path: str, scenario: Dict[str, Any]) -> Dict[str, str]:
    """Word domain: run real WordMaster and WordWriter pipeline with falsifiable oracle."""
    try:
        import docx
        from sidecar.masters.word_master import WordContextExtractor, WordWriter
    except ImportError as exc:
        return _skip(f"Word dependencies missing: {exc}")

    try:
        expected = scenario.get("expected_outcome", {})
        action = expected.get("action")
        if not action:
            return _fail("Missing expected action in expected_outcome")

        # Create a document with some baseline content
        path = os.path.join(sandbox_path, f"scenario_{scenario['id']}.docx")
        doc = docx.Document()
        doc.add_heading("Heading 0", level=1)
        doc.add_paragraph("Paragraph 0 text.")
        doc.add_paragraph("Paragraph 1 text.")
        doc.add_paragraph("Paragraph 2 text.")
        doc.save(path)

        orig_texts = ["Heading 0", "Paragraph 0 text.", "Paragraph 1 text.", "Paragraph 2 text."]

        # Extract context
        extractor = WordContextExtractor()
        ctx = extractor.extract(path, cursor_paragraph_index=0)

        # Apply operations
        writer = WordWriter()
        op = {
            "type": action,
            "style": expected.get("style", "Normal")
        }
        if "after_paragraph_index" in expected:
            op["after_paragraph_index"] = expected["after_paragraph_index"]
        if "paragraph_index" in expected:
            op["paragraph_index"] = expected["paragraph_index"]
        if "text" in expected:
            op["runs"] = [{"text": expected["text"]}]
        if "rows" in expected:
            op["rows"] = [[""] * expected["cols"] for _ in range(expected["rows"])]
            op["cols"] = expected["cols"]
            # Align after_paragraph_index
            if "paragraph_index" in expected:
                op["after_paragraph_index"] = expected["paragraph_index"]

        res = writer.apply_operations(path, [op], ctx)
        if "error" in res:
            return _fail(f"Writer error: {res['error']}")

        # Real oracle: load the modified document and verify changes
        doc2 = docx.Document(path)
        paras = doc2.paragraphs
        tables = doc2.tables
        actual_texts = [p.text for p in paras]

        if action == "insert_paragraph":
            target_text = expected["text"]
            idx = expected["after_paragraph_index"]
            expected_texts = orig_texts.copy()
            if idx == -1:
                expected_texts.append(target_text)
            else:
                expected_texts.insert(idx + 1, target_text)
            if actual_texts != expected_texts:
                return _fail(f"Expected paragraphs {expected_texts}, got {actual_texts}")

        elif action == "replace_paragraph":
            target_text = expected["text"]
            idx = expected["paragraph_index"]
            expected_texts = orig_texts.copy()
            expected_texts[idx] = target_text
            if actual_texts != expected_texts:
                return _fail(f"Expected paragraphs {expected_texts}, got {actual_texts}")

        elif action == "delete_paragraph":
            idx = expected["paragraph_index"]
            expected_texts = [t for i, t in enumerate(orig_texts) if i != idx]
            if actual_texts != expected_texts:
                return _fail(f"Expected paragraphs {expected_texts}, got {actual_texts}")

        elif action == "append_to_run":
            target_text = expected["text"]
            idx = expected["paragraph_index"]
            expected_texts = orig_texts.copy()
            expected_texts[idx] = expected_texts[idx] + target_text
            if actual_texts != expected_texts:
                return _fail(f"Expected paragraphs {expected_texts}, got {actual_texts}")

        elif action == "insert_table":
            if not tables:
                return _fail("No tables found in modified document")
            t = tables[0]
            if len(t.rows) != expected["rows"] or len(t.columns) != expected["cols"]:
                return _fail(f"Expected table {expected['rows']}x{expected['cols']}, got {len(t.rows)}x{len(t.columns)}")
            if actual_texts != orig_texts:
                return _fail(f"Expected paragraphs to remain {orig_texts}, got {actual_texts}")

        return _pass(f"Word oracle passed for action: {action}")
    except Exception as exc:
        return _fail(f"Word executor error: {exc}")


def _exec_excel(sandbox_path: str, scenario: Dict[str, Any]) -> Dict[str, str]:
    """Excel domain: run real ExcelMaster/Writer or forge formula validator with falsifiable oracle."""
    try:
        import openpyxl
        from sidecar.parsers.forge_bridge import validate_formula, explain_formula
        from sidecar.masters.excel_master import ExcelContextExtractor, ExcelWriter
    except ImportError as exc:
        return _skip(f"Excel dependencies missing: {exc}")

    try:
        expected = scenario.get("expected_outcome", {})
        action = expected.get("action")
        if not action:
            return _fail("Missing expected action in expected_outcome")

        if action == "validate_formula":
            res = validate_formula(expected["formula"])
            valid = res.get("valid", False) and not res.get("fix_applied")
            if valid != expected["expected_valid"]:
                return _fail(f"Expected formula validity {expected['expected_valid']} for {expected['formula']}, got {valid}")
            return _pass(f"Formula validation oracle passed: {expected['formula']} is {valid}")

        elif action == "explain_formula":
            res = explain_formula(expected["formula"])
            explanation = str(res or "").lower()
            if expected["contains"] not in explanation:
                return _fail(f"Expected explanation for {expected['formula']} to contain '{expected['contains']}', got '{explanation}'")
            return _pass(f"Formula explanation oracle passed: {expected['formula']} explained correctly")

        elif action in ("write_cell", "write_range"):
            path = os.path.join(sandbox_path, f"scenario_{scenario['id']}.xlsx")
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = "Sheet"
            # Add some dummy cells
            ws["A1"] = 10
            ws["A2"] = 20
            ws["B1"] = 30
            ws["B2"] = 40
            wb.save(path)

            writer = ExcelWriter()
            op = {
                "type": action,
                "sheet": "Sheet"
            }
            if "cell" in expected:
                op["cell"] = expected["cell"]
            if "range" in expected:
                op["range"] = expected["range"]
            if "value" in expected:
                op["value"] = expected["value"]
            if "formula" in expected:
                op["formula"] = expected["formula"]
            if "values" in expected:
                op["values"] = expected["values"]
            if "formulas" in expected:
                op["formulas"] = expected["formulas"]

            res = writer.apply_operations(path, [op])
            if "errors" in res and res["errors"]:
                return _fail(f"Writer errors: {res['errors']}")
            if "error" in res:
                return _fail(f"Writer error: {res['error']}")

            # Load and verify
            wb2 = openpyxl.load_workbook(path, data_only=False)
            ws2 = wb2["Sheet"]

            if action == "write_cell":
                cell_ref = expected["cell"]
                cell_obj = ws2[cell_ref]
                if "formula" in expected:
                    if cell_obj.value != expected["formula"]:
                        return _fail(f"Expected formula '{expected['formula']}' in cell {cell_ref}, got '{cell_obj.value}'")
                elif "value" in expected:
                    if cell_obj.value != expected["value"]:
                        return _fail(f"Expected value '{expected['value']}' in cell {cell_ref}, got '{cell_obj.value}'")

            elif action == "write_range":
                range_ref = expected["range"]
                if ":" not in range_ref:
                    # Single cell reference in range
                    cell_obj = ws2[range_ref]
                    if "values" in expected:
                        expected_val = expected["values"][0][0]
                        if cell_obj.value != expected_val:
                            return _fail(f"Expected '{expected_val}' at cell {range_ref}, got '{cell_obj.value}'")
                    elif "formulas" in expected:
                        expected_form = expected["formulas"][0][0]
                        if cell_obj.value != expected_form:
                            return _fail(f"Expected formula '{expected_form}' at cell {range_ref}, got '{cell_obj.value}'")
                else:
                    if "values" in expected:
                        expected_vals = expected["values"]
                        rows = list(ws2[range_ref])
                        for r_idx, row in enumerate(rows):
                            for c_idx, cell_obj in enumerate(row):
                                expected_val = expected_vals[r_idx][c_idx]
                                if cell_obj.value != expected_val:
                                    return _fail(f"Expected '{expected_val}' at cell {cell_obj.coordinate}, got '{cell_obj.value}'")
                    elif "formulas" in expected:
                        expected_forms = expected["formulas"]
                        rows = list(ws2[range_ref])
                        for r_idx, row in enumerate(rows):
                            for c_idx, cell_obj in enumerate(row):
                                expected_form = expected_forms[r_idx][c_idx]
                                if cell_obj.value != expected_form:
                                    return _fail(f"Expected formula '{expected_form}' at cell {cell_obj.coordinate}, got '{cell_obj.value}'")

            return _pass(f"Excel oracle passed for action: {action}")
    except Exception as exc:
        return _fail(f"Excel executor error: {exc}")


def _exec_ppt(sandbox_path: str, scenario: Dict[str, Any]) -> Dict[str, str]:
    """PPT domain: run PowerPointMaster/pptx context extraction with falsifiable oracle."""
    try:
        import pptx
        from sidecar.masters.other_masters import PowerPointMaster
    except ImportError as exc:
        return _skip(f"PPT dependencies missing: {exc}")

    try:
        expected = scenario.get("expected_outcome", {})
        action = expected.get("action")
        if not action:
            return _fail("Missing expected action in expected_outcome")

        path = os.path.join(sandbox_path, f"scenario_{scenario['id']}.pptx")

        if action == "create_presentation":
            prs = pptx.Presentation()
            for _ in range(expected["slides"]):
                prs.slides.add_slide(prs.slide_layouts[5])
            prs.save(path)

            prs2 = pptx.Presentation(path)
            if len(prs2.slides) != expected["slides"]:
                return _fail(f"Expected {expected['slides']} slides, got {len(prs2.slides)}")
            return _pass(f"PowerPoint creation oracle passed: created {expected['slides']} slides")

        elif action == "extract_context":
            prs = pptx.Presentation()
            for _ in range(expected["expected_slide_count"]):
                prs.slides.add_slide(prs.slide_layouts[5])
            prs.save(path)

            master = PowerPointMaster()
            ctx = master.extract_context(path, expected["cursor_slide"])
            if ctx.get("total_slides") != expected["expected_slide_count"]:
                return _fail(f"Expected total_slides={expected['expected_slide_count']} in context, got {ctx.get('total_slides')}")
            return _pass(f"PowerPoint extract_context oracle passed: total_slides matches")

        elif action == "slide_layout":
            prs = pptx.Presentation()
            layouts_map = {
                "Title Slide": prs.slide_layouts[0],
                "Title and Content": prs.slide_layouts[1],
                "Blank": prs.slide_layouts[6],
                "Section Header": prs.slide_layouts[2]
            }
            layout_name = expected["expected_layout"]
            layout = layouts_map.get(layout_name, prs.slide_layouts[5])
            
            for idx in range(expected["cursor_slide"] + 1):
                if idx == expected["cursor_slide"]:
                    prs.slides.add_slide(layout)
                else:
                    prs.slides.add_slide(prs.slide_layouts[5])
            prs.save(path)

            master = PowerPointMaster()
            ctx = master.extract_context(path, expected["cursor_slide"])
            actual_layout = ctx.get("layout_name", "")
            if expected["expected_layout"] not in actual_layout:
                return _fail(f"Expected layout '{expected['expected_layout']}', got '{actual_layout}'")
            return _pass(f"PowerPoint layout oracle passed: layout is '{actual_layout}'")

        return _fail(f"Unknown PPT action: {action}")
    except Exception as exc:
        return _fail(f"PPT executor error: {exc}")


def _exec_legal(sandbox_path: str, scenario: Dict[str, Any]) -> Dict[str, str]:
    """Legal domain: run CUAD clause detection on prompt with falsifiable oracle."""
    try:
        from sidecar.parsers.legal_redline import detect_cuad_clauses
    except ImportError as exc:
        return _skip(f"Legal dependencies missing: {exc}")

    try:
        expected = scenario.get("expected_outcome", {})
        cid = expected.get("clause_id")
        if not cid:
            return _fail("Missing expected clause_id in expected_outcome")

        prompt_text = scenario.get("prompt", "")
        res = detect_cuad_clauses(prompt_text)
        clauses = res.get("data", {}).get("detected_clauses", [])
        detected_ids = [clause.get("id") for clause in clauses if isinstance(clause, dict)]

        if cid not in detected_ids:
            return _fail(f"Expected to detect clause '{cid}', but detected only {detected_ids}")
        return _pass(f"Legal oracle passed: detected '{cid}' in prompt")
    except Exception as exc:
        return _fail(f"Legal executor error: {exc}")


def _exec_cua(sandbox_path: str, scenario: Dict[str, Any]) -> Dict[str, str]:
    """CUA domain: verify safety and IPC limits and fallback behaviors."""
    try:
        from sidecar.cua.canva_cua import CanvaCUAAgent, SAFETY_LIMITS
        from sidecar.ipc import MAX_MESSAGE_BYTES, MAX_CONCURRENT_REQUESTS
        from sidecar.cua.driver_service import CuaDriverService
    except ImportError as exc:
        return _skip(f"CUA dependencies missing: {exc}")

    try:
        expected = scenario.get("expected_outcome", {})
        group = expected.get("group")
        field = expected.get("field")
        expected_val = expected.get("expected")

        field_clean = field
        for suffix in ("_value", "_verify", "_fallback"):
            if field_clean.endswith(suffix):
                field_clean = field_clean[:-len(suffix)]

        if group == "SAFETY_LIMITS":
            actual = SAFETY_LIMITS.get(field_clean)
            if actual != expected_val:
                return _fail(f"Expected SAFETY_LIMITS[{field_clean}] == {expected_val}, got {actual}")
            return _pass(f"CUA Safety Limit oracle passed: {field_clean} matches")

        elif group == "IPC_LIMITS":
            if field_clean == "MAX_MESSAGE_BYTES":
                if MAX_MESSAGE_BYTES != expected_val:
                    return _fail(f"Expected MAX_MESSAGE_BYTES == {expected_val}, got {MAX_MESSAGE_BYTES}")
            elif field_clean == "MAX_CONCURRENT_REQUESTS":
                if MAX_CONCURRENT_REQUESTS != expected_val:
                    return _fail(f"Expected MAX_CONCURRENT_REQUESTS == {expected_val}, got {MAX_CONCURRENT_REQUESTS}")
            return _pass(f"IPC Limit oracle passed: {field_clean} matches")

        elif group == "CLIPBOARD_FALLBACK":
            agent = CanvaCUAAgent()
            res = agent.execute_text_replacement(expected_val)
            if hasattr(res, "clipboard_content") and res.clipboard_content != expected_val:
                return _fail(f"Expected clipboard content '{expected_val}', got '{res.clipboard_content}'")
            return _pass("CUA Clipboard Fallback oracle passed")

        elif group == "DRIVER_STATUS":
            service = CuaDriverService()
            if field_clean == "screenshot":
                res = service.screenshot()
                if res is not None:
                    return _fail("Expected screenshot to return None on headless fallback")
            elif field_clean == "click":
                res = service.click(0, 0)
                if res is not False:
                    return _fail("Expected click to return False on headless fallback")
            elif field_clean == "type":
                res = service.type_text("test")
                if res is not False:
                    return _fail("Expected type_text to return False on headless fallback")
            return _pass(f"CUA Driver Status oracle passed: {field_clean} behaves correctly")

        return _fail(f"Unknown CUA group: {group}")
    except Exception as exc:
        return _fail(f"CUA executor error: {exc}")


def _exec_security(sandbox_path: str, scenario: Dict[str, Any]) -> Dict[str, str]:
    """Security domain: verify integrity guard rejects known violations."""
    import subprocess
    guard_path = os.path.join(_REPO_ROOT, "scripts", "ci", "eval_integrity_guard.py")
    if not os.path.exists(guard_path):
        return _skip("eval_integrity_guard.py not found")

    try:
        expected = scenario.get("expected_outcome", {})
        target = expected.get("target")
        should_reject = expected.get("expected_rejected")

        if os.path.exists(os.path.join(_REPO_ROOT, target)):
            file_to_scan = os.path.join(_REPO_ROOT, target)
        else:
            file_to_scan = os.path.join(sandbox_path, "scan_target.py")
            content = ""
            if target == "secrets_password":
                content = "password = 'hardcoded_secret_123'"
            elif target == "secrets_eval":
                content = "eval('print(1)')"
            elif target == "secrets_shell":
                content = "import os; os.system('ls')"
            elif target == "secrets_traversal":
                content = "with open('../../../etc/passwd') as f: pass"
            elif target == "secrets_apikey":
                content = "API_KEY = 'SG.1234567890'"
            elif target == "secrets_exec":
                content = "exec('x = 5')"
            elif target == "secrets_os_system":
                content = "import os; os.system('rm -rf /')"
            elif target == "secrets_path_traversal":
                content = "open('/etc/shadow', 'w')"
            elif target == "secrets_token":
                content = "GITHUB_TOKEN = 'ghp_abcdef'"
            elif target == "secrets_dynamic_import":
                content = "__import__('os').system('ls')"
            elif target == "secrets_popen":
                content = "import subprocess; subprocess.Popen('ls')"
            elif target == "secrets_abs_path":
                content = "open('/absolute/path/file.txt')"
            else:
                content = "import math; print(math.sqrt(4))"

            with open(file_to_scan, "w", encoding="utf-8") as fh:
                fh.write(content)

        result = subprocess.run(
            [sys.executable, guard_path, "--paths", file_to_scan],
            capture_output=True, text=True, timeout=10,
        )
        is_rejected = result.returncode != 0

        if is_rejected != should_reject:
            return _fail(f"Integrity guard check for {target}: expected_rejected={should_reject}, got rejected={is_rejected} (stdout: {result.stdout.strip()})")
        return _pass(f"Security integrity oracle passed for {target} (rejected={is_rejected})")
    except Exception as exc:
        return _fail(f"Security executor error: {exc}")


def _exec_memory(sandbox_path: str, scenario: Dict[str, Any]) -> Dict[str, str]:
    """Memory domain: verify DP noise meets CSPRNG expectations."""
    try:
        from sidecar.mem_sync import add_gaussian_noise
    except ImportError as exc:
        return _skip(f"Memory dependencies missing: {exc}")

    try:
        expected = scenario.get("expected_outcome", {})
        vector = expected.get("vector")
        std_dev = expected.get("std_dev")

        noisy = add_gaussian_noise(vector, std_dev=std_dev)
        if not isinstance(noisy, list):
            return _fail(f"Expected list, got {type(noisy)}")
        if len(noisy) != len(vector):
            return _fail(f"Length mismatch: {len(noisy)} != {len(vector)}")
        if noisy == vector:
            return _fail("add_gaussian_noise returned identical vector")
        for val in noisy:
            if not isinstance(val, (int, float)):
                return _fail(f"Non-numeric noisy value: {val}")
        return _pass("Memory DP noise oracle passed: noisy vector verified")
    except Exception as exc:
        return _fail(f"Memory executor error: {exc}")


import threading
_offline_env_lock = threading.Lock()


def _exec_offline(sandbox_path: str, scenario: Dict[str, Any]) -> Dict[str, str]:
    """Offline domain: verify sidecar self check responds offline_mode=True under KAIRO_OFFLINE=1."""
    import asyncio
    with _offline_env_lock:
        old_val = os.environ.get("KAIRO_OFFLINE")
        try:
            os.environ["KAIRO_OFFLINE"] = "1"
            import importlib
            import sidecar.main as km
            importlib.reload(km)
            
            result = asyncio.run(km.handle_request({"action": "self_check"}))
            if not isinstance(result, dict) or not result.get("ok"):
                return _fail("Self check failed under KAIRO_OFFLINE=1")
            
            data = result.get("data", {})
            if data.get("offline_mode") is not True:
                return _fail("Expected offline_mode=True in self check data")
            
            return _pass("Offline mode self check oracle passed")
        except Exception as exc:
            return _fail(f"Offline executor error: {exc}")
        finally:
            if old_val is None:
                os.environ.pop("KAIRO_OFFLINE", None)
            else:
                os.environ["KAIRO_OFFLINE"] = old_val


def _exec_degradation(sandbox_path: str, scenario: Dict[str, Any]) -> Dict[str, str]:
    """Degradation domain: verify graceful rejection of nonexistent domains/actions."""
    import asyncio
    try:
        import sidecar.main as km
        expected = scenario.get("expected_outcome", {})
        target = expected.get("target")
        type_ = expected.get("type")

        if type_ == "invalid_domain":
            result = asyncio.run(km.handle_request({
                "action": "apply_operations",
                "domain": target,
                "operations": []
            }))
        elif type_ == "invalid_action":
            result = asyncio.run(km.handle_request({
                "action": target,
                "domain": "Word",
                "payload": {}
            }))
        else:
            return _fail(f"Unknown degradation type: {type_}")

        if not isinstance(result, dict):
            return _fail("Expected dict response, got %s" % type(result))
        if result.get("ok") is not False or "error" not in result:
            return _fail("Expected ok=False and error in response on invalid request")
            
        return _pass(f"Degradation oracle passed: gracefully rejected {target} ({type_})")
    except Exception as exc:
        return _fail(f"Degradation executor error: {exc}")


def _exec_performance(sandbox_path: str, scenario: Dict[str, Any]) -> Dict[str, str]:
    """Performance domain: verify execution speed boundaries are honored."""
    try:
        import docx
        import openpyxl
        from sidecar.masters.word_master import WordContextExtractor, WordWriter
    except ImportError as exc:
        return _skip(f"Performance dependencies missing: {exc}")

    try:
        expected = scenario.get("expected_outcome", {})
        type_ = expected.get("type")
        count = expected.get("count")
        max_sec = expected.get("max_seconds")

        if type_ == "perf_docx":
            path = os.path.join(sandbox_path, f"perf_{scenario['id']}.docx")
            doc = docx.Document()
            doc.add_heading("Perf Test Title", level=1)
            for j in range(count):
                doc.add_paragraph(f"This is paragraph {j} of performance testing document. Let's make it long enough to measure.")
            doc.save(path)

            extractor = WordContextExtractor()
            t0 = time.perf_counter()
            ctx = extractor.extract(path, cursor_paragraph_index=0)
            elapsed = time.perf_counter() - t0

            if elapsed > max_sec:
                return _fail(f"DOCX context extraction took {elapsed:.3f}s, exceeding {max_sec}s threshold for {count} paragraphs")
            return _pass(f"DOCX performance oracle passed: {elapsed:.3f}s < {max_sec}s for {count} paragraphs")

        elif type_ == "perf_xlsx":
            path = os.path.join(sandbox_path, f"perf_{scenario['id']}.xlsx")
            wb = openpyxl.Workbook()
            ws = wb.active
            for r in range(1, count + 1):
                ws.cell(row=r, column=1, value=r)
                ws.cell(row=r, column=2, value=f"Row {r} cell value")
            wb.save(path)

            from sidecar.masters.excel_master import ExcelContextExtractor
            extractor = ExcelContextExtractor()
            t0 = time.perf_counter()
            ctx = extractor.extract(path, active_cell="A1")
            elapsed = time.perf_counter() - t0

            if elapsed > max_sec:
                return _fail(f"Excel context extraction took {elapsed:.3f}s, exceeding {max_sec}s threshold for {count} rows")
            return _pass(f"Excel performance oracle passed: {elapsed:.3f}s < {max_sec}s for {count} rows")

        return _fail(f"Unknown performance type: {type_}")
    except Exception as exc:
        return _fail(f"Performance executor error: {exc}")


# ── Dispatch table ────────────────────────────────────────────────────────────
_CATEGORY_EXECUTORS: Dict[str, Any] = {
    "Word":        _exec_word,
    "Excel":       _exec_excel,
    "PPT":         _exec_ppt,
    "Legal":       _exec_legal,
    "CUA":         _exec_cua,
    "Security":    _exec_security,
    "Memory":      _exec_memory,
    "Offline":     _exec_offline,
    "Degradation": _exec_degradation,
    "Performance": _exec_performance,
}


# ── Core runner ───────────────────────────────────────────────────────────────

def run_scenario(scenario: Dict[str, Any], base_dir: str) -> Dict[str, Any]:
    """Execute one scenario and return a result dict. Never raises."""
    sid = scenario.get("id", "?")
    cat = scenario.get("category", "Unknown")
    status = scenario.get("status", "pending")
    t0 = time.perf_counter()

    if status == "excluded":
        return {
            "id": sid, "category": cat, "status": status,
            "oracle_verdict": "SKIP",
            "reason": "Scenario excluded from gauntlet",
            "elapsed_s": 0.0,
        }

    executor = _CATEGORY_EXECUTORS.get(cat)
    if executor is None:
        return {
            "id": sid, "category": cat, "status": status,
            "oracle_verdict": "SKIP",
            "reason": f"No executor for category '{cat}'",
            "elapsed_s": round(time.perf_counter() - t0, 4),
        }

    sandbox = tempfile.mkdtemp(prefix=f"gauntlet_{sid}_", dir=base_dir)
    try:
        verdict_dict = executor(sandbox, scenario)
        elapsed = time.perf_counter() - t0
        return {
            "id": sid, "category": cat, "status": status,
            "oracle_verdict": verdict_dict.get("oracle_verdict", "FAIL"),
            "reason": verdict_dict.get("reason", ""),
            "elapsed_s": round(elapsed, 4),
        }
    except Exception as exc:
        elapsed = time.perf_counter() - t0
        return {
            "id": sid, "category": cat, "status": status,
            "oracle_verdict": "FAIL",
            "reason": f"Executor crashed: {exc}",
            "elapsed_s": round(elapsed, 4),
        }
    finally:
        shutil.rmtree(sandbox, ignore_errors=True)


def run_gauntlet(
    scenarios: List[Dict[str, Any]],
    workers: int = 4,
    output_path: str = "task_completion_rate.json",
) -> Dict[str, Any]:
    """Run all scenarios, build report, write JSON, return report dict."""
    base_dir = tempfile.mkdtemp(prefix="kairo_gauntlet_")
    log.info("Running %d scenarios with %d workers", len(scenarios), workers)
    t_start = time.perf_counter()
    results: List[Dict[str, Any]] = []

    with concurrent.futures.ThreadPoolExecutor(max_workers=workers) as ex:
        futures = {ex.submit(run_scenario, sc, base_dir): sc for sc in scenarios}
        for fut in concurrent.futures.as_completed(futures):
            try:
                res = fut.result()
            except Exception as exc:
                sc = futures[fut]
                res = {
                    "id": sc.get("id", "?"), "category": sc.get("category", "?"),
                    "status": sc.get("status", "?"), "oracle_verdict": "FAIL",
                    "reason": f"Future error: {exc}", "elapsed_s": 0.0,
                }
            results.append(res)
            log.info("[%s] %s -> %s (%.3fs)",
                     res["status"], res["id"], res["oracle_verdict"], res["elapsed_s"])

    shutil.rmtree(base_dir, ignore_errors=True)
    elapsed_total = round(time.perf_counter() - t_start, 2)

    # ── Counts ────────────────────────────────────────────────────────────────
    total = len(results)
    active_count   = sum(1 for r in results if r["status"] == "active")
    pending_count  = sum(1 for r in results if r["status"] == "pending")
    excluded_count = sum(1 for r in results if r["status"] == "excluded")
    passed  = sum(1 for r in results if r["oracle_verdict"] == "PASS")
    failed  = sum(1 for r in results if r["oracle_verdict"] == "FAIL")
    skipped = sum(1 for r in results if r["oracle_verdict"] == "SKIP")
    active_passed = sum(
        1 for r in results if r["status"] == "active" and r["oracle_verdict"] == "PASS"
    )
    active_failed = sum(
        1 for r in results if r["status"] == "active" and r["oracle_verdict"] == "FAIL"
    )
    pass_rate_active = round(active_passed / active_count * 100, 2) if active_count else 0.0
    pass_rate_all    = round(passed / total * 100, 2) if total else 0.0

    if pass_rate_active >= GATE_THRESHOLD:
        verdict = "PASS"
    elif pass_rate_active >= 50.0:
        verdict = "PARTIAL"
    else:
        verdict = "FAIL"

    # ── Per-category ──────────────────────────────────────────────────────────
    cat_stats: Dict[str, Dict] = {}
    for r in results:
        cat = r["category"]
        cs = cat_stats.setdefault(cat, {"total": 0, "active": 0, "passed": 0,
                                        "failed": 0, "skipped": 0})
        cs["total"] += 1
        if r["status"] == "active":
            cs["active"] += 1
        ov = r["oracle_verdict"]
        cs["passed"] += ov == "PASS"
        cs["failed"] += ov == "FAIL"
        cs["skipped"] += ov == "SKIP"

    report = {
        "product": "Kairo Phantom",
        "gauntlet_version": GAUNTLET_VERSION,
        "run_at": datetime.now(timezone.utc).isoformat(),
        "elapsed_seconds": elapsed_total,
        "total": total,
        "active": active_count,
        "pending": pending_count,
        "excluded": excluded_count,
        "passed": passed,
        "failed": failed,
        "skipped": skipped,
        "active_passed": active_passed,
        "active_failed": active_failed,
        "pass_rate_active": pass_rate_active,
        "pass_rate_all": pass_rate_all,
        "verdict": verdict,
        "gate_threshold": GATE_THRESHOLD,
        "categories": cat_stats,
        "results": sorted(results, key=lambda r: r["id"]),
    }

    with open(output_path, "w", encoding="utf-8") as fh:
        json.dump(report, fh, indent=2)
    log.info("Report written -> %s", output_path)
    return report


# ── Entry point ───────────────────────────────────────────────────────────────

def main() -> int:
    parser = argparse.ArgumentParser(description="KairoReal headless 200-scenario gauntlet")
    parser.add_argument("--workers", type=int, default=4)
    parser.add_argument("--output",
                        default=os.path.join(_REPO_ROOT, "task_completion_rate.json"))
    parser.add_argument("--scenarios", default=SCENARIOS_JSON)
    args = parser.parse_args()

    # Ensure stdout uses UTF-8 on Windows to avoid cp1252 encode errors
    try:
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
        sys.stderr.reconfigure(encoding="utf-8", errors="replace")
    except AttributeError:
        pass  # Python < 3.7 or non-TextIOWrapper stdout

    if not os.path.exists(args.scenarios):
        log.error("scenarios.json not found: %s", args.scenarios)
        return 2

    with open(args.scenarios, encoding="utf-8") as fh:
        scenarios = json.load(fh)
    log.info("Loaded %d scenarios from %s", len(scenarios), args.scenarios)

    report = run_gauntlet(scenarios, workers=args.workers, output_path=args.output)

    print()
    print("=" * 68)
    print("  KAIROREAL GAUNTLET - %s" % report["verdict"])
    print("=" * 68)
    print("  Total    : %d  (active=%d pending=%d excl=%d)" % (
          report["total"], report["active"], report["pending"], report["excluded"]))
    print("  Results  : PASS=%d FAIL=%d SKIP=%d" % (
          report["passed"], report["failed"], report["skipped"]))
    print("  Active   : %d passed / %d -> %.1f%%  (gate=%.0f%%)" % (
          report["active_passed"], report["active"],
          report["pass_rate_active"], report["gate_threshold"]))
    print("  Elapsed  : %.1fs" % report["elapsed_seconds"])
    print("  Report   : %s" % args.output)
    print("=" * 68)

    return 0 if report["verdict"] == "PASS" else 1


if __name__ == "__main__":
    sys.exit(main())
