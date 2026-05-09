#!/usr/bin/env python3
"""
Office PowerPoint MCP Bridge — Phase 2 of Kairo Phantom v3.0
Full .pptx creation, slide management, image injection via python-pptx.
JSON-RPC stdio transport — spawned by kairo-phantom via mcp_bridge.rs.

Requirements:
    pip install python-pptx pillow

Usage:
    python server.py
"""

import json
import sys
import os
import base64
import io
import tempfile
import subprocess
from pathlib import Path
from typing import Any, Dict, List, Optional

_BRIDGE_DIR = Path(__file__).parent

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("WARNING: python-pptx not installed. Run: pip install python-pptx pillow", file=sys.stderr)

# ─── Tool Implementations ─────────────────────────────────────────────────────

def check_pptx() -> Optional[Dict]:
    if not PPTX_AVAILABLE:
        return {"error": "python-pptx not installed. Run: pip install python-pptx pillow"}
    return None

def create_presentation(args: Dict) -> Dict:
    """Create a full .pptx presentation with multiple slides."""
    err = check_pptx()
    if err: return err

    output_path = args.get("output_path", "")
    slides_data = args.get("slides", [])
    title_only = args.get("title_only", False)

    if not output_path:
        # Default: save to temp dir
        output_path = os.path.join(tempfile.gettempdir(), "kairo_presentation.pptx")

    prs = Presentation()

    # Slide layouts:
    # 0 = Title Slide, 1 = Title and Content, 5 = Blank, 6 = Title Only
    slide_layouts = {
        "title": 0,
        "content": 1,
        "blank": 5,
        "title_only": 6,
    }

    for i, slide_spec in enumerate(slides_data):
        layout_name = slide_spec.get("layout", "title" if i == 0 else "content")
        layout_idx = slide_layouts.get(layout_name, 1)
        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)

        # Set title
        title_text = slide_spec.get("title", "")
        if title_text and slide.shapes.title:
            slide.shapes.title.text = title_text

        # Set content/body
        content_text = slide_spec.get("content", "")
        if content_text:
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.idx != 0:  # Not the title
                    placeholder.text = content_text
                    break

        # Add image if provided
        image_b64 = slide_spec.get("image_base64", "")
        if image_b64:
            try:
                img_bytes = base64.b64decode(image_b64)
                img_stream = io.BytesIO(img_bytes)
                # Place image on right side of slide
                left = Inches(6)
                top = Inches(1.5)
                width = Inches(3.5)
                slide.shapes.add_picture(img_stream, left, top, width=width)
            except Exception as e:
                print(f"WARNING: Failed to add image to slide {i}: {e}", file=sys.stderr)

        # Add speaker notes if provided
        notes_text = slide_spec.get("notes", "")
        if notes_text:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text

    # Save
    prs.save(output_path)

    return {
        "status": "ok",
        "output_path": output_path,
        "slide_count": len(slides_data),
        "file_size_bytes": os.path.getsize(output_path)
    }

def add_slide(args: Dict) -> Dict:
    """Add a new slide to an existing presentation."""
    err = check_pptx()
    if err: return err

    pptx_path = args.get("pptx_path", "")
    title = args.get("title", "")
    content = args.get("content", "")
    layout = args.get("layout", "content")

    if not pptx_path or not os.path.exists(pptx_path):
        return {"error": f"PPTX file not found: {pptx_path}"}

    prs = Presentation(pptx_path)

    layout_map = {"title": 0, "content": 1, "blank": 5, "title_only": 6}
    layout_idx = layout_map.get(layout, 1)
    slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(slide_layout)

    if title and slide.shapes.title:
        slide.shapes.title.text = title

    if content:
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.idx != 0:
                placeholder.text = content
                break

    prs.save(pptx_path)
    return {
        "status": "ok",
        "slide_count": len(prs.slides),
        "new_slide_index": len(prs.slides) - 1
    }

def add_image_to_slide(args: Dict) -> Dict:
    """Add an image to a specific slide in a PPTX file."""
    err = check_pptx()
    if err: return err

    pptx_path = args.get("pptx_path", "")
    slide_index = args.get("slide_index", 0)
    image_base64 = args.get("image_base64", "")
    left_inches = args.get("left_inches", 1.0)
    top_inches = args.get("top_inches", 1.5)
    width_inches = args.get("width_inches", 4.0)

    if not pptx_path or not os.path.exists(pptx_path):
        return {"error": f"PPTX file not found: {pptx_path}"}
    if not image_base64:
        return {"error": "image_base64 is required"}

    prs = Presentation(pptx_path)
    if slide_index >= len(prs.slides):
        return {"error": f"Slide index {slide_index} out of range (total: {len(prs.slides)})"}

    slide = prs.slides[slide_index]
    try:
        img_bytes = base64.b64decode(image_base64)
        img_stream = io.BytesIO(img_bytes)
        slide.shapes.add_picture(
            img_stream,
            Inches(left_inches),
            Inches(top_inches),
            width=Inches(width_inches)
        )
    except Exception as e:
        return {"error": f"Failed to add image: {str(e)}"}

    prs.save(pptx_path)
    return {
        "status": "ok",
        "pptx_path": pptx_path,
        "slide_index": slide_index
    }

def get_slide_info(args: Dict) -> Dict:
    """Get information about a specific slide."""
    err = check_pptx()
    if err: return err

    pptx_path = args.get("pptx_path", "")
    slide_index = args.get("slide_index", 0)

    if not pptx_path or not os.path.exists(pptx_path):
        return {"error": f"PPTX file not found: {pptx_path}"}

    prs = Presentation(pptx_path)
    if slide_index >= len(prs.slides):
        return {"error": f"Slide index out of range"}

    slide = prs.slides[slide_index]
    shapes_info = []
    for shape in slide.shapes:
        shape_info = {
            "name": shape.name,
            "shape_type": str(shape.shape_type),
        }
        if shape.has_text_frame:
            shape_info["text"] = shape.text_frame.text
        shapes_info.append(shape_info)

    return {
        "status": "ok",
        "slide_index": slide_index,
        "total_slides": len(prs.slides),
        "shapes": shapes_info
    }

def list_slides(args: Dict) -> Dict:
    """List all slides in a PPTX file with their titles."""
    err = check_pptx()
    if err: return err

    pptx_path = args.get("pptx_path", "")
    if not pptx_path or not os.path.exists(pptx_path):
        return {"error": f"PPTX file not found: {pptx_path}"}

    prs = Presentation(pptx_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text
        slides.append({"index": i, "title": title})

    return {
        "status": "ok",
        "pptx_path": pptx_path,
        "slide_count": len(prs.slides),
        "slides": slides
    }

def generate_ai_presentation(args: Dict) -> Dict:
    """
    High-level: generate a full AI-powered presentation from a topic string.
    This uses the provided slide_specs (from Kairo's Swarm Brain) to create the deck.
    """
    err = check_pptx()
    if err: return err

    topic = args.get("topic", "")
    slide_specs = args.get("slide_specs", [])
    output_path = args.get("output_path", "")

    if not slide_specs:
        # Generate minimal template slides if no specs provided
        slide_specs = [
            {"title": topic, "content": "", "layout": "title"},
            {"title": "Overview", "content": f"Key points about {topic}", "layout": "content"},
            {"title": "Details", "content": "Add your content here", "layout": "content"},
            {"title": "Conclusion", "content": f"Summary of {topic}", "layout": "content"},
        ]

    if not output_path:
        safe_topic = "".join(c if c.isalnum() or c in " _-" else "_" for c in topic)[:30]
        output_path = os.path.join(tempfile.gettempdir(), f"kairo_{safe_topic}.pptx")

    return create_presentation({"output_path": output_path, "slides": slide_specs})


def add_text_box(args: Dict) -> Dict:
    """Add a styled text box to a specific slide."""
    err = check_pptx()
    if err: return err

    pptx_path = args.get("pptx_path", "")
    slide_index = args.get("slide_index", 0)
    text = args.get("text", "")
    left = args.get("left_inches", 1.0)
    top = args.get("top_inches", 1.0)
    width = args.get("width_inches", 6.0)
    height = args.get("height_inches", 1.5)
    font_size = args.get("font_size_pt", 18)
    bold = args.get("bold", False)
    color_hex = args.get("color_hex", "000000")

    if not pptx_path or not os.path.exists(pptx_path):
        return {"error": f"PPTX not found: {pptx_path}"}

    prs = Presentation(pptx_path)
    if slide_index >= len(prs.slides):
        return {"error": f"Slide {slide_index} out of range"}

    slide = prs.slides[slide_index]
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    try:
        r, g, b = int(color_hex[0:2], 16), int(color_hex[2:4], 16), int(color_hex[4:6], 16)
        run.font.color.rgb = RGBColor(r, g, b)
    except Exception:
        pass

    prs.save(pptx_path)
    return {"status": "ok", "slide_index": slide_index, "text_preview": text[:50]}


def apply_theme_colors(args: Dict) -> Dict:
    """Apply a color theme to a presentation (background + accent colors via XML)."""
    err = check_pptx()
    if err: return err

    pptx_path = args.get("pptx_path", "")
    theme_name = args.get("theme", "dark")  # dark | light | corporate | minimal

    themes = {
        "dark": {"bg": "0f172a", "title": "e2e8f0", "accent": "6366f1"},
        "light": {"bg": "ffffff", "title": "1e293b", "accent": "3b82f6"},
        "corporate": {"bg": "f8fafc", "title": "0f172a", "accent": "dc2626"},
        "minimal": {"bg": "fafafa", "title": "374151", "accent": "6b7280"},
        "ocean": {"bg": "0c4a6e", "title": "e0f2fe", "accent": "38bdf8"},
        "forest": {"bg": "14532d", "title": "dcfce7", "accent": "4ade80"},
    }

    if not pptx_path or not os.path.exists(pptx_path):
        return {"error": f"PPTX not found: {pptx_path}"}

    colors = themes.get(theme_name, themes["corporate"])
    prs = Presentation(pptx_path)

    def hex_to_rgb(h):
        return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)

    bg_r, bg_g, bg_b = hex_to_rgb(colors["bg"])
    title_r, title_g, title_b = hex_to_rgb(colors["title"])
    accent_r, accent_g, accent_b = hex_to_rgb(colors["accent"])

    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    for slide in prs.slides:
        # Set slide background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(bg_r, bg_g, bg_b)

        # Recolor title shapes
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if shape.name.lower().startswith("title"):
                            run.font.color.rgb = RGBColor(title_r, title_g, title_b)

    prs.save(pptx_path)
    return {"status": "ok", "theme_applied": theme_name, "colors": colors}


def add_chart_slide(args: Dict) -> Dict:
    """Add a slide with a bar/line/pie chart from data."""
    err = check_pptx()
    if err: return err

    pptx_path = args.get("pptx_path", "")
    title = args.get("title", "Chart")
    chart_type = args.get("chart_type", "bar")  # bar | line | pie
    categories = args.get("categories", ["A", "B", "C"])
    series_data = args.get("series", [{"name": "Value", "data": [1, 2, 3]}])

    if not pptx_path or not os.path.exists(pptx_path):
        return {"error": f"PPTX not found: {pptx_path}"}

    try:
        from pptx.chart.data import ChartData
        from pptx.enum.chart import XL_CHART_TYPE

        chart_type_map = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "bar_stacked": XL_CHART_TYPE.COLUMN_STACKED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE,
            "area": XL_CHART_TYPE.AREA,
        }
        xl_type = chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

        prs = Presentation(pptx_path)
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        if slide.shapes.title:
            slide.shapes.title.text = title

        chart_data = ChartData()
        chart_data.categories = categories
        for series in series_data:
            chart_data.add_series(series.get("name", "Series"), tuple(series.get("data", [])))

        chart = slide.shapes.add_chart(xl_type, Inches(1), Inches(1.5), Inches(8), Inches(5), chart_data)
        chart.chart.has_legend = True

        prs.save(pptx_path)
        return {
            "status": "ok",
            "slide_added": len(prs.slides) - 1,
            "chart_type": chart_type,
            "categories_count": len(categories)
        }
    except Exception as e:
        return {"error": f"Chart creation failed: {str(e)}"}


def duplicate_slide(args: Dict) -> Dict:
    """Duplicate a slide by index."""
    err = check_pptx()
    if err: return err

    pptx_path = args.get("pptx_path", "")
    slide_index = args.get("slide_index", 0)

    if not pptx_path or not os.path.exists(pptx_path):
        return {"error": f"PPTX not found: {pptx_path}"}

    import copy
    from lxml import etree

    prs = Presentation(pptx_path)
    if slide_index >= len(prs.slides):
        return {"error": f"Slide {slide_index} out of range"}

    template = prs.slides[slide_index]
    blank_layout = prs.slide_layouts[5]  # Blank
    dup = prs.slides.add_slide(blank_layout)

    # Copy all shapes via XML clone
    for shape in template.shapes:
        el = copy.deepcopy(shape.element)
        dup.shapes._spTree.insert(2, el)

    prs.save(pptx_path)
    return {"status": "ok", "duplicated_from": slide_index, "new_index": len(prs.slides) - 1}


def export_speaker_notes(args: Dict) -> Dict:
    """Export all speaker notes from a presentation as structured text."""
    err = check_pptx()
    if err: return err

    pptx_path = args.get("pptx_path", "")
    if not pptx_path or not os.path.exists(pptx_path):
        return {"error": f"PPTX not found: {pptx_path}"}

    prs = Presentation(pptx_path)
    notes = []
    for i, slide in enumerate(prs.slides):
        title = slide.shapes.title.text if slide.shapes.title else f"Slide {i + 1}"
        note_text = ""
        if slide.has_notes_slide:
            note_text = slide.notes_slide.notes_text_frame.text
        notes.append({"slide": i, "title": title, "notes": note_text})

    return {"status": "ok", "total_slides": len(prs.slides), "notes": notes}


def set_slide_notes(args: Dict) -> Dict:
    """Set speaker notes for a specific slide."""
    err = check_pptx()
    if err: return err

    pptx_path = args.get("pptx_path", "")
    slide_index = args.get("slide_index", 0)
    notes_text = args.get("notes", "")

    if not pptx_path or not os.path.exists(pptx_path):
        return {"error": f"PPTX not found: {pptx_path}"}

    prs = Presentation(pptx_path)
    if slide_index >= len(prs.slides):
        return {"error": f"Slide {slide_index} out of range"}

    slide = prs.slides[slide_index]
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text
    prs.save(pptx_path)
    return {"status": "ok", "slide_index": slide_index}


def canva_autofill_stub(args: Dict) -> Dict:
    """
    Canva Connect API stub — POST /v1/autofills.
    Requires CANVA_ACCESS_TOKEN env var and a template brandTemplateId.
    Full implementation requires Canva Connect API credentials.
    """
    import os
    import urllib.request
    import urllib.error

    token = os.environ.get("CANVA_ACCESS_TOKEN", "")
    template_id = args.get("brand_template_id", "")
    fields = args.get("fields", {})

    if not token:
        return {
            "status": "stub",
            "message": "Set CANVA_ACCESS_TOKEN env var to enable Canva Connect. Get token at: https://www.canva.com/developers/",
            "fields_received": list(fields.keys())
        }

    if not template_id:
        return {"error": "brand_template_id required"}

    payload = {"brandTemplateId": template_id, "data": fields}
    data = json.dumps(payload).encode("utf-8")
    req = urllib.request.Request(
        "https://api.canva.com/rest/v1/autofills",
        data=data,
        headers={"Authorization": f"Bearer {token}", "Content-Type": "application/json"},
        method="POST"
    )
    try:
        with urllib.request.urlopen(req, timeout=30) as resp:
            result = json.loads(resp.read())
        return {"status": "ok", "canva_response": result}
    except urllib.error.HTTPError as e:
        return {"error": f"Canva API error {e.code}: {e.read().decode()}"}
    except Exception as e:
        return {"error": str(e)}



# ─── Tool Registry ────────────────────────────────────────────────────────────

TOOLS = [
    {
        "name": "create_presentation",
        "description": "Create a complete .pptx presentation with multiple slides, images, and speaker notes",
        "inputSchema": {
            "type": "object",
            "properties": {
                "output_path": {"type": "string", "description": "Output .pptx file path"},
                "slides": {
                    "type": "array",
                    "description": "Array of slide specs with title, content, layout, image_base64, notes",
                    "items": {"type": "object"}
                }
            }
        }
    },
    {
        "name": "add_slide",
        "description": "Add a new slide to an existing PPTX file",
        "inputSchema": {
            "type": "object",
            "properties": {
                "pptx_path": {"type": "string"},
                "title": {"type": "string"},
                "content": {"type": "string"},
                "layout": {"type": "string", "enum": ["title", "content", "blank", "title_only"]}
            },
            "required": ["pptx_path"]
        }
    },
    {
        "name": "add_image_to_slide",
        "description": "Add a base64-encoded image to a specific slide in a PPTX file",
        "inputSchema": {
            "type": "object",
            "properties": {
                "pptx_path": {"type": "string"},
                "slide_index": {"type": "integer"},
                "image_base64": {"type": "string"},
                "left_inches": {"type": "number"},
                "top_inches": {"type": "number"},
                "width_inches": {"type": "number"}
            },
            "required": ["pptx_path", "image_base64"]
        }
    },
    {
        "name": "get_slide_info",
        "description": "Get detailed information about a specific slide",
        "inputSchema": {
            "type": "object",
            "properties": {
                "pptx_path": {"type": "string"},
                "slide_index": {"type": "integer"}
            },
            "required": ["pptx_path"]
        }
    },
    {
        "name": "list_slides",
        "description": "List all slides in a PPTX file",
        "inputSchema": {
            "type": "object",
            "properties": {
                "pptx_path": {"type": "string"}
            },
            "required": ["pptx_path"]
        }
    },
    {
        "name": "generate_ai_presentation",
        "description": "High-level: create a full AI-powered presentation from topic + slide specs",
        "inputSchema": {
            "type": "object",
            "properties": {
                "topic": {"type": "string"},
                "slide_specs": {"type": "array", "items": {"type": "object"}},
                "output_path": {"type": "string"}
            },
            "required": ["topic"]
        }
    }
    ,
    {
        "name": "export_revealjs",
        "description": "Export slide data as an interactive Reveal.js HTML presentation (dual export alongside PPTX)",
        "inputSchema": {
            "type": "object",
            "properties": {
                "slides": {"type": "array", "items": {"type": "object"}, "description": "Array of slide objects with title/content/image/notes/layout"},
                "output_path": {"type": "string"},
                "title": {"type": "string"},
                "theme": {"type": "string", "enum": ["black", "white", "moon", "sky", "beige", "night", "serif", "simple", "solarized"]},
                "transition": {"type": "string", "enum": ["none", "fade", "slide", "convex", "zoom"]}
            },
            "required": ["slides", "output_path"]
        }
    },
    {
        "name": "render_transition_video",
        "description": "Render a cinematic physics transition between two slide images as MP4",
        "inputSchema": {
            "type": "object",
            "properties": {
                "from_image": {"type": "string"},
                "to_image": {"type": "string"},
                "effect": {"type": "string", "enum": ["cloth_tear", "glitch_reveal", "gl_wipe_left", "gl_crosszoom", "gl_cube", "gl_ripple", "particle_disintegrate", "cinema_fade", "ascii_dissolve"]},
                "output_path": {"type": "string"}
            },
            "required": ["from_image", "to_image", "effect", "output_path"]
        }
    }
]


def export_revealjs(args: Dict) -> Dict:
    """Export slides as Reveal.js HTML (Advancement 5)."""
    slides = args.get("slides", [])
    output_path = args.get("output_path", "")
    title = args.get("title", "Kairo Presentation")
    theme = args.get("theme", "black")
    transition = args.get("transition", "fade")

    if not slides:
        return {"error": "No slides provided"}
    if not output_path:
        output_path = str(Path(tempfile.mkdtemp()) / "presentation.html")

    export_script = _BRIDGE_DIR / "revealjs_export.py"
    if not export_script.exists():
        return {"error": f"revealjs_export.py not found at {export_script}"}

    import json as _json
    slides_json = _json.dumps(slides)
    # Call the export module directly
    try:
        import sys as _sys
        sys_path_backup = _sys.path[:]
        _sys.path.insert(0, str(_BRIDGE_DIR))
        from revealjs_export import export_revealjs as _export
        result = _export(slides, output_path, title=title, theme=theme, transition=transition)
        _sys.path = sys_path_backup
        return result
    except Exception as e:
        return {"error": str(e)}


def render_transition_video(args: Dict) -> Dict:
    """Render cinematic transition between two slides (Advancement 4)."""
    from_img = args.get("from_image", "")
    to_img = args.get("to_image", "")
    effect = args.get("effect", "cinema_fade")
    output = args.get("output_path", "")

    if not output:
        output = str(Path(tempfile.mkdtemp()) / f"{effect}.mp4")

    effects_script = _BRIDGE_DIR.parent / "kairo-effects" / "effects_engine.py"
    if not effects_script.exists():
        return {"error": f"effects_engine.py not found at {effects_script}"}

    try:
        result = subprocess.run(
            ["python", str(effects_script), "render",
             "--effect", effect, "--from", from_img, "--to", to_img, "--out", output],
            capture_output=True, text=True, timeout=120
        )
        if result.returncode == 0:
            return {"success": True, "output": output, "effect": effect}
        return {"error": result.stderr[-500:]}
    except Exception as e:
        return {"error": str(e)}


TOOL_FNS = {
    "create_presentation": create_presentation,
    "add_slide": add_slide,
    "add_image_to_slide": add_image_to_slide,
    "get_slide_info": get_slide_info,
    "list_slides": list_slides,
    "generate_ai_presentation": generate_ai_presentation,
    "add_text_box": add_text_box,
    "apply_theme_colors": apply_theme_colors,
    "add_chart_slide": add_chart_slide,
    "duplicate_slide": duplicate_slide,
    "export_speaker_notes": export_speaker_notes,
    "set_slide_notes": set_slide_notes,
    "canva_autofill": canva_autofill_stub,
    "export_revealjs": export_revealjs,
    "render_transition_video": render_transition_video,
}

# ─── JSON-RPC Server ──────────────────────────────────────────────────────────

def respond(id, result=None, error=None):
    obj = {"jsonrpc": "2.0", "id": id}
    if result is not None:
        obj["result"] = result
    if error is not None:
        obj["error"] = error
    print(json.dumps(obj), flush=True)

def main():
    status = "ready" if PPTX_AVAILABLE else "degraded (python-pptx missing)"
    print(f"office-pptx-bridge v0.3.0 {status} (stdio transport)", file=sys.stderr, flush=True)

    for raw_line in sys.stdin:
        raw_line = raw_line.strip()
        if not raw_line:
            continue

        try:
            req = json.loads(raw_line)
        except json.JSONDecodeError as e:
            respond(None, error={"code": -32700, "message": f"Parse error: {e}"})
            continue

        req_id = req.get("id")
        method = req.get("method", "")
        params = req.get("params", {})

        if method == "initialize":
            respond(req_id, result={
                "protocolVersion": "2024-11-05",
                "capabilities": {"tools": {}},
                "serverInfo": {
                    "name": "kairo-office-pptx-bridge",
                    "version": "0.3.0",
                    "pptx_available": PPTX_AVAILABLE
                }
            })

        elif method == "tools/list":
            respond(req_id, result={"tools": TOOLS})

        elif method == "tools/call":
            tool_name = params.get("name", "")
            tool_args = params.get("arguments", {})
            fn = TOOL_FNS.get(tool_name)
            if fn is None:
                respond(req_id, error={"code": -32000, "message": f"Unknown tool: {tool_name}"})
            else:
                try:
                    result = fn(tool_args)
                    respond(req_id, result=result)
                except Exception as e:
                    respond(req_id, error={"code": -32000, "message": str(e)})
        else:
            respond(req_id, error={"code": -32601, "message": f"Method not found: {method}"})

if __name__ == "__main__":
    main()
