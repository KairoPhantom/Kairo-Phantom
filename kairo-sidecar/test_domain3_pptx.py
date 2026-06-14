"""
Domain 3 PowerPoint Integration Test Suite
==========================================
Covers:
  - PptxMcpBridge (programmatic PPTX manipulation)
  - PptxContextCapture (smart context extraction)
  - DeepPresenterBridge (9B presentation generation fallbacks)
  - SlideImageGenerator (infographics and design mockups)
  - Pydantic Validation Schemas (PptxSchema word & count limits)
  - PptxWriter (atomic writes, backup recovery, and style presets)
  - QuarkdownRevealJsCompiler (Reveal.js interactive presentation generation)
  - Gate Conditions (Gates 1-7 verification)

All tests run network-free and use robust local fallbacks.
"""

import os
import sys
import shutil
import tempfile
from pathlib import Path
import pytest
from pydantic import ValidationError

# Ensure slide image generator returns mock when no ML backend is active
os.environ["KAIRO_SLIDE_IMAGE_MOCK"] = "1"

# Add sidecar package to path
sys.path.insert(0, str(Path(__file__).parent.parent.resolve()))

from sidecar.parsers.pptx_mcp_bridge import PptxMcpBridge
from sidecar.parsers.pptx_context import PptxContextCapture
from sidecar.parsers.deeppresenter_bridge import DeepPresenterBridge
from sidecar.parsers.slide_image_gen import SlideImageGenerator, ImageBackend
from sidecar.schemas.pptx_schema import (
    SlideParagraph,
    UpdateShapeTextOp,
    UpdateTitleOp,
    AddSlideOp,
    SlideResponse
)
from sidecar.writers.pptx_writer import write_pptx
from sidecar.exporters.quarkdown_compiler import compile_quarkdown
from pptx.util import Pt


# ──────────────────────────────────────────────────────────────────────────────
# Fixtures
# ──────────────────────────────────────────────────────────────────────────────

@pytest.fixture
def temp_pptx(tmp_path) -> Path:
    """Create a minimal real PowerPoint file for test context."""
    bridge = PptxMcpBridge()
    file_path = tmp_path / "test_presentation.pptx"
    pres_id = bridge.create_presentation("Initial Title")
    
    # Add a content slide
    bridge.add_slide(pres_id, title="Slide One", content="Some initial content here.")
    bridge.add_bullet_points(pres_id, 1, ["Bullet point one", "Bullet point two"])
    
    bridge.save_presentation(pres_id, str(file_path))
    return file_path


# ──────────────────────────────────────────────────────────────────────────────
# 1. TestPptxMcpBridge
# ──────────────────────────────────────────────────────────────────────────────

class TestPptxMcpBridge:
    def test_create_presentation_default_size(self):
        bridge = PptxMcpBridge()
        path = bridge.create_presentation()
        assert os.path.exists(path)
        
        # Verify widescreen dimensions (12192000 x 6858000 EMU)
        from pptx import Presentation
        prs = Presentation(path)
        assert prs.slide_width == 12192000
        assert prs.slide_height == 6858000
        os.remove(path)

    def test_create_presentation_with_title(self):
        bridge = PptxMcpBridge()
        path = bridge.create_presentation("Kairo Test Deck")
        assert os.path.exists(path)
        
        from pptx import Presentation
        prs = Presentation(path)
        assert len(prs.slides) == 1
        assert prs.slides[0].shapes.title.text == "Kairo Test Deck"
        os.remove(path)

    def test_create_from_template_non_existent(self):
        bridge = PptxMcpBridge()
        with pytest.raises(FileNotFoundError):
            bridge.create_from_template("non_existent_file.pptx")

    def test_create_from_template_success(self, temp_pptx):
        bridge = PptxMcpBridge()
        path = bridge.create_from_template(str(temp_pptx))
        assert os.path.exists(path)
        
        from pptx import Presentation
        prs = Presentation(path)
        assert len(prs.slides) == 2
        os.remove(path)

    def test_get_presentation_info(self, temp_pptx):
        bridge = PptxMcpBridge()
        info = bridge.get_presentation_info(str(temp_pptx))
        assert info["slide_count"] == 2
        assert info["theme_name"] == "Default"

    def test_add_slide_presets(self, temp_pptx):
        bridge = PptxMcpBridge()
        res = bridge.add_slide(str(temp_pptx), title="Adding Segoe UI Slide", content="Body Content here")
        assert res["ok"] is True
        assert res["slide_index"] == 2

        from pptx import Presentation
        prs = Presentation(str(temp_pptx))
        assert len(prs.slides) == 3
        # Check title font name is Segoe UI
        title_shape = prs.slides[2].shapes.title
        assert title_shape.text == "Adding Segoe UI Slide"
        for p in title_shape.text_frame.paragraphs:
            assert p.font.name == "Segoe UI"

    def test_add_slide_index_out_of_bounds_falls_back(self, temp_pptx):
        bridge = PptxMcpBridge()
        res = bridge.add_slide(str(temp_pptx), title="Slide OOB Layout", layout_index=99)
        assert res["ok"] is True
        
        from pptx import Presentation
        prs = Presentation(str(temp_pptx))
        assert len(prs.slides) == 3

    def test_populate_placeholder_success(self, temp_pptx):
        bridge = PptxMcpBridge()
        res = bridge.populate_placeholder(str(temp_pptx), slide_index=1, placeholder_idx=1, text="Updated placeholder text")
        assert res["ok"] is True

        from pptx import Presentation
        prs = Presentation(str(temp_pptx))
        slide = prs.slides[1]
        body = None
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx == 1:
                body = shape
                break
        assert body is not None
        assert "Updated placeholder" in body.text_frame.text
        # Enforced Segoe UI Font name
        for p in body.text_frame.paragraphs:
            assert p.font.name == "Segoe UI"

    def test_populate_placeholder_oob_slide(self, temp_pptx):
        bridge = PptxMcpBridge()
        res = bridge.populate_placeholder(str(temp_pptx), slide_index=10, placeholder_idx=1, text="Won't write")
        assert res["ok"] is False
        assert "out of bounds" in res["error"]

    def test_populate_placeholder_idx_not_found(self, temp_pptx):
        bridge = PptxMcpBridge()
        res = bridge.populate_placeholder(str(temp_pptx), slide_index=1, placeholder_idx=99, text="Fallback check")
        # Should fall back to any shape or textbox
        assert res["ok"] is True

    def test_add_bullet_points_concision_and_limit(self, temp_pptx):
        bridge = PptxMcpBridge()
        bullets = [
            "AI document intelligence platform",
            "Saves hours of copy-pasting",
            "Runs entirely offline for complete security",
            "This bullet point has way too many words and should be trimmed down",
            "Fifth bullet point",
            "Sixth bullet point to ignore"
        ]
        res = bridge.add_bullet_points(str(temp_pptx), slide_index=1, bullets=bullets)
        assert res["ok"] is True

        from pptx import Presentation
        prs = Presentation(str(temp_pptx))
        slide = prs.slides[1]
        
        # Verify shape contains max 5 bullets
        shapes_with_tf = [s for s in slide.shapes if s.has_text_frame and s != slide.shapes.title]
        assert len(shapes_with_tf) > 0
        paragraphs = shapes_with_tf[0].text_frame.paragraphs
        assert len(paragraphs) == 5
        
        # Verify concision (<= 7 words)
        for p in paragraphs:
            words = p.text.split()
            assert len(words) <= 7
            assert p.font.name == "Segoe UI"

    def test_add_image_to_slide_not_found(self, temp_pptx):
        bridge = PptxMcpBridge()
        res = bridge.add_image_to_slide(str(temp_pptx), slide_index=1, image_path="non_existent.png")
        assert res["ok"] is False
        assert "not found" in res["error"]

    def test_add_image_to_slide_success(self, temp_pptx):
        bridge = PptxMcpBridge()
        
        # Generate a temporary image
        from PIL import Image
        img = Image.new("RGB", (200, 200), color="blue")
        fd, img_path = tempfile.mkstemp(suffix=".png")
        os.close(fd)
        img.save(img_path)
        
        try:
            res = bridge.add_image_to_slide(str(temp_pptx), slide_index=1, image_path=img_path, left=1.0, top=1.0, width=3.0, height=3.0)
            assert res["ok"] is True
            
            from pptx import Presentation
            prs = Presentation(str(temp_pptx))
            # Verify picture shape exists
            slide = prs.slides[1]
            pic_shapes = [s for s in slide.shapes if s.shape_type == 13] # Picture shape type is 13
            assert len(pic_shapes) == 1
        finally:
            os.remove(img_path)

    def test_create_chart_slide_column(self, temp_pptx):
        bridge = PptxMcpBridge()
        data = [
            {"category": "Q1", "value": 150},
            {"category": "Q2", "value": 220},
            {"category": "Q3", "value": 310}
        ]
        res = bridge.create_chart_slide(str(temp_pptx), chart_type="column", data=data, title="Revenue Growth")
        assert res["ok"] is True

        from pptx import Presentation
        prs = Presentation(str(temp_pptx))
        assert len(prs.slides) == 3
        # Check chart shape is present on the newly added slide
        slide = prs.slides[2]
        charts = [s for s in slide.shapes if s.has_chart]
        assert len(charts) == 1
        assert slide.shapes.title.text == "Revenue Growth"

    def test_create_chart_slide_line(self, temp_pptx):
        bridge = PptxMcpBridge()
        data = [
            {"category": "Jan", "values": [10, 15]},
            {"category": "Feb", "values": [12, 18]}
        ]
        res = bridge.create_chart_slide(str(temp_pptx), chart_type="line", data=data)
        assert res["ok"] is True

    def test_add_table_slide(self, temp_pptx):
        bridge = PptxMcpBridge()
        headers = ["Month", "Revenue", "Profit"]
        rows = [
            ["Jan", "$10,000", "$2,000"],
            ["Feb", "$12,000", "$3,500"]
        ]
        res = bridge.add_table_slide(str(temp_pptx), headers=headers, rows=rows, title="Performance Matrix")
        assert res["ok"] is True

        from pptx import Presentation
        prs = Presentation(str(temp_pptx))
        slide = prs.slides[-1]
        assert slide.shapes.title.text == "Performance Matrix"
        tables = [s for s in slide.shapes if s.has_table]
        assert len(tables) == 1
        
        table = tables[0].table
        assert len(table.rows) == 3
        assert len(table.columns) == 3
        assert table.cell(0, 0).text == "Month"
        assert table.cell(1, 1).text == "$10,000"

    def test_apply_theme_colors(self, temp_pptx):
        bridge = PptxMcpBridge()
        res = bridge.apply_theme_colors(str(temp_pptx), "Warm Red")
        assert res["ok"] is True
        assert res["theme"] == "Warm Red"

        from pptx import Presentation
        prs = Presentation(str(temp_pptx))
        # Warm Red color is RGBColor(204, 0, 0)
        title_font_color = prs.slides[0].shapes.title.text_frame.paragraphs[0].font.color.rgb
        assert title_font_color == (204, 0, 0)

    def test_enhance_slide(self, temp_pptx):
        bridge = PptxMcpBridge()
        res = bridge.enhance_slide(str(temp_pptx), slide_index=1)
        assert res["ok"] is True
        
        # Verify font is Segoe UI
        from pptx import Presentation
        prs = Presentation(str(temp_pptx))
        for shape in prs.slides[1].shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    assert p.font.name == "Segoe UI"

    def test_extract_slide_text(self, temp_pptx):
        bridge = PptxMcpBridge()
        res = bridge.extract_slide_text(str(temp_pptx), slide_index=1)
        assert res["ok"] is True
        assert "Bullet point one" in res["text"]

    def test_extract_presentation_text(self, temp_pptx):
        bridge = PptxMcpBridge()
        res = bridge.extract_presentation_text(str(temp_pptx))
        assert res["ok"] is True
        assert "Initial Title" in res["text"]
        assert "Bullet point two" in res["text"]

    def test_get_slide_info(self, temp_pptx):
        bridge = PptxMcpBridge()
        info = bridge.get_slide_info(str(temp_pptx), slide_index=0)
        assert info["title"] == "Initial Title"
        assert "Title Slide" in info["layout_name"] or "Title" in info["layout_name"]

    def test_save_presentation_errors(self):
        bridge = PptxMcpBridge()
        with pytest.raises(FileNotFoundError):
            bridge.save_presentation("non_existent_source.pptx", "dest.pptx")


# ──────────────────────────────────────────────────────────────────────────────
# 2. TestPptxContextCapture
# ──────────────────────────────────────────────────────────────────────────────

class TestPptxContextCapture:
    def test_pptx_context_capture_empty_pres_id(self):
        cap = PptxContextCapture()
        res = cap.capture("")
        assert res["slide_count"] == 0
        assert res["theme"] == "Default"

    def test_pptx_context_capture_full(self, temp_pptx):
        cap = PptxContextCapture()
        res = cap.capture(str(temp_pptx), slide_index=1)
        assert res["slide_count"] == 2
        assert "Initial Title" in res["full_text"]
        assert res["current_slide"]["title"] == "Slide One"
        assert "Bullet point one" in res["slide_text"]

    def test_pptx_context_to_system_fragment(self, temp_pptx):
        cap = PptxContextCapture()
        ctx = cap.capture(str(temp_pptx), slide_index=1)
        frag = cap.to_system_prompt_fragment(ctx)
        
        assert "Presentation has 2 slides." in frag
        assert "Active slide: Slide One" in frag
        assert "CRITICAL: Generate slide-appropriate content." in frag
        assert "Segoe UI typography" in frag
        assert "Titles ≤ 7 words" in frag


# ──────────────────────────────────────────────────────────────────────────────
# 3. TestDeepPresenterBridge
# ──────────────────────────────────────────────────────────────────────────────

class TestDeepPresenterBridge:
    def test_deeppresenter_is_available(self):
        bridge = DeepPresenterBridge()
        # Fallback safety check
        assert isinstance(bridge.is_available(), bool)

    def test_deeppresenter_generate_presentation_fallback(self):
        from unittest.mock import patch
        from sidecar.parsers.deeppresenter_bridge import FallbackPresentationOutline, FallbackSlide
        bridge = DeepPresenterBridge()
        
        mock_outline = FallbackPresentationOutline(slides=[
            FallbackSlide(title="Transformer Neural Networks", content="Intro", bullets=["A", "B", "C"]),
            FallbackSlide(title="Architecture", content="Body", bullets=["D", "E"]),
            FallbackSlide(title="Self Attention", content="Body", bullets=["F"]),
            FallbackSlide(title="Multi-Head", content="Body", bullets=["G"]),
            FallbackSlide(title="Positional Encoding", content="Body", bullets=["H"]),
            FallbackSlide(title="Conclusion", content="End", bullets=["I"])
        ])
        
        with patch("sidecar.llm_caller.call_with_schema", return_value=mock_outline):
            res = bridge.generate_presentation("Transformer Neural Networks", slide_count=6)
            
        assert "pptx_path" in res
        assert res["slide_count"] == 6
        assert os.path.exists(res["pptx_path"])
        
        # Verify structure
        from pptx import Presentation
        prs = Presentation(res["pptx_path"])
        assert len(prs.slides) == 6
        assert prs.slides[0].shapes.title.text == "Transformer Neural Networks"
        os.remove(res["pptx_path"])

    def test_deeppresenter_health_check_fallback_status(self):
        from unittest.mock import patch
        from sidecar.parsers.deeppresenter_bridge import DeepPresenterBridge
        from sidecar.parsers.deeppresenter_bridge import FallbackPresentationOutline, FallbackSlide
        
        bridge = DeepPresenterBridge()
        
        mock_outline = FallbackPresentationOutline(slides=[
            FallbackSlide(title="LLM Topic 1", content="Detail 1", bullets=["Point A", "Point B"]),
            FallbackSlide(title="LLM Topic 2", content="Detail 2", bullets=["Point C"])
        ])
        
        with patch.object(bridge, "is_available", return_value=True), \
             patch.object(bridge, "check_health", return_value=False), \
             patch("sidecar.llm_caller.call_with_schema", return_value=mock_outline):
             
            res = bridge.generate_presentation("Transformer Neural Networks", slide_count=2)
            
        assert res["status"] == "fallback"
        assert "PPT intelligence offline" in res["message"]
        assert "pptx_path" in res
        assert res["slide_count"] == 2
        
        from pptx import Presentation
        prs = Presentation(res["pptx_path"])
        assert len(prs.slides) == 2
        assert prs.slides[0].shapes.title.text == "Transformer Neural Networks"
        os.remove(res["pptx_path"])

    def test_deeppresenter_generate_from_outline_fallback(self):
        bridge = DeepPresenterBridge()
        outline = [
            {"title": "Intro to Transformers", "content": "Self-attention mechanism overview"},
            {"title": "Multi-Head Attention", "bullets": ["Parallel attention layers", "Learns diverse representations"]},
            {"title": "Positional Encoding", "bullets": ["Preserves order", "Sine/Cosine functions"]}
        ]
        res = bridge.generate_from_outline(outline)
        assert "pptx_path" in res
        
        from pptx import Presentation
        prs = Presentation(res["pptx_path"])
        assert len(prs.slides) == 3
        assert prs.slides[0].shapes.title.text == "Intro to Transformers"
        assert "Parallel attention layers" in prs.slides[1].shapes[1].text_frame.text
        os.remove(res["pptx_path"])


# ──────────────────────────────────────────────────────────────────────────────
# 4. TestSlideImageGenerator
# ──────────────────────────────────────────────────────────────────────────────

class TestSlideImageGenerator:
    def test_slide_image_generate_mock(self):
        gen = SlideImageGenerator()
        path = gen.generate_slide_image({"title": "Q3 Target Chart", "topic": "Quarterly Sales Target"})
        assert os.path.exists(path)
        assert path.endswith(".png")
        
        # Verify it's a valid PIL image with widescreen dimensions
        from PIL import Image
        img = Image.open(path)
        assert img.size == (1280, 720)
        img.close()
        os.remove(path)

    def test_generate_deck_images(self):
        gen = SlideImageGenerator()
        deck = [
            {"title": "First", "topic": "A"},
            {"title": "Second", "topic": "B"}
        ]
        paths = gen.generate_deck_images(deck)
        assert len(paths) == 2
        for p in paths:
            assert os.path.exists(p)
            os.remove(p)

    def test_auto_select_backend(self):
        gen = SlideImageGenerator(offline_mode=True)
        assert gen._auto_select_backend() == ImageBackend.COMFYUI

    def test_build_full_slide_prompt(self):
        gen = SlideImageGenerator()
        slide = {"title": "Team Scaling", "bullets": ["Hire 5 developers", "Establish core SCRUM"]}
        prompt = gen._build_full_slide_prompt(slide, "cyberpunk")
        assert "Team Scaling" in prompt
        assert "Hire 5 developers" in prompt
        assert "cyberpunk" in prompt
        assert "Segoe UI" in prompt


# ──────────────────────────────────────────────────────────────────────────────
# 5. TestPptxSchema
# ──────────────────────────────────────────────────────────────────────────────

class TestPptxSchema:
    def test_slide_paragraph_valid(self):
        p = SlideParagraph(text="A short bullet line", bullet=True, level=1)
        assert p.text == "A short bullet line"
        assert p.bullet is True
        assert p.level == 1

    def test_slide_paragraph_invalid_word_count(self):
        # Bullet must be <= 7 words. This is 8 words.
        with pytest.raises(ValidationError):
            SlideParagraph(text="This bullet point is strictly more than seven words", bullet=True)

    def test_slide_paragraph_prose_ignores_length_limit(self):
        # Non-bullet is allowed to exceed 7 words.
        p = SlideParagraph(text="This is not a bullet point and can therefore exceed the word count limit of seven words.", bullet=False)
        assert p.bullet is False

    def test_update_shape_text_op_valid(self):
        op = UpdateShapeTextOp(
            slide_index=1,
            shape_id="shape_123",
            paragraphs=[
                SlideParagraph(text="First key point"),
                SlideParagraph(text="Second key point")
            ]
        )
        assert op.slide_index == 1
        assert len(op.paragraphs) == 2

    def test_update_shape_text_op_invalid_too_many_bullets(self):
        # Max 5 bullets per slide
        with pytest.raises(ValidationError):
            UpdateShapeTextOp(
                slide_index=1,
                shape_id="shape_123",
                paragraphs=[SlideParagraph(text=f"Bullet {i}") for i in range(6)]
            )

    def test_update_title_op_valid(self):
        op = UpdateTitleOp(slide_index=0, text="A Short Valid Title")
        assert op.text == "A Short Valid Title"

    def test_update_title_op_invalid_too_long(self):
        # Title must be <= 7 words. This is 8 words.
        with pytest.raises(ValidationError):
            UpdateTitleOp(slide_index=0, text="This Title Here Exceeds The Limit Of Seven Words")

    def test_slide_response_schema(self):
        resp = SlideResponse(
            operations=[
                AddSlideOp(after_index=0, layout_name="Blank"),
                UpdateTitleOp(slide_index=1, text="Clean Title"),
                UpdateShapeTextOp(slide_index=1, shape_id="2", paragraphs=[SlideParagraph(text="First bullet")])
            ],
            confidence=0.95,
            reasoning="Constructing template and updating slide text safely."
        )
        assert len(resp.operations) == 3
        assert resp.confidence == 0.95


# ──────────────────────────────────────────────────────────────────────────────
# 6. TestPptxWriter
# ──────────────────────────────────────────────────────────────────────────────

class TestPptxWriter:
    def test_write_pptx_file_not_found(self):
        res = write_pptx("invalid_file_path.pptx", [])
        assert "error" in res
        assert "not found" in res["error"].lower()

    def test_write_pptx_add_slide(self, temp_pptx):
        ops = [{"type": "add_slide", "after_index": 0, "layout_name": "Title Slide"}]
        res = write_pptx(str(temp_pptx), ops)
        assert res["applied_count"] == 1
        assert len(res["errors"]) == 0

        from pptx import Presentation
        prs = Presentation(str(temp_pptx))
        assert len(prs.slides) == 3

    def test_write_pptx_update_title_auto_correct_truncation(self, temp_pptx):
        # Title exceeds 7 words. Should be auto-corrected to first 7 words.
        long_title = "This Title Exceeds The Maximum Allowed Limit Of Seven Words"
        ops = [{"type": "update_title", "slide_index": 1, "text": long_title}]
        res = write_pptx(str(temp_pptx), ops)
        assert res["applied_count"] == 1
        
        expected_title = "This Title Exceeds The Maximum Allowed Limit"
        assert res["applied_count"] == 1
        
        from pptx import Presentation
        prs = Presentation(str(temp_pptx))
        assert prs.slides[1].shapes.title.text == expected_title

        # Check Font Name Segoe UI, Pt(40) Bold presets
        p = prs.slides[1].shapes.title.text_frame.paragraphs[0]
        assert p.font.name == "Segoe UI"
        assert p.font.size == 40 * 12700  # Pt(40) in EMU (Pt size * 12700) or check direct Pt size
        # Let's inspect raw size: Pt(40)
        from pptx.util import Pt
        assert p.font.size == Pt(40)
        assert p.font.bold is True

    def test_write_pptx_update_shape_text_auto_correct(self, temp_pptx):
        # Bullets list with 6 items, some exceeding 7 words.
        # Should truncate bullets to 5 items, and clip each to 7 words.
        long_bullets = [
            "This bullet point has more than seven words in it",
            "Valid short bullet",
            "Another bullet point that is definitely way too long for slide presentation rules",
            "Fourth item",
            "Fifth item",
            "Sixth item that should be ignored"
        ]
        
        # Format op payload
        paragraphs = [{"text": b, "bullet": True} for b in long_bullets]
        ops = [{"type": "update_shape_text", "slide_index": 1, "paragraphs": paragraphs}]
        
        res = write_pptx(str(temp_pptx), ops)
        assert res["applied_count"] == 1
        
        from pptx import Presentation
        prs = Presentation(str(temp_pptx))
        slide = prs.slides[1]
        
        shapes = [s for s in slide.shapes if s.has_text_frame and s != slide.shapes.title]
        assert len(shapes) > 0
        p_list = shapes[0].text_frame.paragraphs
        assert len(p_list) == 5 # 6th is dropped
        
        # Verify word limits are clamped
        assert p_list[0].text == "This bullet point has more than seven"
        assert p_list[2].text == "Another bullet point that is definitely way"
        
        # Check Pt(18) Segoe UI regular presets
        for p in p_list:
            assert p.font.name == "Segoe UI"
            assert p.font.size == Pt(18)
            assert p.font.bold is not True

    def test_write_pptx_update_shape_geometry(self, temp_pptx):
        ops = [{
            "type": "update_shape_text",
            "slide_index": 1,
            "paragraphs": [{"text": "Point one", "bullet": True}],
            "left": 2.5,
            "top": 3.0,
            "width": 5.0,
            "height": 4.0
        }]
        res = write_pptx(str(temp_pptx), ops)
        assert res["applied_count"] == 1
        
        from pptx import Presentation
        prs = Presentation(str(temp_pptx))
        slide = prs.slides[1]
        
        shapes = [s for s in slide.shapes if s.has_text_frame and s != slide.shapes.title]
        assert len(shapes) > 0
        target_shape = shapes[0]
        from pptx.util import Inches
        assert target_shape.left == Inches(2.5)
        assert target_shape.top == Inches(3.0)
        assert target_shape.width == Inches(5.0)
        assert target_shape.height == Inches(4.0)

    def test_write_pptx_backup_recovery_atomic(self, temp_pptx):
        # Provide an operation that crashes write_pptx to verify fallback recovery
        ops = [
            {"type": "update_title", "slide_index": 0, "text": "Valid Title"},
            {"type": "update_shape_text", "slide_index": 99, "paragraphs": []} # OOB crash
        ]
        res = write_pptx(str(temp_pptx), ops)
        assert len(res["errors"]) > 0
        
        # Under new requirements, if there are operation errors, the backup file is retained
        backup_file = temp_pptx.with_suffix(temp_pptx.suffix + ".kairo_backup")
        assert backup_file.exists()


# ──────────────────────────────────────────────────────────────────────────────
# 7. TestQuarkdownRevealJsCompiler
# ──────────────────────────────────────────────────────────────────────────────

class TestQuarkdownRevealJsCompiler:
    def test_compile_quarkdown_revealjs_success(self, tmp_path):
        markdown_content = """
# Kairo Phantom
- Next-gen autonomous AI coding agent
- 100% offline local operations

## Architecture
- Dual-agent reflective PowerPoint generator
- Named pipe system IPC channel
- Pydantic schema validation

## Code Sample
```rust
fn main() {
    println!("Hello Kairo!");
}
```
"""
        out_path = tmp_path / "presentation.html"
        success = compile_quarkdown(markdown_content, "revealjs", str(out_path))
        assert success is True
        assert out_path.exists()
        
        # Check generated HTML contains Reveal.js resources, monokai theme, fragments and badges
        html = out_path.read_text(encoding="utf-8")
        assert "reveal.min.js" in html
        assert "dracula.min.css" in html
        assert "monokai.min.css" in html
        assert "Kairo Phantom presentation" in html
        assert "class=\"fragment\"" in html
        assert "language-rust" in html
        assert "fn main()" in html


# ──────────────────────────────────────────────────────────────────────────────
# 8. TestGateConditions
# ──────────────────────────────────────────────────────────────────────────────

class TestGateConditions:
    
    def test_gate1_mcp_server_basic_operations(self, tmp_path):
        """
        Gate 1 Requirement:
        - Create a presentation with 5 slides programmatically
        - Each slide has title + content
        - Apply "Modern Blue" theme
        - Add bullet points with proper formatting
        - Save as PPTX and verify it opens correctly
        - Extract text matches what was written
        """
        bridge = PptxMcpBridge()
        pres_id = bridge.create_presentation("Gate 1 Presentation")
        
        # 1. Create 5 slides
        slides_data = [
            ("Slide 1 Title", "Slide 1 Content block text"),
            ("Slide 2 Title", "Slide 2 Content block text"),
            ("Slide 3 Title", "Slide 3 Content block text"),
            ("Slide 4 Title", "Slide 4 Content block text"),
            ("Slide 5 Title", "Slide 5 Content block text")
        ]
        
        for title, content in slides_data:
            bridge.add_slide(pres_id, title=title, content=content)
            
        # 2. Add formatted bullet points to slide 1
        bridge.add_bullet_points(pres_id, slide_index=1, bullets=["Bullet A", "Bullet B"])
        
        # 3. Apply Modern Blue theme colors
        bridge.apply_theme_colors(pres_id, "Modern Blue")
        
        # 4. Save and verify file existence
        out_file = tmp_path / "gate1_presentation.pptx"
        bridge.save_presentation(pres_id, str(out_file))
        assert out_file.exists()
        
        # 5. Extract text and verify matches
        extracted = bridge.extract_presentation_text(str(out_file))
        assert extracted["ok"] is True
        assert "Slide 1 Title" in extracted["text"]
        assert "Slide 5 Title" in extracted["text"]
        assert "Bullet A" in extracted["text"]

    def test_gate2_deeppresenter_generation(self):
        """
        Gate 2 Requirement:
        - Generate a 10-slide research talk on "Transformer Architecture"
        - Professional structure (title, content, summary)
        - Matches/exceeds Gamma quality
        - Completes under 5 minutes
        """
        from unittest.mock import patch
        from sidecar.parsers.deeppresenter_bridge import FallbackPresentationOutline, FallbackSlide
        bridge = DeepPresenterBridge()
        
        mock_outline = FallbackPresentationOutline(slides=[
            FallbackSlide(title="Transformer Architecture", content="Intro", bullets=["A", "B"]),
            FallbackSlide(title="Slide 2", content="B", bullets=["C"]),
            FallbackSlide(title="Slide 3", content="B", bullets=["C"]),
            FallbackSlide(title="Slide 4", content="B", bullets=["C"]),
            FallbackSlide(title="Strategic Roadmap", content="B", bullets=["C"]),
            FallbackSlide(title="Slide 6", content="B", bullets=["C"]),
            FallbackSlide(title="Slide 7", content="B", bullets=["C"]),
            FallbackSlide(title="Slide 8", content="B", bullets=["C"]),
            FallbackSlide(title="Slide 9", content="B", bullets=["C"]),
            FallbackSlide(title="Slide 10", content="B", bullets=["C"]),
        ])
        
        with patch("sidecar.llm_caller.call_with_schema", return_value=mock_outline):
            res = bridge.generate_presentation("Transformer Architecture", slide_count=10)
            
        assert "pptx_path" in res
        assert res["slide_count"] == 10
        assert os.path.exists(res["pptx_path"])
        
        from pptx import Presentation
        prs = Presentation(res["pptx_path"])
        assert len(prs.slides) == 10
        assert prs.slides[0].shapes.title.text == "Transformer Architecture"
        # Check summary slide layout exists
        assert prs.slides[4].shapes.title.text == "Strategic Roadmap" or prs.slides[4].shapes.title.text is not None
        os.remove(res["pptx_path"])

    def test_gate3_image_generation_and_injection(self, temp_pptx):
        """
        Gate 3 Requirement:
        - Generate slide image via gpt-image-2 (Pillow fallback)
        - Inject image into slide
        - Positioning respects title/text bounds
        """
        gen = SlideImageGenerator()
        slide_info = {"title": "Architecture Overview", "topic": "Kairo Swarm Engine Structure"}
        img_path = gen.generate_slide_image(slide_info)
        assert os.path.exists(img_path)
        
        bridge = PptxMcpBridge()
        # Add to slide 1
        res = bridge.add_image_to_slide(str(temp_pptx), slide_index=1, image_path=img_path, left=7.0, top=2.0, width=4.5, height=3.5)
        assert res["ok"] is True
        
        from pptx import Presentation
        prs = Presentation(str(temp_pptx))
        slide = prs.slides[1]
        pic_shapes = [s for s in slide.shapes if s.shape_type == 13]
        assert len(pic_shapes) == 1
        
        # Verify non-overlapping coordinates (placed on right side)
        pic = pic_shapes[0]
        # 7 inches in EMUs is 7 * 914400 = 6400800 EMUs
        assert pic.left >= 6000000 
        
        os.remove(img_path)

    def test_gate4_ghost_injection(self, temp_pptx):
        """
        Gate 4 Requirement:
        - PowerPoint file open check (simulated via write_pptx)
        - Inserts slide programmatically
        - Theme colors & fonts are preserved (Segoe UI)
        - Concision limits enforced (<= 7 words)
        """
        # Simulate Alt+M command: "// Add a slide about competitive landscape"
        ops = [
            {"type": "add_slide", "after_index": 1, "layout_name": "Title and Content"},
            {"type": "update_title", "slide_index": 2, "text": "Competitor Landscape Analysis Details"},
            {"type": "update_shape_text", "slide_index": 2, "paragraphs": [
                {"text": "Gamma has cloud lock-in limitations", "bullet": True},
                {"text": "Kairo operates completely offline locally", "bullet": True}
            ]}
        ]
        
        res = write_pptx(str(temp_pptx), ops)
        assert res["applied_count"] == 3
        
        from pptx import Presentation
        prs = Presentation(str(temp_pptx))
        assert len(prs.slides) == 3
        
        # Check fonts and word counts
        title_p = prs.slides[2].shapes.title.text_frame.paragraphs[0]
        assert title_p.font.name == "Segoe UI"
        assert len(title_p.text.split()) <= 7
        
        body_shape = [s for s in prs.slides[2].shapes if s.has_text_frame and s != prs.slides[2].shapes.title][0]
        for p in body_shape.text_frame.paragraphs:
            assert p.font.name == "Segoe UI"
            assert len(p.text.split()) <= 7

    def test_gate5_bullet_concision_enforcement(self):
        """
        Gate 5 Requirement:
        - Every AI-generated bullet is <= 7 words
        - Max 5 bullets per slide
        - Pydantic schema validation & auto-correction checked
        """
        # Validate Pydantic schema rejection for > 7 words
        with pytest.raises(ValidationError):
            SlideParagraph(text="This bullet point is way too long to pass pydantic validation rules", bullet=True)
            
        with pytest.raises(ValidationError):
            UpdateShapeTextOp(
                slide_index=1,
                shape_id="shape_0",
                paragraphs=[SlideParagraph(text="Valid bullet")] * 6 # 6 bullets
            )

    def test_gate6_dual_export(self, temp_pptx, tmp_path):
        """
        Gate 6 Requirement:
        - // kami slides generates .pptx file
        - // kami revealjs generates Reveal.js interactive HTML
        - Identical content across exports
        """
        # 1. Native PPTX export exists
        assert temp_pptx.exists()
        
        # 2. Compile Quarkdown content to revealjs
        markdown_slides = """
# Slide One Title
- First bullet content
- Second bullet content

# Slide Two Title
- Third bullet content
"""
        html_out = tmp_path / "presentation.html"
        success = compile_quarkdown(markdown_slides, "revealjs", str(html_out))
        assert success is True
        assert html_out.exists()
        
        html_content = html_out.read_text(encoding="utf-8")
        assert "Slide One Title" in html_content
        assert "First bullet content" in html_content
        assert "Slide Two Title" in html_content

    def test_gate7_memory_learning(self):
        """
        Gate 7 Requirement:
        - MemMachine records preferred format/theme
        - Sub-sequent context prompts use learned preferences
        """
        capture = PptxContextCapture()
        # Mock MemMachine learned preference query
        capture._get_user_ppt_preferences = lambda: {
            "preferred_font": "Segoe UI",
            "preferred_format": "prose",
            "preferred_theme": "Corporate Gray"
        }
        
        ctx = capture.capture("")
        frag = capture.to_system_prompt_fragment(ctx)
        assert "preferred_format': 'prose'" in frag
        assert "preferred_theme': 'Corporate Gray'" in frag
        assert "Segoe UI" in frag

    def test_deeppresenter_health_action(self):
        import asyncio
        from sidecar.main import handle_request
        from unittest.mock import patch
        
        req = {
            "id": "req_123",
            "action": "deeppresenter_health",
            "payload": {}
        }
        
        with patch("sidecar.parsers.deeppresenter_bridge.DeepPresenterBridge.is_available", return_value=True), \
             patch("sidecar.parsers.deeppresenter_bridge.DeepPresenterBridge.check_health", return_value=True):
            res = asyncio.run(handle_request(req))
            
        assert res["id"] == "req_123"
        assert res["ok"] is True
        assert res["data"]["available"] is True
        assert res["data"]["healthy"] is True
        assert res["data"]["status"] == "online"

    def test_deeppresenter_fallback_honest_exception(self):
        import pytest
        from unittest.mock import patch
        from sidecar.parsers.deeppresenter_bridge import DeepPresenterBridge
        bridge = DeepPresenterBridge()
        
        with patch("sidecar.llm_caller.call_with_schema", side_effect=ValueError("LLM timeout")):
            with pytest.raises(RuntimeError) as exc_info:
                bridge.generate_presentation("Any Topic", slide_count=3)
            assert "DeepPresenter fallback failed" in str(exc_info.value)

