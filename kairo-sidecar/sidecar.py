#!/usr/bin/env python3
"""
Kairo Phantom Python Sidecar v1.0
==================================
Persistent document I/O process. Rust daemon spawns this on startup and
communicates via TCP socket (localhost:7438) using newline-delimited JSON.

Protocol:
  Request:  {"id":"<uuid>","action":"read_docx|write_docx|read_xlsx|write_xlsx|read_pptx|write_pptx|ping","path":"...","payload":{...}}
  Response: {"id":"<uuid>","ok":true,"data":{...}}  |  {"id":"<uuid>","ok":false,"error":"..."}

This sidecar owns ALL document I/O. Rust owns hotkey, UIA, injection, orchestration.
"""

import asyncio
import json
import logging
import os
import sys
import traceback
import tempfile
import shutil
from pathlib import Path
from typing import Any, Optional

# Enforce clean package imports
sys.path.insert(0, str(Path(__file__).parent.resolve()))

# ─── Logging ─────────────────────────────────────────────────────────────────

log_path = Path.home() / ".kairo-phantom" / "sidecar.log"
log_path.parent.mkdir(parents=True, exist_ok=True)

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [SIDECAR] %(levelname)s %(message)s",
    handlers=[
        logging.FileHandler(log_path, encoding="utf-8"),
        logging.StreamHandler(sys.stderr),
    ],
)
log = logging.getLogger("kairo-sidecar")

# ─── Lazy imports (graceful degradation if package missing) ──────────────────

def _try_import(name):
    try:
        return __import__(name)
    except ImportError:
        log.warning(f"Optional package '{name}' not installed — related features disabled")
        return None

# ─── DOCX Operations ─────────────────────────────────────────────────────────

def _read_docx(path: str) -> dict:
    """
    Read a .docx file and return structured document JSON.
    Returns: {paragraphs, tables, headings, full_text}
    """
    try:
        from docx import Document
        from docx.oxml.ns import qn
    except ImportError:
        return {"error": "python-docx not installed — run: pip install python-docx"}

    doc = Document(path)
    paragraphs = []
    headings = []
    full_text_parts = []

    for i, para in enumerate(doc.paragraphs):
        text = para.text.strip()
        style = para.style.name if para.style else "Normal"
        is_heading = style.startswith("Heading")
        entry = {
            "index": i,
            "text": text,
            "style": style,
            "is_heading": is_heading,
            "heading_level": int(style.replace("Heading ", "")) if is_heading and style != "Heading" else 0,
            "bold": any(run.bold for run in para.runs if run.bold is not None),
            "italic": any(run.italic for run in para.runs if run.italic is not None),
        }
        paragraphs.append(entry)
        if is_heading and text:
            headings.append({"index": i, "level": entry["heading_level"], "text": text})
        if text:
            full_text_parts.append(text)

    tables = []
    for t_idx, table in enumerate(doc.tables):
        rows = []
        for row in table.rows:
            rows.append([cell.text.strip() for cell in row.cells])
        tables.append({"index": t_idx, "rows": rows})

    return {
        "paragraphs": paragraphs,
        "headings": headings,
        "tables": tables,
        "full_text": "\n".join(full_text_parts),
        "paragraph_count": len(paragraphs),
    }


def _write_docx(path: str, operations: list) -> dict:
    """
    Apply DocxOperation list to a .docx file using python-docx.
    Operations: [{"action":"insert_after_heading","heading_text":"...","style":"ListBullet","content":"..."}]

    Atomic write: write to temp file then rename.
    Preserves all existing content — only touches specified sections.
    """
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except ImportError:
        return {"error": "python-docx not installed — run: pip install python-docx"}

    path = Path(path)
    if not path.exists():
        return {"error": f"Document not found: {path}"}

    # Backup original
    backup_path = path.with_suffix(".docx.kairo_backup")
    shutil.copy2(path, backup_path)

    doc = Document(str(path))
    applied = []
    errors = []

    for op in operations:
        action = op.get("action", "")
        try:
            if action == "insert_after_heading":
                _op_insert_after_heading(doc, op)
                applied.append(op)
            elif action == "insert_paragraph":
                _op_insert_paragraph(doc, op)
                applied.append(op)
            elif action == "replace_paragraph":
                _op_replace_paragraph(doc, op)
                applied.append(op)
            elif action == "append":
                _op_append(doc, op)
                applied.append(op)
            elif action == "insert_table":
                _op_insert_table(doc, op)
                applied.append(op)
            else:
                errors.append(f"Unknown action: {action}")
        except Exception as e:
            errors.append(f"Operation {action} failed: {e}")
            log.error(f"Operation error: {e}\n{traceback.format_exc()}")

    # Atomic write to temp then rename
    tmp_path = path.with_suffix(".docx.tmp")
    doc.save(str(tmp_path))
    try:
        os.replace(str(tmp_path), str(path))
    except PermissionError:
        # Word has the file open/locked
        try:
            tmp_path.unlink(missing_ok=True)
        except Exception:
            pass
        return {
            "error": "Word has this file open. Save and close the document first, then press Alt+M again.",
            "path": str(path),
        }

    # Remove backup on success
    if not errors:
        try:
            backup_path.unlink(missing_ok=True)
        except Exception:
            pass

    return {
        "applied": len(applied),
        "errors": errors,
        "path": str(path),
    }


def _op_insert_after_heading(doc, op: dict):
    """Insert paragraphs after a heading matched by text."""
    from docx.oxml import OxmlElement
    target_heading = op.get("heading_text", "")
    style = op.get("style", "Normal")
    content = op.get("content", "")
    if isinstance(content, str):
        content = [content]

    # Find heading paragraph index
    target_idx = None
    for i, para in enumerate(doc.paragraphs):
        if para.style.name.startswith("Heading") and target_heading.lower() in para.text.lower():
            target_idx = i
            break

    if target_idx is None:
        # Append to end if heading not found
        for text in content:
            if text.strip():
                _add_paragraph_with_style(doc, text, style)
        return

    # Insert after target_idx — find last para before next heading
    insert_after = target_idx
    for j in range(target_idx + 1, len(doc.paragraphs)):
        if doc.paragraphs[j].style.name.startswith("Heading"):
            break
        insert_after = j

    # FIX-5: Insert in FORWARD order by advancing ref_para after each insert
    ref_para = doc.paragraphs[insert_after]._element
    for text in content:
        if text.strip():
            new_para = _make_paragraph(doc, text, style)
            ref_para.addnext(new_para._element)
            ref_para = new_para._element  # advance so next item goes AFTER this one


def _op_insert_paragraph(doc, op: dict):
    """Insert a paragraph at a given index."""
    idx = op.get("index", -1)
    style = op.get("style", "Normal")
    content = op.get("content", "")
    if idx < 0 or idx >= len(doc.paragraphs):
        _add_paragraph_with_style(doc, content, style)
    else:
        ref = doc.paragraphs[idx]._element
        new_para = _make_paragraph(doc, content, style)
        ref.addprevious(new_para._element)


def _op_replace_paragraph(doc, op: dict):
    """Replace the text of paragraph at index."""
    idx = op.get("index", 0)
    content = op.get("content", "")
    style = op.get("style", None)
    if 0 <= idx < len(doc.paragraphs):
        para = doc.paragraphs[idx]
        for run in para.runs:
            run.text = ""
        if para.runs:
            para.runs[0].text = content
        else:
            para.add_run(content)
        if style:
            try:
                para.style = doc.styles[style]
            except Exception:
                pass


def _op_append(doc, op: dict):
    """Append paragraph(s) to end of document."""
    style = op.get("style", "Normal")
    content = op.get("content", "")
    if isinstance(content, str):
        content = [content]
    for text in content:
        if text.strip():
            _add_paragraph_with_style(doc, text, style)


def _op_insert_table(doc, op: dict):
    """Insert a table at end of document."""
    rows_data = op.get("rows", [])
    if not rows_data:
        return
    cols = max(len(r) for r in rows_data)
    table = doc.add_table(rows=0, cols=cols)
    table.style = "Table Grid"
    for row_data in rows_data:
        row = table.add_row()
        for i, cell_text in enumerate(row_data):
            if i < cols:
                row.cells[i].text = str(cell_text)


def _make_paragraph(doc, text: str, style: str):
    """Create a new paragraph element (not added to doc yet)."""
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    para = doc.add_paragraph(text)
    try:
        para.style = doc.styles[style]
    except Exception:
        pass
    # Remove from doc body (we'll insert it manually)
    para._element.getparent().remove(para._element)
    return para


def _add_paragraph_with_style(doc, text: str, style: str):
    """Add paragraph at end of document with given style."""
    para = doc.add_paragraph(text)
    try:
        para.style = doc.styles[style]
    except Exception:
        log.warning(f"Style '{style}' not found in document — using Normal")
    return para


# ─── XLSX Operations ─────────────────────────────────────────────────────────

def _read_xlsx(path: str, active_cell: Optional[str] = None) -> dict:
    """Read Excel file and return surrounding context."""
    try:
        import openpyxl
        from openpyxl.utils import column_index_from_string, get_column_letter
    except ImportError:
        return {"error": "openpyxl not installed — run: pip install openpyxl"}

    wb = openpyxl.load_workbook(path, data_only=False)
    ws = wb.active

    # Parse active cell (e.g. "G5")
    if active_cell:
        col_str = ''.join(c for c in active_cell if c.isalpha())
        row = int(''.join(c for c in active_cell if c.isdigit()))
        col = column_index_from_string(col_str)
    else:
        row, col = 1, 1

    # Extract 10x10 window around active cell
    r_start = max(1, row - 5)
    r_end = min(ws.max_row, row + 5)
    c_start = max(1, col - 5)
    c_end = min(ws.max_column, col + 5)

    grid = []
    for r in range(r_start, r_end + 1):
        row_data = []
        for c in range(c_start, c_end + 1):
            cell = ws.cell(row=r, column=c)
            row_data.append({
                "ref": f"{get_column_letter(c)}{r}",
                "value": str(cell.value) if cell.value is not None else "",
                "formula": str(cell.value) if str(cell.value).startswith("=") else "",
                "is_active": (r == row and c == col),
            })
        grid.append(row_data)

    # Column headers (row 1)
    headers = {}
    for c in range(1, ws.max_column + 1):
        val = ws.cell(row=1, column=c).value
        if val:
            headers[get_column_letter(c)] = str(val)

    # Named ranges (openpyxl 3.x compatible)
    try:
        named_ranges = list(wb.defined_names.keys())
    except Exception:
        named_ranges = []

    return {
        "active_cell": active_cell or "A1",
        "active_row": row,
        "active_col": col,
        "sheet_name": ws.title,
        "sheet_names": wb.sheetnames,
        "grid": grid,
        "headers": headers,
        "named_ranges": named_ranges,
        "max_row": ws.max_row,
        "max_col": ws.max_column,
    }


def _write_xlsx(path: str, operations: list) -> dict:
    """Apply ExcelOperation list to .xlsx file."""
    try:
        import openpyxl
    except ImportError:
        return {"error": "openpyxl not installed — run: pip install openpyxl"}

    path = Path(path)
    if not path.exists():
        return {"error": f"File not found: {path}"}

    backup_path = path.with_suffix(".xlsx.kairo_backup")
    shutil.copy2(path, backup_path)

    wb = openpyxl.load_workbook(str(path))
    ws = wb.active
    applied = []
    errors = []

    for op in operations:
        try:
            cell_ref = op.get("cell", "A1")
            value = op.get("value", "")
            formula = op.get("formula", "")
            write_val = formula if formula else value
            ws[cell_ref] = write_val
            applied.append({"cell": cell_ref, "written": write_val})
        except Exception as e:
            errors.append(f"Excel write error {op}: {e}")

    tmp_path = path.with_suffix(".xlsx.tmp")
    wb.save(str(tmp_path))
    try:
        os.replace(str(tmp_path), str(path))
    except PermissionError:
        try:
            tmp_path.unlink(missing_ok=True)
        except Exception:
            pass
        return {
            "error": "Excel has this file open. Save and close the workbook first, then press Alt+M again.",
            "path": str(path),
        }
    if not errors:
        try:
            backup_path.unlink(missing_ok=True)
        except Exception:
            pass

    return {"applied": len(applied), "cells": applied, "errors": errors}



# ─── PPTX Operations ─────────────────────────────────────────────────────────

def _read_pptx(path: str) -> dict:
    """Read PowerPoint and return slide/shape inventory."""
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError:
        return {"error": "python-pptx not installed — run: pip install python-pptx"}

    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides):
        shapes = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = "\n".join(p.text for p in shape.text_frame.paragraphs)
                shapes.append({
                    "id": shape.shape_id,
                    "name": shape.name,
                    "text": text,
                    "left": shape.left,
                    "top": shape.top,
                })
        # Slide title
        title = slide.shapes.title.text if slide.shapes.title else f"Slide {i+1}"
        slides.append({"index": i, "title": title, "shapes": shapes})

    return {"slides": slides, "slide_count": len(slides)}


def _write_pptx(path: str, operations: list) -> dict:
    """Apply SlideOperation list to .pptx file."""
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError:
        return {"error": "python-pptx not installed — run: pip install python-pptx"}

    path = Path(path)
    if not path.exists():
        return {"error": f"File not found: {path}"}

    backup_path = path.with_suffix(".pptx.kairo_backup")
    shutil.copy2(path, backup_path)

    prs = Presentation(str(path))
    applied = []
    errors = []

    for op in operations:
        try:
            slide_idx = op.get("slide_index", 0)
            shape_id = op.get("shape_id", None)
            bullets = op.get("bullets", [])
            # Enforce ≤7 words per bullet
            bullets = [" ".join(b.split()[:7]) for b in bullets]

            if slide_idx >= len(prs.slides):
                errors.append(f"Slide {slide_idx} out of range")
                continue

            slide = prs.slides[slide_idx]
            target_shape = None
            for shape in slide.shapes:
                if shape_id and shape.shape_id == shape_id:
                    target_shape = shape
                    break
                elif not shape_id and shape.has_text_frame and shape != slide.shapes.title:
                    target_shape = shape

            if not target_shape or not target_shape.has_text_frame:
                if bullets:
                    slide.notes_slide.notes_text_frame.text = "\n".join(bullets)
                    applied.append({"slide": slide_idx, "notes": bullets})
                    continue
                else:
                    errors.append(f"Shape {shape_id} not found on slide {slide_idx}")
                    continue

            tf = target_shape.text_frame

            # FIX-9: Preserve theme font — update existing paragraphs/runs instead of
            # clearing the text frame (tf.clear() destroys run formatting/theme font).
            # Strategy: reuse existing paragraphs where possible, add/remove as needed.
            from pptx.util import Pt

            # Snapshot existing run format from first paragraph (theme font, size, color)
            existing_fmt = None
            if tf.paragraphs and tf.paragraphs[0].runs:
                r0 = tf.paragraphs[0].runs[0]
                existing_fmt = {
                    "bold": r0.font.bold,
                    "size": r0.font.size,
                    "color": r0.font.color.rgb if r0.font.color and r0.font.color.type else None,
                }

            # Clear text content but keep XML structure for theme font
            for para in tf.paragraphs:
                for run in para.runs:
                    run.text = ""

            # Write new bullets — reuse existing <a:p> elements, add/remove as needed
            existing_paras = list(tf.paragraphs)
            for i, bullet_text in enumerate(bullets):
                if i < len(existing_paras):
                    p = existing_paras[i]
                    if p.runs:
                        p.runs[0].text = bullet_text
                        if existing_fmt:
                            if existing_fmt["bold"] is not None:
                                p.runs[0].font.bold = existing_fmt["bold"]
                    else:
                        run = p.add_run()
                        run.text = bullet_text
                else:
                    p = tf.add_paragraph()
                    p.text = bullet_text
                    p.level = 0

            # Remove excess paragraphs beyond bullet count
            for i in range(len(bullets), len(existing_paras)):
                para_elem = existing_paras[i]._p
                para_elem.getparent().remove(para_elem)

            applied.append({"slide": slide_idx, "shape": shape_id, "bullets": bullets})
        except Exception as e:
            errors.append(f"PPTX write error: {e}")

    tmp_path = path.with_suffix(".pptx.tmp")
    prs.save(str(tmp_path))
    try:
        os.replace(str(tmp_path), str(path))
    except PermissionError:
        try:
            tmp_path.unlink(missing_ok=True)
        except Exception:
            pass
        return {
            "error": "PowerPoint has this file open. Save and close the presentation first, then press Alt+M again.",
            "path": str(path),
        }
    if not errors:
        try:
            backup_path.unlink(missing_ok=True)
        except Exception:
            pass

    return {"applied": len(applied), "errors": errors}



# ─── Context extractor (for LLM prompt building) ─────────────────────────────

def _extract_context(path: str, active_cell: Optional[str] = None) -> dict:
    """
    Route to correct reader based on file extension.
    Returns unified context dict for LLM prompt building.
    """
    ext = Path(path).suffix.lower()
    if ext == ".docx":
        return _read_docx(path)
    elif ext in (".xlsx", ".xlsm"):
        return _read_xlsx(path, active_cell)
    elif ext == ".pptx":
        return _read_pptx(path)
    elif ext in (".txt", ".md"):
        text = Path(path).read_text(encoding="utf-8", errors="replace")
        return {"full_text": text, "format": ext.lstrip(".")}
    elif ext == ".pdf":
        return _read_pdf(path)
    else:
        return {"error": f"Unsupported format: {ext}", "ext": ext}


def _read_pdf(path: str) -> dict:
    """3-tier PDF extraction: fitz → docling → mineru_vlm."""
    from sidecar.parsers.pdf_parser import parse_pdf
    try:
        res = parse_pdf(path)
        # Convert structured paragraphs to flat full_text if expected by old callers
        if "full_text" not in res:
            res["full_text"] = "\n".join(p["text"] for p in res.get("paragraphs", []))
        return res
    except Exception as e:
        log.error(f"parse_pdf failed: {e}")
        return {"error": f"PDF parsing failed: {e}", "format": "pdf"}



# ─── TCP Server ──────────────────────────────────────────────────────────────

HOST = "127.0.0.1"
PORT = 7438


async def _check_already_running() -> bool:
    """Check if a sidecar is already running on our port."""
    import asyncio
    try:
        reader, writer = await asyncio.wait_for(
            asyncio.open_connection(HOST, PORT), timeout=1.0
        )
        # Send ping
        writer.write(b'{"id":"startup-check","action":"ping"}\n')
        await writer.drain()
        line = await asyncio.wait_for(reader.readline(), timeout=1.0)
        writer.close()
        data = json.loads(line.decode().strip())
        if data.get("ok"):
            log.info(f"✅ Sidecar already running on port {PORT} — this instance will exit")
            return True
    except Exception:
        pass
    return False


async def handle_request(data: dict) -> dict:
    req_id = data.get("id", "unknown")
    action = data.get("action", "")
    path = data.get("path", "")
    payload = data.get("payload", {})

    log.info(f"Request [{req_id}] action={action} path={path}")

    try:
        if action == "ping":
            return {"id": req_id, "ok": True, "data": {"pong": True, "version": "1.0.0"}}

        elif action == "read_docx":
            data_out = _read_docx(path)
            return {"id": req_id, "ok": "error" not in data_out, "data": data_out}

        elif action == "write_docx":
            ops = payload.get("operations", [])
            data_out = _write_docx(path, ops)
            return {"id": req_id, "ok": "error" not in data_out, "data": data_out}

        elif action == "read_xlsx":
            active_cell = payload.get("active_cell")
            data_out = _read_xlsx(path, active_cell)
            return {"id": req_id, "ok": "error" not in data_out, "data": data_out}

        elif action == "write_xlsx":
            ops = payload.get("operations", [])
            data_out = _write_xlsx(path, ops)
            return {"id": req_id, "ok": "error" not in data_out, "data": data_out}

        elif action == "read_pptx":
            data_out = _read_pptx(path)
            return {"id": req_id, "ok": "error" not in data_out, "data": data_out}

        elif action == "write_pptx":
            ops = payload.get("operations", [])
            data_out = _write_pptx(path, ops)
            return {"id": req_id, "ok": "error" not in data_out, "data": data_out}

        elif action == "extract_context":
            active_cell = payload.get("active_cell")
            data_out = _extract_context(path, active_cell)
            return {"id": req_id, "ok": "error" not in data_out, "data": data_out}

        # ─── Domain 5: Design & Figma Actions ───────────────────────────────
        elif action == "figma_create":
            from sidecar.parsers.design_bridge import UnifiedDesignBridge
            bridge = UnifiedDesignBridge(offline_mode=True)
            node_type = payload.get("node_type", "FRAME").upper()
            name = payload.get("name", "Unnamed Node")
            x = payload.get("x", 0)
            y = payload.get("y", 0)
            width = payload.get("width", 100)
            height = payload.get("height", 100)
            parent_id = payload.get("parent_id", "canvas-root")
            characters = payload.get("characters", "")
            fontSize = payload.get("fontSize", 14)
            color_hex = payload.get("color_hex")
            layout_mode = payload.get("layout_mode")
            spacing = payload.get("spacing", 10)
            padding_tb = payload.get("padding_tb", 0)
            padding_lr = payload.get("padding_lr", 0)

            res = {}
            if node_type == "FRAME":
                res = bridge.figma.create_frame(name, x, y, width, height, parent_id)
            elif node_type == "TEXT":
                res = bridge.figma.create_text_node(name, characters, fontSize, parent_id)
            elif node_type == "RECTANGLE":
                res = bridge.figma.create_rectangle(name, x, y, width, height, parent_id)
            elif node_type == "COMPONENT":
                res = bridge.figma.create_component(name, parent_id)
            elif node_type == "SECTION":
                res = bridge.figma.create_section(name, parent_id)
            else:
                res = {"ok": False, "error": f"Unsupported node type: {node_type}"}

            if res.get("ok") and color_hex:
                bridge.figma.set_fills(res["node_id"], color_hex)
            if res.get("ok") and layout_mode and node_type in ("FRAME", "COMPONENT"):
                bridge.figma.set_auto_layout(res["node_id"], layout_mode, spacing, padding_tb, padding_lr)

            return {"id": req_id, "ok": res.get("ok", False), "data": res}

        elif action == "design_ghost_write":
            from sidecar.parsers.design_bridge import UnifiedDesignBridge
            bridge = UnifiedDesignBridge(offline_mode=True)
            window_title = payload.get("window_title", "")
            tool = bridge.detect_active_design_tool(window_title)
            
            res = {"tool_detected": tool}
            if tool == "figma":
                # Create a sample text node to represent ghost design action
                node = bridge.figma.create_text_node("Ghost Written Element", "Generated via Kairo Swarm", 16)
                res.update(node)
            elif tool == "penpot":
                svg = payload.get("svg", "<svg><rect width='100' height='100' fill='#6140f0'/></svg>")
                penpot_res = bridge.penpot.draw_svg(svg)
                res.update(penpot_res)
            elif tool == "tldraw":
                shapes = payload.get("shapes", [{"type": "geo", "x": 100, "y": 100, "props": {"w": 100, "h": 50, "text": "Ghost"}}])
                created = []
                for s in shapes:
                    s_type = s.get("type", "geo")
                    s_x = s.get("x", 0.0)
                    s_y = s.get("y", 0.0)
                    s_props = s.get("props", {})
                    created.append(bridge.tldraw.create_shape(s_type, s_x, s_y, s_props))
                res["created_shapes"] = created
            elif tool == "openpencil":
                stroke = payload.get("stroke", {"points": [(0,0), (10,10)]})
                op_res = bridge.openpencil.record_stroke(stroke)
                res.update(op_res)
            elif tool == "frameground":
                filename = payload.get("filename", "frameground_canvas.html")
                html = payload.get("html", "<div>Frameground</div>")
                path_out = bridge.frameground.create_canvas(filename, html)
                res["canvas_path"] = path_out
            else:
                res["error"] = "No active design tool matched window title."
                res["ok"] = False
                return {"id": req_id, "ok": False, "data": res}
            
            res["ok"] = True
            return {"id": req_id, "ok": True, "data": res}

        elif action == "generate_design_asset":
            from sidecar.parsers.design_bridge import UnifiedDesignBridge
            bridge = UnifiedDesignBridge(offline_mode=True)
            prompt = payload.get("prompt", "")
            style = payload.get("style", "default")
            output_path = payload.get("output_path")
            res = bridge.comfyui.generate_asset(prompt, style, output_path)
            return {"id": req_id, "ok": res.get("ok", False), "data": res}

        elif action == "tldraw_canvas":
            from sidecar.parsers.design_bridge import UnifiedDesignBridge
            bridge = UnifiedDesignBridge(offline_mode=True)
            operation = payload.get("operation", "get_shapes")
            
            if operation == "create_shape":
                shape_type = payload.get("shape_type", "geo")
                x = payload.get("x", 0.0)
                y = payload.get("y", 0.0)
                props = payload.get("props", {})
                res = bridge.tldraw.create_shape(shape_type, x, y, props)
            elif operation == "update_shape":
                shape_id = payload.get("shape_id", "")
                x = payload.get("x")
                y = payload.get("y")
                props = payload.get("props")
                res = bridge.tldraw.update_shape(shape_id, x, y, props)
            elif operation == "delete_shape":
                shape_id = payload.get("shape_id", "")
                res = bridge.tldraw.delete_shape(shape_id)
            elif operation == "draw_flowchart":
                nodes = payload.get("nodes", [])
                edges = payload.get("edges", [])
                res = bridge.tldraw.draw_flowchart(nodes, edges)
            else:
                res = {"shapes": bridge.tldraw.get_canvas_shapes(), "ok": True}
                
            return {"id": req_id, "ok": res.get("ok", True), "data": res}

        elif action == "extract_design_code":
            from sidecar.parsers.design_bridge import UnifiedDesignBridge
            bridge = UnifiedDesignBridge(offline_mode=True)
            root_id = payload.get("root_id", "canvas-root")
            html = bridge.transpile_figma_to_tailwind(root_id)
            return {"id": req_id, "ok": True, "data": {"html": html}}

        elif action == "learn_design_preference":
            from sidecar.parsers.design_bridge import UnifiedDesignBridge
            bridge = UnifiedDesignBridge(offline_mode=True)
            tool = payload.get("tool", "default")
            key = payload.get("key", "")
            value = payload.get("value")
            bridge.memory.learn_preference(tool, key, value)
            return {"id": req_id, "ok": True}

        elif action == "get_design_preference":
            from sidecar.parsers.design_bridge import UnifiedDesignBridge
            bridge = UnifiedDesignBridge(offline_mode=True)
            tool = payload.get("tool", "default")
            key = payload.get("key", "")
            fallback = payload.get("fallback")
            val = bridge.memory.get_preference(tool, key, fallback)
            return {"id": req_id, "ok": True, "data": val}

        else:
            return {"id": req_id, "ok": False, "error": f"Unknown action: {action}"}

    except Exception as e:
        log.error(f"Handler error for [{req_id}]: {e}\n{traceback.format_exc()}")
        return {"id": req_id, "ok": False, "error": str(e)}


async def handle_client(reader: asyncio.StreamReader, writer: asyncio.StreamWriter):
    addr = writer.get_extra_info("peername")
    log.info(f"Client connected: {addr}")
    try:
        while True:
            line = await reader.readline()
            if not line:
                break
            try:
                request = json.loads(line.decode("utf-8").strip())
                response = await handle_request(request)
                writer.write((json.dumps(response) + "\n").encode("utf-8"))
                await writer.drain()
            except json.JSONDecodeError as e:
                err = json.dumps({"ok": False, "error": f"Invalid JSON: {e}"}) + "\n"
                writer.write(err.encode("utf-8"))
                await writer.drain()
    except (asyncio.IncompleteReadError, ConnectionResetError):
        pass
    finally:
        writer.close()
        log.info(f"Client disconnected: {addr}")


async def main():
    log.info(f"Kairo sidecar v1.0 starting on {HOST}:{PORT}")
    # Check if another sidecar instance is already running
    if await _check_already_running():
        # Already running — this daemon-spawned copy can exit; the existing one serves
        sys.exit(0)
    try:
        server = await asyncio.start_server(handle_client, HOST, PORT)
    except OSError as e:
        if "address already in use" in str(e).lower() or e.errno == 10048:
            log.info(f"Port {PORT} already in use — another sidecar is running, exiting")
            sys.exit(0)
        raise
    async with server:
        log.info(f"✅ Sidecar ready on {HOST}:{PORT}")
        await server.serve_forever()


if __name__ == "__main__":
    asyncio.run(main())
