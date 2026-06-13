"""
sidecar/creators/pptx_creator.py — Kairo Phantom Create-From-Scratch PPTX
==========================================================================
Creates a new .pptx presentation from a structured content dict and opens it
in PowerPoint via os.startfile().

Usage
-----
    from sidecar.creators.pptx_creator import PptxCreator

    creator = PptxCreator()
    path = creator.create_and_open({
        "title": "Kairo Phantom — Investor Pitch",
        "author": "Kairo Phantom",
        "theme": "dark",  # "dark" | "light" | "minimal" (future)
        "slides": [
            {
                "layout": "title",
                "title": "Kairo Phantom",
                "subtitle": "AI that writes into your documents. 100% offline."
            },
            {
                "layout": "content",
                "title": "The Problem",
                "bullets": [
                    "AI assistants suggest — they don't execute",
                    "Cloud-only: your documents leave your machine",
                    "No memory: same errors every session"
                ]
            },
            {
                "layout": "two_column",
                "title": "Before vs After Kairo",
                "left": ["Clipboard paste → wrong style", "No memory", "Cloud required"],
                "right": ["python-docx write-back → correct style", "MemMachine recall", "100% offline"]
            },
            {
                "layout": "blank",
                "title": "Thank You",
                "body": "github.com/KairoPhantom/Kairo-Phantom"
            }
        ]
    })
"""

import os
import logging
from pathlib import Path
from typing import Dict, Any, Optional, List

log = logging.getLogger("kairo-sidecar.pptx_creator")

KAIRO_DOCS_DIR = Path.home() / "Documents" / "Kairo"


class PptxCreator:
    """
    Creates new .pptx presentations from a structured content dict.
    Saves to ~/Documents/Kairo/ and opens in PowerPoint via os.startfile().
    """

    def create(
        self,
        content: Dict[str, Any],
        output_path: Optional[str] = None,
    ) -> str:
        """
        Creates a .pptx from structured content.

        Parameters
        ----------
        content     : Dict with keys:
                      - title (str): Presentation title
                      - author (str): Author name
                      - slides (list): Each slide dict may have:
                          - layout (str): "title" | "content" | "two_column" | "blank"
                          - title (str): Slide title
                          - subtitle (str): Subtitle (title slides)
                          - bullets (list[str]): Bullet points (content slides)
                          - body (str): Plain body text (blank slides)
                          - left (list[str]): Left column bullets (two_column)
                          - right (list[str]): Right column bullets (two_column)
        output_path : Full path to save. Auto-generated if None.

        Returns
        -------
        str — absolute path of the created file.
        """
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()

        # Set slide dimensions to widescreen 16:9
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        title = content.get("title", "Untitled Presentation")
        slides_data = content.get("slides", [])

        # If no slides provided, add a default title slide
        if not slides_data:
            slides_data = [{"layout": "title", "title": title, "subtitle": ""}]

        for slide_data in slides_data:
            layout_name = slide_data.get("layout", "content")
            slide_title = slide_data.get("title", "")

            if layout_name == "title":
                # Use title slide layout (index 0)
                slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                if slide.shapes.title:
                    slide.shapes.title.text = slide_title
                subtitle_text = slide_data.get("subtitle", "")
                # Subtitle placeholder is usually index 1
                if len(slide.placeholders) > 1 and subtitle_text:
                    slide.placeholders[1].text = subtitle_text

            elif layout_name == "content":
                # Use content/bullet layout (index 1)
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                if slide.shapes.title:
                    slide.shapes.title.text = slide_title
                bullets = slide_data.get("bullets", [])
                body_text = slide_data.get("body", "")
                if len(slide.placeholders) > 1:
                    tf = slide.placeholders[1].text_frame
                    tf.clear()
                    for i, bullet in enumerate(bullets):
                        if i == 0:
                            tf.paragraphs[0].text = bullet
                            tf.paragraphs[0].level = 0
                        else:
                            p = tf.add_paragraph()
                            p.text = bullet
                            p.level = 0
                    if body_text and not bullets:
                        tf.paragraphs[0].text = body_text

            elif layout_name == "two_column":
                # Use blank layout and manually add two text boxes
                slide_layout = prs.slide_layouts[6]  # blank
                slide = prs.slides.add_slide(slide_layout)

                # Add title text box
                txBox = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.2), Inches(12.33), Inches(1.0)
                )
                tf = txBox.text_frame
                tf.text = slide_title
                tf.paragraphs[0].runs[0].font.size = Pt(28)
                tf.paragraphs[0].runs[0].font.bold = True

                # Left column
                left_items = slide_data.get("left", [])
                left_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.5)
                )
                ltf = left_box.text_frame
                ltf.word_wrap = True
                for i, item in enumerate(left_items):
                    if i == 0:
                        ltf.paragraphs[0].text = f"• {item}"
                    else:
                        p = ltf.add_paragraph()
                        p.text = f"• {item}"

                # Right column
                right_items = slide_data.get("right", [])
                right_box = slide.shapes.add_textbox(
                    Inches(7.0), Inches(1.4), Inches(5.8), Inches(5.5)
                )
                rtf = right_box.text_frame
                rtf.word_wrap = True
                for i, item in enumerate(right_items):
                    if i == 0:
                        rtf.paragraphs[0].text = f"• {item}"
                    else:
                        p = rtf.add_paragraph()
                        p.text = f"• {item}"

            else:
                # Blank layout with optional title + body
                slide_layout = prs.slide_layouts[6]  # blank
                slide = prs.slides.add_slide(slide_layout)
                body = slide_data.get("body", slide_data.get("subtitle", ""))

                if slide_title:
                    txBox = slide.shapes.add_textbox(
                        Inches(0.5), Inches(0.2), Inches(12.33), Inches(1.0)
                    )
                    tf = txBox.text_frame
                    tf.text = slide_title
                    if tf.paragraphs:
                        tf.paragraphs[0].runs[0].font.size = Pt(28)
                        tf.paragraphs[0].runs[0].font.bold = True

                if body:
                    bBox = slide.shapes.add_textbox(
                        Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.5)
                    )
                    btf = bBox.text_frame
                    btf.word_wrap = True
                    btf.text = body

        # Set core properties
        try:
            prs.core_properties.author = content.get("author", "Kairo Phantom")
            prs.core_properties.title = title
        except Exception:
            pass

        # Resolve output path
        if not output_path:
            KAIRO_DOCS_DIR.mkdir(parents=True, exist_ok=True)
            safe_title = "".join(
                c if c.isalnum() or c in " _-" else "_" for c in title
            ).strip()[:50] or "Presentation"
            output_path = str(KAIRO_DOCS_DIR / f"{safe_title}.pptx")

        prs.save(output_path)
        log.info(f"PptxCreator: saved to {output_path}")
        return output_path

    def create_and_open(self, content: Dict[str, Any]) -> str:
        """Creates a .pptx and opens it in PowerPoint."""
        path = self.create(content)
        try:
            os.startfile(path)
            log.info(f"PptxCreator: opened {path}")
        except Exception as e:
            log.warning(f"PptxCreator: could not auto-open {path}: {e}")
        return path
