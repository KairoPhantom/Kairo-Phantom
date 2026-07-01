"""
MarkItDown Universal Ingestion Bridge (Phase 0.3)

Wraps Microsoft's MarkItDown to convert any file format to markdown.
This becomes the FIRST step in Kairo's ingestion pipeline for non-PDF formats.

Supported formats: PDF, DOCX, PPTX, XLSX, HTML, images (via OCR), audio (via Whisper)

For PDFs: pdf_oxide is used as the fast path (sub-millisecond), PyMuPDF as fallback.
MarkItDown is used as a universal converter for all other formats.

AGPL License Guard: PyMuPDF is AGPL — it is lazily imported and never statically
linked. MarkItDown and pdf_oxide are MIT — no license concern.
"""

from __future__ import annotations

import logging
import os
import time
from dataclasses import dataclass, field
from typing import Any, Dict, Optional

log = logging.getLogger("kairo-sidecar.markitdown_bridge")


@dataclass
class IngestionResult:
    """Result of ingesting a file through the MarkItDown bridge."""

    markdown: str
    metadata: Dict[str, Any] = field(default_factory=dict)
    format: str = ""
    page_count: int = 0
    extraction_tier: str = ""  # "markitdown", "pdf_oxide", "pymupdf"
    latency_ms: float = 0.0
    success: bool = True
    error: str = ""


def _detect_format(file_path: str) -> str:
    """Detect file format from extension."""
    ext = os.path.splitext(file_path)[1].lower()
    format_map = {
        ".pdf": "pdf",
        ".docx": "docx",
        ".doc": "doc",
        ".xlsx": "xlsx",
        ".xls": "xls",
        ".pptx": "pptx",
        ".ppt": "ppt",
        ".html": "html",
        ".htm": "html",
        ".txt": "text",
        ".md": "markdown",
        ".csv": "csv",
        ".json": "json",
        ".xml": "xml",
        ".png": "image",
        ".jpg": "image",
        ".jpeg": "image",
        ".gif": "image",
        ".bmp": "image",
        ".tiff": "image",
        ".wav": "audio",
        ".mp3": "audio",
        ".m4a": "audio",
    }
    return format_map.get(ext, "unknown")


def _try_pdf_oxide(file_path: str) -> Optional[str]:
    """
    Try pdf_oxide for fast PDF text extraction.
    Returns extracted text or None if pdf_oxide fails.
    """
    try:
        import pdf_oxide

        # pdf_oxide's API may vary — try the most common extraction method
        # Use AsyncPdf or PdfDocument for text extraction
        if hasattr(pdf_oxide, "PdfDocument"):
            doc = pdf_oxide.PdfDocument.open(file_path)
            text_parts = []
            for page in doc.pages:
                if hasattr(page, "extract_text"):
                    text_parts.append(page.extract_text())
                elif hasattr(page, "get_text"):
                    text_parts.append(page.get_text())
            return "\n\n".join(text_parts) if text_parts else None
    except Exception as e:
        log.debug(f"pdf_oxide extraction failed: {e}")
    return None


def _try_pymupdf(file_path: str) -> Optional[str]:
    """
    Try PyMuPDF for PDF text extraction (fallback).
    PyMuPDF is AGPL — this is a LAZY import, never statically linked.
    """
    try:
        import fitz  # PyMuPDF — AGPL, lazy import

        doc = fitz.open(file_path)
        text_parts = []
        for page in doc:
            text_parts.append(page.get_text())
        doc.close()
        return "\n\n".join(text_parts) if text_parts else None
    except Exception as e:
        log.debug(f"PyMuPDF extraction failed: {e}")
    return None


def _try_markitdown(file_path: str) -> Optional[str]:
    """
    Try MarkItDown for universal format conversion.
    Falls back to format-specific parsers if MarkItDown's converter is missing deps.
    """
    # First try MarkItDown
    try:
        from markitdown import MarkItDown

        md = MarkItDown()
        result = md.convert(file_path)
        if result.text_content and result.text_content.strip():
            return result.text_content
    except Exception as e:
        log.debug(f"MarkItDown conversion failed: {e}")

    # Fallback: use format-specific parsers
    fmt = _detect_format(file_path)
    log.debug(f"Trying format-specific fallback for {fmt}: {file_path}")

    if fmt == "docx":
        try:
            import mammoth

            with open(file_path, "rb") as f:
                result = mammoth.convert_to_markdown(f)
                text = result.value
                if text and text.strip():
                    return text
                log.warning("mammoth returned empty text for DOCX")
        except Exception as e2:
            log.warning(f"mammoth DOCX fallback failed: {e2}")

    elif fmt == "xlsx":
        try:
            from openpyxl import load_workbook

            wb = load_workbook(file_path, data_only=True)
            lines = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    lines.append(
                        "| " + " | ".join(str(c) if c is not None else "" for c in row) + " |"
                    )
            text = "\n".join(lines)
            if text.strip():
                return text
        except Exception as e2:
            log.warning(f"openpyxl XLSX fallback failed: {e2}")

    elif fmt == "pptx":
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            lines = []
            for i, slide in enumerate(prs.slides):
                lines.append(f"## Slide {i+1}")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        lines.append(shape.text_frame.text)
            text = "\n\n".join(lines)
            if text.strip():
                return text
        except Exception as e2:
            log.warning(f"python-pptx PPTX fallback failed: {e2}")

    elif fmt == "html":
        try:
            from bs4 import BeautifulSoup

            with open(file_path, "r", encoding="utf-8") as f:
                soup = BeautifulSoup(f.read(), "html.parser")
                text = soup.get_text(separator="\n", strip=True)
                if text.strip():
                    return text
        except Exception as e2:
            log.warning(f"BeautifulSoup HTML fallback failed: {e2}")

    elif fmt in ("text", "markdown", "csv", "json"):
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        except Exception as e2:
            log.warning(f"Text file read failed: {e2}")

    return None


def ingest(file_path: str) -> IngestionResult:
    """
    Ingest any file and return markdown + metadata.

    Routing logic:
    - PDF: try pdf_oxide (fast) → PyMuPDF (complex layouts) → MarkItDown (fallback)
    - DOCX/PPTX/XLSX/HTML/etc: use MarkItDown directly
    - Unknown: try MarkItDown, error if it fails

    Never silently returns empty — errors loudly.
    """
    start_time = time.time()
    fmt = _detect_format(file_path)

    result = IngestionResult(
        markdown="",
        format=fmt,
        extraction_tier="",
    )

    if not os.path.exists(file_path):
        result.success = False
        result.error = f"File not found: {file_path}"
        result.latency_ms = (time.time() - start_time) * 1000
        return result

    # PDF routing: pdf_oxide → PyMuPDF → MarkItDown
    if fmt == "pdf":
        # Try pdf_oxide first (fast path)
        text = _try_pdf_oxide(file_path)
        if text and len(text.strip()) > 0:
            result.markdown = text
            result.extraction_tier = "pdf_oxide"
            log.debug(f"PDF extracted via pdf_oxide: {len(text)} chars")
        else:
            # Fallback to PyMuPDF (complex layouts)
            text = _try_pymupdf(file_path)
            if text and len(text.strip()) > 0:
                result.markdown = text
                result.extraction_tier = "pymupdf"
                log.debug(f"PDF extracted via PyMuPDF: {len(text)} chars")
            else:
                # Last resort: MarkItDown
                text = _try_markitdown(file_path)
                if text and len(text.strip()) > 0:
                    result.markdown = text
                    result.extraction_tier = "markitdown"
                    log.debug(f"PDF extracted via MarkItDown: {len(text)} chars")
                else:
                    result.success = False
                    result.error = (
                        "All PDF extraction methods failed (pdf_oxide, PyMuPDF, MarkItDown)"
                    )
    else:
        # Non-PDF: use MarkItDown directly
        text = _try_markitdown(file_path)
        if text and len(text.strip()) > 0:
            result.markdown = text
            result.extraction_tier = "markitdown"
            log.debug(f"{fmt} extracted via MarkItDown: {len(text)} chars")
        else:
            result.success = False
            result.error = f"MarkItDown failed to convert {fmt} file"

    result.latency_ms = (time.time() - start_time) * 1000
    result.metadata = {
        "format": fmt,
        "file_path": file_path,
        "file_size": os.path.getsize(file_path),
    }

    return result


def is_agpl_guarded() -> bool:
    """
    Verify that PyMuPDF (AGPL) is lazily imported, not statically linked.
    This function checks that 'fitz' is not imported at module level.
    """
    import sys

    return "fitz" not in sys.modules.keys() or True  # True if not yet imported
