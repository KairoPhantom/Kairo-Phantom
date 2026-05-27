"""PPTX MCP Bridge: Programmatic PowerPoint for Kairo Phantom."""
import subprocess
import json
import os
import tempfile
import shutil
import logging
from pathlib import Path
from typing import List, Dict, Any, Optional

log = logging.getLogger("kairo-sidecar.pptx_mcp_bridge")

class PptxMcpBridge:
    """Provides PowerPoint creation/manipulation via MCP server with python-pptx fallbacks."""

    def __init__(self, mcp_command: str = "uvx --from office-powerpoint-mcp-server ppt_mcp_server"):
        self.mcp_command = mcp_command
        # We can check if MCP server is available/operational, otherwise fallback directly
        self.use_fallback = True # Always use robust python-pptx local implementation to guarantee offline reliability and avoid subprocess overhead

    # ── Presentation Management ──

    def create_presentation(self, title: str = None) -> str:
        """Create a new presentation. Returns presentation ID (file path)."""
        if self.use_fallback:
            from pptx import Presentation
            prs = Presentation()
            # Set default slide size to 16:9 widescreen
            prs.slide_width = 12192000 # 13.33 inches
            prs.slide_height = 6858000 # 7.5 inches
            
            if title:
                # Add title slide layout
                slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                title_shape = slide.shapes.title
                if title_shape:
                    title_shape.text = title
            
            fd, path = tempfile.mkstemp(suffix=".pptx", prefix="kairo_pres_")
            os.close(fd)
            prs.save(path)
            return path
        
        res = self._call_tool("create_presentation", {"title": title})
        return res.get("presentation_id")

    def create_from_template(self, template_path: str) -> str:
        """Create presentation from template preserving theme."""
        if self.use_fallback:
            from pptx import Presentation
            if not os.path.exists(template_path):
                raise FileNotFoundError(f"Template path {template_path} not found")
            prs = Presentation(template_path)
            fd, path = tempfile.mkstemp(suffix=".pptx", prefix="kairo_pres_")
            os.close(fd)
            prs.save(path)
            return path

        res = self._call_tool("create_presentation_from_template", {"template_path": template_path})
        return res.get("presentation_id")

    def get_presentation_info(self, pres_id: str) -> dict:
        """Get comprehensive presentation information."""
        if self.use_fallback:
            from pptx import Presentation
            prs = Presentation(pres_id)
            return {
                "slide_count": len(prs.slides),
                "theme_name": "Default"
            }
        return self._call_tool("get_presentation_info", {"presentation_id": pres_id})

    # ── Slide Operations ──

    def add_slide(self, pres_id: str, title: str = None,
                  content: str = None, layout_index: int = 1,
                  background_style: str = None) -> dict:
        """Add a slide with optional background styling."""
        if self.use_fallback:
            from pptx import Presentation
            prs = Presentation(pres_id)
            
            # Select layout_index safely
            idx = layout_index if layout_index < len(prs.slide_layouts) else 1
            layout = prs.slide_layouts[idx]
            slide = prs.slides.add_slide(layout)
            
            if title and slide.shapes.title:
                slide.shapes.title.text = title
                # Format slide title to Segoe UI
                for p in slide.shapes.title.text_frame.paragraphs:
                    p.font.name = "Segoe UI"
            
            if content:
                # Find body placeholder or add a text box
                body_shape = None
                for shape in slide.shapes:
                    if shape.is_placeholder and shape.placeholder_format.idx == 1:
                        body_shape = shape
                        break
                if not body_shape:
                    for shape in slide.shapes:
                        if shape.has_text_frame and shape != slide.shapes.title:
                            body_shape = shape
                            break
                if body_shape:
                    body_shape.text_frame.text = content
                    for p in body_shape.text_frame.paragraphs:
                        p.font.name = "Segoe UI"
                else:
                    # add text box
                    from pptx.util import Inches
                    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(4))
                    tf = txBox.text_frame
                    tf.text = content
                    for p in tf.paragraphs:
                        p.font.name = "Segoe UI"
                        
            prs.save(pres_id)
            return {"slide_index": len(prs.slides) - 1, "ok": True}

        return self._call_tool("add_slide", {
            "presentation_id": pres_id,
            "title": title,
            "content": content,
            "layout_index": layout_index,
            "background_style": background_style
        })

    def populate_placeholder(self, pres_id: str, slide_index: int,
                             placeholder_idx: int, text: str) -> dict:
        """Populate placeholder with text."""
        if self.use_fallback:
            from pptx import Presentation
            prs = Presentation(pres_id)
            if slide_index >= len(prs.slides):
                return {"ok": False, "error": f"Slide index {slide_index} out of bounds"}
            slide = prs.slides[slide_index]
            
            target_shape = None
            for shape in slide.shapes:
                if shape.is_placeholder and shape.placeholder_format.idx == placeholder_idx:
                    target_shape = shape
                    break
            
            if not target_shape:
                # Try finding any shape with text frame by index
                shapes_with_tf = [s for s in slide.shapes if s.has_text_frame]
                if placeholder_idx < len(shapes_with_tf):
                    target_shape = shapes_with_tf[placeholder_idx]
                elif shapes_with_tf:
                    target_shape = shapes_with_tf[0]
                else:
                    from pptx.util import Inches
                    target_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(4))
            
            if target_shape:
                target_shape.text_frame.text = text
                for p in target_shape.text_frame.paragraphs:
                    p.font.name = "Segoe UI"
                prs.save(pres_id)
                return {"ok": True}
            return {"ok": False, "error": f"Placeholder {placeholder_idx} not found"}

        return self._call_tool("populate_placeholder", {
            "presentation_id": pres_id,
            "slide_index": slide_index,
            "placeholder_index": placeholder_idx,
            "text": text
        })

    def add_bullet_points(self, pres_id: str, slide_index: int,
                          bullets: list, placeholder_idx: int = 1) -> dict:
        """Add formatted bullet points. Max 5 bullets, <= 7 words per bullet."""
        # Clean bullets to max 5 bullets, and max 7 words each
        cleaned_bullets = []
        for b in bullets[:5]:
            words = b.strip().split()
            cleaned_bullets.append(" ".join(words[:7]))

        if self.use_fallback:
            from pptx import Presentation
            prs = Presentation(pres_id)
            if slide_index >= len(prs.slides):
                return {"ok": False, "error": f"Slide index {slide_index} out of bounds"}
            slide = prs.slides[slide_index]
            
            # Find appropriate text shape
            target_shape = None
            for shape in slide.shapes:
                if shape.is_placeholder and shape.placeholder_format.idx == placeholder_idx:
                    target_shape = shape
                    break
            if not target_shape:
                # Try body placeholder (idx=1) or title (idx=0)
                for shape in slide.shapes:
                    if shape.is_placeholder and shape.placeholder_format.idx == 1:
                        target_shape = shape
                        break
            if not target_shape:
                # Find any text frame
                for shape in slide.shapes:
                    if shape.has_text_frame and shape != slide.shapes.title:
                        target_shape = shape
                        break
            
            if not target_shape:
                from pptx.util import Inches
                target_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(4))

            tf = target_shape.text_frame
            tf.clear() # clear text
            
            for i, b_text in enumerate(cleaned_bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = b_text
                p.level = 0
                p.font.name = "Segoe UI"
                
            prs.save(pres_id)
            return {"ok": True}

        return self._call_tool("add_bullet_points", {
            "presentation_id": pres_id,
            "slide_index": slide_index,
            "bullets": cleaned_bullets,
            "placeholder_index": placeholder_idx
        })

    # ── Visual Elements ──

    def add_image_to_slide(self, pres_id: str, slide_index: int,
                           image_path: str, left: float = None,
                           top: float = None, width: float = None,
                           height: float = None) -> dict:
        """Add image to slide with positioning."""
        if self.use_fallback:
            from pptx import Presentation
            from pptx.util import Inches
            prs = Presentation(pres_id)
            if slide_index >= len(prs.slides):
                return {"ok": False, "error": f"Slide index {slide_index} out of bounds"}
            slide = prs.slides[slide_index]
            
            # Default placement: center-right
            l = Inches(left) if left is not None else Inches(7.5)
            t = Inches(top) if top is not None else Inches(2.0)
            w = Inches(width) if width is not None else Inches(5.0)
            h = Inches(height) if height is not None else Inches(4.0)
            
            if os.path.exists(image_path):
                slide.shapes.add_picture(image_path, l, t, width=w, height=h)
                prs.save(pres_id)
                return {"ok": True}
            return {"ok": False, "error": f"Image path {image_path} not found"}

        return self._call_tool("add_image_to_slide", {
            "presentation_id": pres_id,
            "slide_index": slide_index,
            "image_path": image_path,
            "left": left,
            "top": top,
            "width": width,
            "height": height
        })

    def create_chart_slide(self, pres_id: str, chart_type: str,
                           data: list, title: str = None) -> dict:
        """Create data‑driven chart."""
        if self.use_fallback:
            from pptx import Presentation
            from pptx.chart.data import CategoryChartData
            from pptx.enum.chart import XL_CHART_TYPE
            from pptx.util import Inches
            prs = Presentation(pres_id)
            slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title Only or blank
            
            if title and slide.shapes.title:
                slide.shapes.title.text = title
                slide.shapes.title.text_frame.paragraphs[0].font.name = "Segoe UI"
                
            chart_data = CategoryChartData()
            # data schema: [{"category": "A", "values": [10, 20]}, ...] or similar. Let's assume list of dicts.
            categories = [d.get("category", "") for d in data]
            chart_data.categories = categories
            
            # Series names
            if data and "values" in data[0]:
                series_count = len(data[0]["values"])
                for s in range(series_count):
                    series_vals = [d["values"][s] for d in data]
                    chart_data.add_series(f"Series {s+1}", tuple(series_vals))
            else:
                series_vals = [d.get("value", 0) for d in data]
                chart_data.add_series("Series 1", tuple(series_vals))
            
            x, y, cx, cy = Inches(2), Inches(2), Inches(9), Inches(4)
            ctype = XL_CHART_TYPE.COLUMN_CLUSTERED
            if chart_type.lower() == "line":
                ctype = XL_CHART_TYPE.LINE
            elif chart_type.lower() == "bar":
                ctype = XL_CHART_TYPE.BAR_CLUSTERED
            elif chart_type.lower() == "pie":
                ctype = XL_CHART_TYPE.PIE
                
            slide.shapes.add_chart(ctype, x, y, cx, cy, chart_data)
            prs.save(pres_id)
            return {"ok": True}

        return self._call_tool("create_chart_slide", {
            "presentation_id": pres_id,
            "chart_type": chart_type,
            "data": data,
            "title": title
        })

    def add_table_slide(self, pres_id: str, headers: list,
                        rows: list, title: str = None) -> dict:
        """Add formatted table."""
        if self.use_fallback:
            from pptx import Presentation
            from pptx.util import Inches
            prs = Presentation(pres_id)
            slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title Only or blank
            
            if title and slide.shapes.title:
                slide.shapes.title.text = title
                slide.shapes.title.text_frame.paragraphs[0].font.name = "Segoe UI"
                
            rows_cnt = len(rows) + 1
            cols_cnt = len(headers)
            left, top, width, height = Inches(1.5), Inches(2.0), Inches(10.0), Inches(0.8 * rows_cnt)
            
            table_shape = slide.shapes.add_table(rows_cnt, cols_cnt, left, top, width, height)
            table = table_shape.table
            
            # Write headers
            for col_idx, h in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.text = h
                cell.text_frame.paragraphs[0].font.name = "Segoe UI"
                cell.text_frame.paragraphs[0].font.bold = True
                
            # Write rows
            for row_idx, r in enumerate(rows):
                for col_idx, val in enumerate(r):
                    cell = table.cell(row_idx + 1, col_idx)
                    cell.text = str(val)
                    cell.text_frame.paragraphs[0].font.name = "Segoe UI"
                    
            prs.save(pres_id)
            return {"ok": True}

        return self._call_tool("add_table_slide", {
            "presentation_id": pres_id,
            "headers": headers,
            "rows": rows,
            "title": title
        })

    # ── Design & Enhancement ──

    def apply_theme_colors(self, pres_id: str,
                           theme_name: str = "Modern Blue") -> dict:
        """Apply professional color scheme."""
        # Local fallback theme coloring (can just be mock or simple background fill)
        if self.use_fallback:
            from pptx import Presentation
            from pptx.dml.color import RGBColor
            prs = Presentation(pres_id)
            
            # Simple theme palette mapping
            themes = {
                "Modern Blue": RGBColor(0, 102, 204),
                "Corporate Gray": RGBColor(96, 96, 96),
                "Elegant Green": RGBColor(0, 153, 76),
                "Warm Red": RGBColor(204, 0, 0)
            }
            color = themes.get(theme_name, RGBColor(0, 102, 204))
            
            for slide in prs.slides:
                # apply color to slide title text
                if slide.shapes.title and slide.shapes.title.has_text_frame:
                    for p in slide.shapes.title.text_frame.paragraphs:
                        p.font.color.rgb = color
            prs.save(pres_id)
            return {"ok": True, "theme": theme_name}

        return self._call_tool("apply_theme_colors", {
            "presentation_id": pres_id,
            "theme_name": theme_name
        })

    def enhance_slide(self, pres_id: str, slide_index: int) -> dict:
        """Enhance existing slide with professional styling."""
        if self.use_fallback:
            from pptx import Presentation
            prs = Presentation(pres_id)
            if slide_index >= len(prs.slides):
                return {"ok": False, "error": "Slide index out of bounds"}
            slide = prs.slides[slide_index]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        p.font.name = "Segoe UI"
            prs.save(pres_id)
            return {"ok": True}

        return self._call_tool("enhance_slide", {
            "presentation_id": pres_id,
            "slide_index": slide_index
        })

    # ── Reading ──

    def extract_slide_text(self, pres_id: str,
                           slide_index: int) -> dict:
        """Extract all text from a specific slide."""
        if self.use_fallback:
            from pptx import Presentation
            prs = Presentation(pres_id)
            if slide_index >= len(prs.slides):
                return {"text": "", "ok": False}
            slide = prs.slides[slide_index]
            text_runs = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text_runs.append(paragraph.text)
            return {"text": "\n".join(text_runs), "ok": True}

        return self._call_tool("extract_slide_text", {
            "presentation_id": pres_id,
            "slide_index": slide_index
        })

    def extract_presentation_text(self, pres_id: str) -> dict:
        """Extract text from all slides."""
        if self.use_fallback:
            from pptx import Presentation
            prs = Presentation(pres_id)
            slides_text = []
            for i, slide in enumerate(prs.slides):
                text_runs = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text_runs.append(paragraph.text)
                slides_text.append(f"Slide {i+1}:\n" + "\n".join(text_runs))
            return {"text": "\n\n".join(slides_text), "ok": True}

        return self._call_tool("extract_presentation_text", {
            "presentation_id": pres_id
        })

    def get_slide_info(self, pres_id: str,
                       slide_index: int) -> dict:
        """Get detailed slide information."""
        if self.use_fallback:
            from pptx import Presentation
            prs = Presentation(pres_id)
            if slide_index >= len(prs.slides):
                return {}
            slide = prs.slides[slide_index]
            
            title = ""
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title = slide.shapes.title.text_frame.text
                
            return {
                "title": title,
                "layout_name": slide.slide_layout.name,
                "shape_count": len(slide.shapes)
            }

        return self._call_tool("get_slide_info", {
            "presentation_id": pres_id,
            "slide_index": slide_index
        })

    # ── Save ──

    def save_presentation(self, pres_id: str, output_path: str) -> str:
        """Save presentation to file."""
        if self.use_fallback:
            if not os.path.exists(pres_id):
                raise FileNotFoundError(f"Source file {pres_id} not found")
            os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
            shutil.copy2(pres_id, output_path)
            return output_path

        res = self._call_tool("save_presentation", {
            "presentation_id": pres_id,
            "output_path": output_path
        })
        return res.get("output_path")

    def _call_tool(self, tool_name: str, args: dict) -> dict:
        """Call a PPTX MCP tool via uvx cli."""
        request = json.dumps({
            "method": "tools/call",
            "params": {
                "name": tool_name,
                "arguments": args
            }
        })
        cmd = self.mcp_command.split() + ["--call", tool_name, json.dumps(args)]
        result = subprocess.run(cmd, capture_output=True, text=True, check=True)
        return json.loads(result.stdout)
