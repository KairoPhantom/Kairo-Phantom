import logging
from pathlib import Path

log = logging.getLogger("kairo-sidecar.pptx_parser")

def parse_pptx(file_path: str) -> dict:
    """
    Read PowerPoint and return slide/shape inventory.
    """
    try:
        from pptx import Presentation
    except ImportError:
        return {"error": "python-pptx not installed — run: pip install python-pptx"}

    path = Path(file_path)
    if not path.exists():
        return {"error": f"File not found: {path}"}

    try:
        prs = Presentation(str(path))
        slides = []
        for i, slide in enumerate(prs.slides):
            shapes = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = "\n".join(p.text for p in shape.text_frame.paragraphs)
                    shapes.append({
                        "id": shape.shape_id,
                        "name": shape.name,
                        "text": text,
                        "left": shape.left,
                        "top": shape.top,
                    })
            # Slide title
            title = slide.shapes.title.text if slide.shapes.title else f"Slide {i+1}"
            slides.append({"index": i, "title": title, "shapes": shapes})

        return {"slides": slides, "slide_count": len(slides)}
    except Exception as e:
        log.error(f"Failed to parse PowerPoint file {file_path}: {e}")
        return {"error": f"Failed to parse PowerPoint: {e}"}
