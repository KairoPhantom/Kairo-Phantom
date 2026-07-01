# sidecar/oracles.py
import os
import sys
import time
import subprocess
import threading
import ipaddress
import re
import unicodedata
from pathlib import Path
from typing import Dict, List, Any, Optional

import docx
import openpyxl
import pptx
import pdfplumber
from PIL import Image
import imagehash

try:
    import fitz  # PyMuPDF — AGPL, lazy import to preserve licensing boundary
    HAS_FITZ = True
except ImportError:
    HAS_FITZ = False

try:
    import scapy.all as scapy
    HAS_SCAPY = True
except ImportError:
    HAS_SCAPY = False

try:
    import psutil
    HAS_PSUTIL = True
except ImportError:
    HAS_PSUTIL = False


def normalize_text(text: str) -> str:
    """Apply Unicode NFC normalization and fold/clean whitespace."""
    if not isinstance(text, str):
        text = str(text)
    # Unicode NFC normalization
    normalized = unicodedata.normalize("NFC", text)
    # Clean/fold whitespace: replace any whitespace sequence with a single space
    folded = re.sub(r"\s+", " ", normalized).strip()
    return folded


def normalize_xlsx_value(val: Any) -> str:
    """Normalize Excel cell values handles float formatting/rounding and string normalization."""
    if val is None:
        return ""
    try:
        if not isinstance(val, bool):
            f_val = float(val)
            rounded = round(f_val, 4)
            if rounded.is_integer():
                return str(int(rounded))
            return str(rounded)
    except (ValueError, TypeError):
        pass
    return normalize_text(val)


# ─── 1. Word / DOCX Oracle ────────────────────────────────────────────────────

def verify_docx(path: str, expected_text_substrings: Optional[List[str]] = None) -> bool:
    """Verifies DOCX structure and content."""
    from docx.text.paragraph import Paragraph
    from docx.table import Table

    doc = docx.Document(path)
    full_text = []

    # Traverse XML body elements of the document to extract text and tables in exact chronological layout order.
    for element in doc.element.body:
        tag = element.tag
        if '}' in tag:
            tag = tag.split('}', 1)[1]

        if tag == 'p':
            p = Paragraph(element, doc)
            full_text.append(p.text)
        elif tag == 'tbl':
            t = Table(element, doc)
            for row in t.rows:
                for cell in row.cells:
                    full_text.append(cell.text)

    combined_text = "\n".join(full_text)
    normalized_doc = normalize_text(combined_text)

    if expected_text_substrings:
        for substring in expected_text_substrings:
            normalized_substring = normalize_text(substring)
            if normalized_substring not in normalized_doc:
                raise AssertionError(f"Expected DOCX to contain: '{normalized_substring}' but not found.")
    return True


# ─── 2. Excel / XLSX Oracle ────────────────────────────────────────────────────

def verify_xlsx(path: str, cell_values: Optional[Dict[str, Any]] = None, cell_formulas: Optional[Dict[str, str]] = None) -> bool:
    """Verifies XLSX sheet structure, exact cell values, and formulas."""
    wb = openpyxl.load_workbook(path, data_only=False)
    try:
        sheet = wb.active

        if cell_formulas:
            for cell_ref, expected_formula in cell_formulas.items():
                val = sheet[cell_ref].value
                if not isinstance(val, str) or not val.startswith("="):
                    raise AssertionError(f"Expected formula in {cell_ref} but found value: {val}")
                if normalize_text(val).upper() != normalize_text(expected_formula).upper():
                    raise AssertionError(f"Formula mismatch in {cell_ref}: expected '{expected_formula}', found '{val}'")

        if cell_values:
            # Load workbook again with data_only=True to read evaluated values
            wb_val = openpyxl.load_workbook(path, data_only=True)
            try:
                sheet_val = wb_val.active
                for cell_ref, expected_val in cell_values.items():
                    val = sheet_val[cell_ref].value
                    normalized_actual = normalize_xlsx_value(val)
                    normalized_expected = normalize_xlsx_value(expected_val)
                    if normalized_actual != normalized_expected:
                        raise AssertionError(f"Value mismatch in {cell_ref}: expected '{expected_val}', found '{val}'")
            finally:
                wb_val.close()
    finally:
        wb.close()

    return True


# ─── 3. PowerPoint / PPTX Oracle ──────────────────────────────────────────────

def verify_pptx(path: str, expected_slide_count: Optional[int] = None, check_placeholders: bool = True, bullet_word_limit: int = 12, expected_text_substrings: Optional[List[str]] = None) -> bool:
    """Verifies PPTX structure, slide counts, placeholder slop, and word limits."""
    prs = pptx.Presentation(path)

    if expected_slide_count is not None:
        if len(prs.slides) != expected_slide_count:
            raise AssertionError(f"Slide count mismatch: expected {expected_slide_count}, found {len(prs.slides)}")

    all_slide_texts = []
    for i, slide in enumerate(prs.slides):
        # Sort all slide shapes visually (top-to-bottom, then left-to-right) with binning to prevent float flakiness
        def shape_sort_key(s):
            try:
                return (round(s.top / 5) * 5, s.left)
            except AttributeError:
                return (0, 0)
        sorted_shapes = sorted(slide.shapes, key=shape_sort_key)

        for shape in sorted_shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    # Perform Unicode NFC normalization and whitespace folding
                    text = normalize_text(paragraph.text)
                    if text:
                        all_slide_texts.append(text)

                    raw_text = paragraph.text.strip()
                    # Check for placeholders
                    if check_placeholders:
                        placeholders = ["click to add", "lorem ipsum", "[enter title]", "[insert text]"]
                        for pl in placeholders:
                            if pl in raw_text.lower():
                                raise AssertionError(f"Placeholder slop found on slide {i}: '{raw_text}'")
                    
                    # Check word limit on bullet points (including level 0 bullets in BODY/OBJECT placeholders)
                    is_bullet = False
                    if raw_text.startswith("-") or paragraph.level > 0:
                        is_bullet = True
                    else:
                        try:
                            ph_type = shape.placeholder_format.type
                            from pptx.enum.shapes import PP_PLACEHOLDER
                            if ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
                                is_bullet = True
                        except (AttributeError, ImportError, ValueError):
                            pass

                    if is_bullet:
                        words = len(raw_text.split())
                        if words > bullet_word_limit:
                            raise AssertionError(f"Slide {i} bullet exceeds word limit ({words} > {bullet_word_limit}): '{raw_text}'")

    if expected_text_substrings:
        combined_text = normalize_text(" ".join(all_slide_texts))
        for substring in expected_text_substrings:
            normalized_substring = normalize_text(substring)
            if normalized_substring not in combined_text:
                raise AssertionError(f"Expected PPTX to contain: '{normalized_substring}' but not found.")

    return True


# ─── 4. PDF Oracle ────────────────────────────────────────────────────────────

def verify_pdf(path: str, expected_substrings: Optional[List[str]] = None) -> bool:
    """Verifies PDF content by reading it back using both pdfplumber and PyMuPDF."""
    if not HAS_FITZ:
        raise ImportError("PyMuPDF (fitz) is not installed — cannot verify PDF content")
    # 1. PyMuPDF (fitz) text extraction
    extracted_fitz = []
    with fitz.open(path) as doc:
        for page in doc:
            blocks = page.get_text("blocks")
            sorted_blocks = sorted(blocks, key=lambda b: (round(b[1] / 5) * 5, b[0]))
            for b in sorted_blocks:
                extracted_fitz.append(b[4])
    combined_fitz = "\n".join(extracted_fitz)
    normalized_fitz = normalize_text(combined_fitz)

    # 2. pdfplumber text extraction
    extracted_plumber = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                extracted_plumber.append(t)
    combined_plumber = "\n".join(extracted_plumber)
    normalized_plumber = normalize_text(combined_plumber)

    if expected_substrings:
        for substring in expected_substrings:
            normalized_substring = normalize_text(substring)
            if normalized_substring not in normalized_fitz:
                raise AssertionError(f"Expected PDF (PyMuPDF) to contain: '{normalized_substring}' but not found.")
            if normalized_substring not in normalized_plumber:
                raise AssertionError(f"Expected PDF (pdfplumber) to contain: '{normalized_substring}' but not found.")
    return True


# ─── 5. LibreOffice Headless Excel Recompute Oracle ───────────────────────────

def excel_libreoffice_recompute(file_path: str, output_pdf_dir: str) -> str:
    """Runs LibreOffice headless to convert Excel sheet to PDF (which triggers formula recomputation).
    Returns the path to the converted PDF file. If LibreOffice is not installed, raises FileNotFoundError.
    """
    # 1. Search for soffice
    soffice_path = None
    common_paths = [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]
    for p in common_paths:
        if os.path.exists(p):
            soffice_path = p
            break
            
    if not soffice_path:
        # Check in PATH
        try:
            subprocess.run(["soffice", "--version"], stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            soffice_path = "soffice"
        except FileNotFoundError:
            pass
            
    if not soffice_path:
        # Halt and report blocked
        raise FileNotFoundError(
            "LibreOffice (soffice) not found on system. This is a critical system dependency. "
            "Mark this item as BLOCKED until LibreOffice is installed."
        )
        
    # Run LibreOffice headless conversion
    cmd = [
        soffice_path,
        "--headless",
        "--convert-to", "pdf",
        "--outdir", output_pdf_dir,
        file_path
    ]
    
    res = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, timeout=30)
    if res.returncode != 0:
        raise RuntimeError(f"LibreOffice conversion failed: {res.stderr}\n{res.stdout}")
        
    pdf_filename = Path(file_path).with_suffix(".pdf").name
    pdf_path = os.path.join(output_pdf_dir, pdf_filename)
    if not os.path.exists(pdf_path):
        raise RuntimeError(f"Converted PDF file not found at: {pdf_path}")
        
    return pdf_path


def verify_xlsx_recomputed_values(xlsx_path: str, temp_dir: str, expected_values: List[str]) -> bool:
    """Converts xlsx to PDF to trigger recomputation, reads PDF back, asserts expected values exist."""
    pdf_path = None
    try:
        # excel_libreoffice_recompute will raise FileNotFoundError if LibreOffice is missing
        pdf_path = excel_libreoffice_recompute(xlsx_path, temp_dir)
        verify_pdf(pdf_path, expected_values)
    finally:
        if pdf_path and os.path.exists(pdf_path):
            try:
                os.remove(pdf_path)
            except Exception:
                pass
    return True


# ─── 6. Network Sniffer Oracle ────────────────────────────────────────────────

class NetworkSnifferOracle:
    """A zero-flake network sniffer that monitors external egress connections.
    Uses scapy if available/privileged, falling back to psutil scanning of active TCP sockets.
    """
    def __init__(self):
        self.external_destinations = set()
        self.stop_sniffing = threading.Event()
        self._thread = None

    def _is_private(self, ip_str: str) -> bool:
        try:
            ip = ipaddress.ip_address(ip_str)
            return ip.is_private or ip.is_loopback or ip.is_multicast or ip.is_link_local
        except ValueError:
            return True # If not parseable, ignore or treat as safe for routing
            
    def _run_scapy(self):
        def prn(pkt):
            if self.stop_sniffing.is_set():
                return
            if pkt.haslayer(scapy.IP):
                dst = pkt[scapy.IP].dst
                if not self._is_private(dst):
                    self.external_destinations.add(dst)
            elif pkt.haslayer(scapy.IPv6):
                dst = pkt[scapy.IPv6].dst
                if not self._is_private(dst):
                    self.external_destinations.add(dst)
        try:
            scapy.sniff(prn=prn, stop_filter=lambda x: self.stop_sniffing.is_set(), timeout=60)
        except Exception:
            pass # Fall back silently to psutil

    def _get_process_connections(self) -> List[Any]:
        conns = []
        try:
            current_proc = psutil.Process()
            try:
                conns.extend(current_proc.connections(kind="inet"))
            except (psutil.AccessDenied, psutil.NoSuchProcess):
                pass
            
            try:
                children = current_proc.children(recursive=True)
            except Exception:
                children = []
                
            for child in children:
                try:
                    conns.extend(child.connections(kind="inet"))
                except (psutil.AccessDenied, psutil.NoSuchProcess):
                    pass
        except Exception:
            pass
        return conns

    def _run_psutil(self):
        use_fallback = False
        while not self.stop_sniffing.is_set():
            try:
                connections = []
                if not use_fallback:
                    try:
                        connections = psutil.net_connections(kind="inet")
                    except (psutil.AccessDenied, PermissionError):
                        use_fallback = True
                        connections = self._get_process_connections()
                else:
                    connections = self._get_process_connections()
                
                for conn in connections:
                    if conn.raddr:
                        r_ip = conn.raddr.ip
                        if not self._is_private(r_ip):
                            self.external_destinations.add(r_ip)
            except Exception:
                pass
            time.sleep(0.1)

    def start(self):
        self.external_destinations.clear()
        self.stop_sniffing.clear()
        
        # Prefer psutil because it doesn't require admin/Npcap on Windows
        if HAS_PSUTIL:
            self._thread = threading.Thread(target=self._run_psutil, daemon=True)
            self._thread.start()
        elif HAS_SCAPY:
            self._thread = threading.Thread(target=self._run_scapy, daemon=True)
            self._thread.start()

    def stop(self) -> List[str]:
        self.stop_sniffing.set()
        if self._thread:
            self._thread.join(timeout=1.0)
        return list(self.external_destinations)


# ─── 7. Screenshot Diff / Image Oracle ────────────────────────────────────────

def verify_screenshot_diff(path_a: str, path_b: str, max_hash_diff: int = 2) -> bool:
    """Verifies that two screenshots match visually using perceptual hash."""
    from PIL import ImageStat
    
    with Image.open(path_a) as raw_a, Image.open(path_b) as raw_b:
        with raw_a.convert("RGB") as img_a, raw_b.convert("RGB") as img_b:
            if img_a.size != img_b.size:
                raise AssertionError(f"Image dimensions mismatch: {img_a.size} vs {img_b.size}")
                
            hash_a = imagehash.average_hash(img_a)
            hash_b = imagehash.average_hash(img_b)
            
            diff = hash_a - hash_b
            
            # Check overall color/luminance difference to handle flat solid color images
            stat_a = ImageStat.Stat(img_a)
            stat_b = ImageStat.Stat(img_b)
            mean_diff = sum(abs(a - b) for a, b in zip(stat_a.mean, stat_b.mean)) / len(stat_a.mean)
            if mean_diff > 10:
                diff += int(mean_diff)
                
            if diff > max_hash_diff:
                raise AssertionError(f"Visual diff failed: hash difference {diff} > {max_hash_diff}")
            
    return True


verify_screenshot_diff_hash = verify_screenshot_diff

