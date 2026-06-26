import os
import shutil
import logging
from pathlib import Path
from typing import List, Dict, Any

log = logging.getLogger("kairo-sidecar.pptx_writer")


def _enforce_title_words(text: str, max_words: int = 7) -> str:
    """Trim title/bullet to max_words words."""
    words = text.strip().split()
    return " ".join(words[:max_words])


def _write_title(shape, text: str):
    """Write title text with Segoe UI 40pt bold."""
    from pptx.util import Pt

    tf = shape.text_frame
    tf.text = _enforce_title_words(text, 7)
    for p in tf.paragraphs:
        p.font.name = "Segoe UI"
        p.font.size = Pt(40)
        p.font.bold = True


def _write_bullets(shape, bullets: list):
    """Write up to 5 bullets (≤ 7 words each) into a content shape."""
    from pptx.util import Pt

    tf = shape.text_frame
    tf.clear()
    for i, b in enumerate(bullets[:5]):
        text = _enforce_title_words(b, 7)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.name = "Segoe UI"
        p.font.size = Pt(18)


def _find_body_shape(slide, exclude_title=True):
    """Find the content/body placeholder on a slide."""
    # Prefer placeholder idx=1 (body/content)
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == 1:
            return shape
    # Fall back to any text frame that isn't the title
    title_shape = slide.shapes.title
    for shape in slide.shapes:
        if shape.has_text_frame and (not exclude_title or shape is not title_shape):
            return shape
    return None


def write_pptx(file_path: str, operations: List[Dict[str, Any]]) -> dict:
    """
    Apply SlideOperation list to .pptx file atomically.

    Supported operation modes:
      add_new=True   → Append a brand-new slide with title + bullets in one shot.
                        Fields: title (str), bullets (list[str]), layout_index (int, default 1)
      add_new=False  → Update an existing slide at slide_index.
                        Supports sub-types: update_title, update_shape_text (bullets field),
                        or the legacy bullets-only form.

    Title constraint  : ≤ 7 words (enforced, excess trimmed)
    Bullet constraint : ≤ 7 words per bullet; max 5 bullets per slide (enforced)
    Font              : Segoe UI throughout (title 40pt bold, body 18pt)
    """
    try:
        from pptx import Presentation
        from pptx.util import Pt, Inches
    except ImportError:
        return {"error": "python-pptx not installed — run: pip install python-pptx"}

    path = Path(file_path)
    if not path.exists():
        return {"error": f"File not found: {path}"}

    # ── Atomic backup ────────────────────────────────────────────────────────
    backup_path = path.with_suffix(path.suffix + ".kairo_backup")
    try:
        shutil.copy2(path, backup_path)
    except Exception as e:
        return {"error": f"Failed to create backup: {e}"}

    try:
        prs = Presentation(str(path))
    except Exception as e:
        return {"error": f"Failed to load presentation: {e}"}

    applied: list = []
    errors: list = []

    for op in operations:
        try:
            # ── Mode 1: Add new slide ─────────────────────────────────────────
            if op.get("add_new"):
                layout_idx = int(op.get("layout_index", 1))
                if layout_idx >= len(prs.slide_layouts):
                    layout_idx = 1  # fallback to Title+Content
                layout = prs.slide_layouts[layout_idx]
                slide = prs.slides.add_slide(layout)

                # Title
                title_text = op.get("title", "")
                title_shape = slide.shapes.title
                if title_text:
                    if title_shape:
                        _write_title(title_shape, title_text)
                    else:
                        # No title placeholder on this layout → add text box
                        tb = slide.shapes.add_textbox(
                            Inches(0.5), Inches(0.3), Inches(12), Inches(1)
                        )
                        _write_title(tb, title_text)

                # Bullets
                bullets = op.get("bullets", [])
                if bullets:
                    body = _find_body_shape(slide)
                    if not body:
                        # No body placeholder → add text box
                        body = slide.shapes.add_textbox(
                            Inches(0.5), Inches(1.5), Inches(12), Inches(5)
                        )
                    _write_bullets(body, bullets)

                applied.append(
                    {
                        "type": "add_slide",
                        "title": title_text,
                        "bullets": len(bullets),
                    }
                )
                continue

            # ── Mode 2: Update existing slide ─────────────────────────────────
            op_type = op.get("type")
            slide_idx = int(op.get("slide_index", 0))

            # Infer op_type from fields when not explicitly given
            if not op_type:
                if "bullets" in op and op["bullets"]:
                    op_type = "update_shape_text"
                    op["paragraphs"] = [{"text": b, "bullet": True} for b in op["bullets"]]
                elif "text" in op:
                    op_type = "update_title"

            if slide_idx >= len(prs.slides):
                errors.append(f"Slide index {slide_idx} out of range ({len(prs.slides)} slides)")
                continue

            slide = prs.slides[slide_idx]

            if op_type == "update_title":
                title_shape = slide.shapes.title
                if not title_shape:
                    for shape in slide.shapes:
                        if shape.is_placeholder and shape.placeholder_format.idx == 0:
                            title_shape = shape
                            break
                if not title_shape:
                    title_shape = slide.shapes.add_textbox(
                        Inches(1), Inches(0.5), Inches(11.3), Inches(1)
                    )
                _write_title(title_shape, op.get("text", ""))
                applied.append({"type": "update_title", "slide_index": slide_idx})

            elif op_type == "update_shape_text":
                shape_id = op.get("shape_id")
                paragraphs_data = op.get("paragraphs", [])
                if not paragraphs_data and op.get("bullets"):
                    paragraphs_data = [{"text": b, "bullet": True} for b in op["bullets"]]

                # Speaker notes fallback (e.g. shape_id == 9999)
                if shape_id is not None and str(shape_id) == "9999":
                    bullets = op.get("bullets", [])
                    if not bullets and paragraphs_data:
                        bullets = [p.get("text", "") for p in paragraphs_data]
                    slide.notes_slide.notes_text_frame.text = "\n".join(bullets)
                    applied.append(
                        {"type": "update_shape_text", "slide_index": slide_idx, "notes": bullets}
                    )
                    continue

                target_shape = None
                for shape in slide.shapes:
                    if shape_id is not None and str(shape.shape_id) == str(shape_id):
                        target_shape = shape
                        break
                if not target_shape:
                    target_shape = _find_body_shape(slide)
                if not target_shape:
                    target_shape = slide.shapes.add_textbox(
                        Inches(0.5), Inches(1.5), Inches(12), Inches(5)
                    )

                if op.get("left") is not None:
                    target_shape.left = Inches(float(op["left"]))
                if op.get("top") is not None:
                    target_shape.top = Inches(float(op["top"]))
                if op.get("width") is not None:
                    target_shape.width = Inches(float(op["width"]))
                if op.get("height") is not None:
                    target_shape.height = Inches(float(op["height"]))

                from pptx.util import Pt

                tf = target_shape.text_frame
                tf.clear()
                bullets_written = 0
                for i, para_data in enumerate(paragraphs_data):
                    is_bullet = para_data.get("bullet", True)
                    if is_bullet and bullets_written >= 5:
                        continue
                    p_text = (
                        _enforce_title_words(para_data.get("text", ""), 7)
                        if is_bullet
                        else para_data.get("text", "")
                    )
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = p_text
                    p.level = para_data.get("level", 0)
                    p.font.name = "Segoe UI"
                    p.font.size = Pt(18)
                    if is_bullet:
                        bullets_written += 1
                applied.append({"type": "update_shape_text", "slide_index": slide_idx})

            elif op_type == "add_slide":
                after_idx = int(op.get("after_index", 0))
                layout_name = op.get("layout_name", "Title and Content")
                layout = prs.slide_layouts[1]
                for l in prs.slide_layouts:
                    if l.name.lower() == layout_name.lower():
                        layout = l
                        break
                slide = prs.slides.add_slide(layout)
                if after_idx < len(prs.slides) - 2:
                    slides_id_list = prs.slides._sldIdLst
                    slide_id_element = slides_id_list[-1]
                    slides_id_list.remove(slide_id_element)
                    slides_id_list.insert(after_idx + 1, slide_id_element)

                # Write title if provided
                title_text = op.get("title")
                if title_text:
                    title_shape = slide.shapes.title
                    if title_shape:
                        _write_title(title_shape, title_text)
                    else:
                        tb = slide.shapes.add_textbox(
                            Inches(0.5), Inches(0.3), Inches(12), Inches(1)
                        )
                        _write_title(tb, title_text)

                # Write bullets if provided
                bullets = op.get("bullets")
                if bullets:
                    body = _find_body_shape(slide)
                    if not body:
                        body = slide.shapes.add_textbox(
                            Inches(0.5), Inches(1.5), Inches(12), Inches(5)
                        )
                    _write_bullets(body, bullets)

                applied.append({"type": "add_slide", "after_index": after_idx})

            elif op_type == "update_notes":
                slide.notes_slide.notes_text_frame.text = op.get("text", "")
                applied.append({"type": "update_notes", "slide_index": slide_idx})

            else:
                errors.append(f"Unknown operation type '{op_type}' — op: {op}")

        except Exception as e:
            errors.append(f"Error on op {op}: {e}")

    # ── Atomic save ──────────────────────────────────────────────────────────
    tmp_path = path.with_suffix(path.suffix + ".pptx.tmp")
    try:
        prs.save(str(tmp_path))
        os.replace(str(tmp_path), str(path))
    except PermissionError:
        try:
            if backup_path.exists():
                shutil.copy2(str(backup_path), str(path))
        except Exception:
            pass
        try:
            backup_path.unlink(missing_ok=True)
        except Exception:
            pass
        try:
            tmp_path.unlink()
        except Exception:
            pass
        return {
            "error": "PowerPoint has this file open. Save and close the presentation first, then press Alt+Ctrl+M again.",
            "path": str(path),
        }
    except Exception as e:
        try:
            if backup_path.exists():
                shutil.copy2(str(backup_path), str(path))
        except Exception:
            pass
        try:
            backup_path.unlink(missing_ok=True)
        except Exception:
            pass
        try:
            tmp_path.unlink()
        except Exception:
            pass
        return {"error": f"Failed to save or replace presentation: {e}"}

    # ── Clean up backup (only restore on total failure above) ────────────────
    if not errors:
        try:
            backup_path.unlink(missing_ok=True)
        except Exception:
            pass

    return {
        "applied_count": len(applied),
        "errors": errors,
        "slides": [a for a in applied],
    }
