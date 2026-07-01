import os
import re
import logging
from typing import List, Any

# Import schemas
from sidecar.schemas.domain_schemas import (
    CodeResponse,
    PDFResponse,
    BrowserResponse,
    TerminalResponse,
    EmailResponse,
    NotesResponse,
    DesignResponse,
    MediaResponse,
    DataResponse,
)
from sidecar.schemas.pptx_schema import SlideResponse
from sidecar.observability.opik_tracer import track

log = logging.getLogger("kairo-sidecar.other_masters")


# Helper result class
class ValidationResult:
    def __init__(self, valid: bool, error: str = "", op: dict = None):
        self.valid = valid
        self.error = error
        self.op = op


# ==========================================
# 1. PowerPoint Master
# ==========================================
class PowerPointMaster:
    def extract_context(self, file_path: str, cursor_info: Any) -> dict:
        try:
            from sidecar.parsers.pptx_context import PptxContextCapture
            import json

            capture = PptxContextCapture()
            slide_idx = (
                int(cursor_info)
                if isinstance(cursor_info, (int, str)) and str(cursor_info).isdigit()
                else 0
            )

            ctx = capture.capture(file_path, slide_idx)
            ctx["slide_index"] = slide_idx
            ctx["total_slides"] = ctx.get("slide_count", 0)
            ctx["layout_name"] = ctx.get("current_slide", {}).get(
                "layout_name", "Title and Content"
            )
            ctx["major_font"] = "Segoe UI Light"
            ctx["minor_font"] = "Segoe UI"
            ctx["deck_purpose"] = "sales_deck"
            ctx["shapes_json"] = "[]"

            try:
                from pptx import Presentation

                prs = Presentation(file_path)
                slide_count = len(prs.slides)
                ctx["total_slides"] = slide_count
                if slide_count > 0:
                    slide_idx = max(0, min(slide_idx, slide_count - 1))
                    slide = prs.slides[slide_idx]
                    ctx["layout_name"] = slide.slide_layout.name

                    shapes_list = []
                    for shape in slide.shapes:
                        shape_info = {
                            "shape_id": str(shape.shape_id)
                            if hasattr(shape, "shape_id")
                            else str(shape.id),
                            "name": shape.name,
                            "type": str(shape.shape_type)
                            if hasattr(shape, "shape_type")
                            else "Unknown",
                        }
                        if shape.has_text_frame:
                            shape_info["text"] = shape.text_frame.text
                        shapes_list.append(shape_info)
                    ctx["shapes_json"] = json.dumps(shapes_list, indent=2)
            except Exception as e:
                log.warning(f"PowerPointMaster python-pptx extraction failed: {e}")

            return ctx
        except Exception as e:
            log.warning(f"PowerPointMaster context extraction failed: {e}")
            return {
                "full_text": "",
                "current_slide": {},
                "slide_text": "",
                "slide_count": 0,
                "theme": "Default",
                "user_preferences": {},
                "slide_index": 0,
                "total_slides": 0,
                "layout_name": "Title and Content",
                "major_font": "Segoe UI Light",
                "minor_font": "Segoe UI",
                "deck_purpose": "sales_deck",
                "shapes_json": "[]",
            }

    def build_prompt(
        self, user_instruction: str, context: dict, mem_context: str, classification: Any = None
    ) -> str:
        # Build prompt using PPTX prompt fragment from Capture
        from sidecar.parsers.pptx_context import PptxContextCapture

        capture = PptxContextCapture()
        fragment = capture.to_system_prompt_fragment(context)

        app_ctx = f"""=== APP CONTEXT ===
Application Name: Microsoft PowerPoint
Application Type: Presentation
File Path: {context.get('file_path', 'Unknown')}"""

        # Document Context
        doc_ctx = f"""=== DOCUMENT CONTEXT ===
SLIDE CONTEXT (injected):
Current slide: {context.get('slide_index', 0)} of {context.get('total_slides', 0)}
Slide layout: {context.get('layout_name', 'Title and Content')}
Theme fonts — Major: {context.get('major_font', 'Segoe UI Light')} / Minor: {context.get('minor_font', 'Segoe UI')}
Shapes on this slide:
{context.get('shapes_json', '[]')}
Deck purpose: {context.get('deck_purpose', 'sales_deck')}

## Presentation Overview Details:
{fragment}"""

        mem_ctx = f"""=== MEMORY CONTEXT ===
User Writing Preferences:
{mem_context or "No writing preferences learned yet. Use professional defaults."}"""

        intent_part = """=== INTENT CLASSIFICATION ===
Intent: PowerPoint Presentation / Slide Design Action"""

        system_rules = """SYSTEM:
You are the PowerPoint Presentation Master for Kairo Phantom. You understand slide design principles,
presentation narrative, and the python-pptx shape model deeply.

CONSTRAINTS:
- Output ONLY valid JSON matching SlideResponse schema. No prose. No explanation.
- NEVER modify shapes not explicitly targeted
- NEVER hardcode font names or sizes — use theme-relative values or null
- NEVER exceed 7 words per bullet point — this rule has no exceptions
- NEVER exceed 6 bullet points per slide

SLIDE DESIGN RULES:
1. One idea per slide — if content exceeds one idea, split into two slides
2. Titles: 5-8 words, active voice, specific (not "Overview" — say "Revenue Grew 40% in Q3")
3. Bullets: max 7 words, present tense, parallel structure, start with verbs when possible
4. Data: always include the unit ($, %, x) and timeframe (Q3, 2025, YoY)
5. Never use complete sentences in bullets — fragments only
6. Three types of slides: (1) claim + evidence bullets, (2) visual + one-line caption, (3) full-bleed stat

DECK PURPOSE BEHAVIOR:
investor_pitch: Bold claims. Specific numbers. Short slides. Each slide = one investment thesis.
sales_deck: Customer pain first. Solution second. Social proof. Clear CTA.
training: Learning objectives upfront. Build complexity gradually. Summary at end.
status_update: RAG status (Red/Amber/Green). Decisions needed. Next steps with owner/date.
conference: Story arc. Section dividers. Quotable moments. Speaker-friendly notes.

OUTPUT SCHEMA:
{
  "operations": [
    {
      "type": "update_shape_text",
      "slide_index": 2,
      "shape_id": "sp-1",
      "paragraphs": [
        {"text": "Revenue Up 40% YoY", "bold": false, "italic": false, "bullet": false, "level": 0, "font_size": null},
        {"text": "Fastest growth in company history", "bold": false, "italic": false, "bullet": true, "level": 0, "font_size": null}
      ]
    },
    {
      "type": "update_notes",
      "slide_index": 2,
      "text": "Speaker notes text here"
    },
    {
      "type": "add_slide",
      "after_index": 3,
      "layout_name": "Title and Content",
      "title": "Slide Title",
      "bullets": ["First point max 7 words", "Second point here", "Third point here"]
    }
  ],
  "confidence": 0.0-1.0,
  "needs_clarification": false,
  "clarification_question": null
}

7-WORD BULLET ENFORCEMENT — EXAMPLES:
WRONG (10 words): "Our platform increases user productivity by reducing manual work"
RIGHT (5 words): "Reduces manual work by 60%"
WRONG (9 words): "We have successfully expanded into three new European markets"
RIGHT (4 words): "Three new European markets"
"""

        json_reminder = "REMINDER: Your entire response must be a single JSON object. First character must be {. Last character must be }."

        return f"""{system_rules}

{app_ctx}

{doc_ctx}

{mem_ctx}

{intent_part}

{json_reminder}
USER INSTRUCTION: {user_instruction}
OUTPUT (JSON only):
"""

    def validate_operations(self, response: SlideResponse, context: dict) -> List[dict]:
        validated_ops = []
        for op in response.operations:
            op_dict = op.model_dump()
            op_type = op_dict.get("type")

            # Slide index bounds check
            slide_count = context.get("slide_count", 0)
            if op_type in ("update_shape_text", "update_title", "update_notes"):
                idx = op_dict.get("slide_index", 0)
                if slide_count > 0:
                    op_dict["slide_index"] = max(0, min(idx, slide_count - 1))

            # Bullet word count check
            if op_type == "update_shape_text":
                for p in op_dict.get("paragraphs", []):
                    text = p.get("text", "")
                    words = text.split()
                    if p.get("bullet") and len(words) > 7:
                        p["text"] = " ".join(words[:7])  # Clamping

            # Title word count check
            if op_type == "update_title":
                text = op_dict.get("text", "")
                words = text.split()
                if len(words) > 7:
                    op_dict["text"] = " ".join(words[:7])

            validated_ops.append(op_dict)
        return validated_ops

    @track("pptx", "apply_operations")
    def apply_operations(self, file_path: str, operations: list) -> dict:
        """Apply validated slide operations to a .pptx file via the PPTX writer.

        Args:
            file_path: Absolute path to the .pptx file.
            operations: Validated ops list produced by validate_operations().

        Returns:
            dict with keys applied_count, errors, and slides (same shape as write_pptx).
        """
        try:
            from sidecar.writers.pptx_writer import write_pptx

            result = write_pptx(file_path, operations)
            log.info(
                f"PowerPointMaster.apply_operations: applied={result.get('applied_count', 0)} "
                f"errors={len(result.get('errors', []))} path={file_path}"
            )
            return result
        except Exception as e:
            log.error(f"PowerPointMaster.apply_operations failed: {e}")
            return {"applied_count": 0, "errors": [str(e)], "slides": []}

    def get_schema_class(self):
        return SlideResponse


# ==========================================
# 2. Code Master
# ==========================================
class CodeMaster:
    @track("code", "extract_context")
    def extract_context(self, file_path: str, cursor_info: Any) -> dict:
        context = {
            "language": "python",
            "file_path": file_path,
            "cursor_line": 1,
            "indent_style": "spaces",
            "indent_size": 4,
            "line_endings": "LF",
            "enclosing_function_signature": "None",
            "enclosing_class_name": "None",
            "existing_imports": [],
            "surrounding_code": "",
        }
        if not file_path or not os.path.exists(file_path):
            return context

        try:
            # Simple line/context extractor
            ext = os.path.splitext(file_path)[1].lower()
            lang_map = {
                ".py": "python",
                ".rs": "rust",
                ".ts": "typescript",
                ".js": "javascript",
                ".go": "go",
                ".java": "java",
                ".cpp": "cpp",
                ".c": "c",
                ".sh": "bash",
            }
            context["language"] = lang_map.get(ext, "python")

            cursor_line = (
                int(cursor_info)
                if isinstance(cursor_info, (int, str)) and str(cursor_info).isdigit()
                else 1
            )
            context["cursor_line"] = cursor_line

            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                lines = f.readlines()

            context["line_endings"] = "CRLF" if any("\r\n" in l for l in lines[:10]) else "LF"

            # Capture surrounding code
            start = max(0, cursor_line - 15)
            end = min(len(lines), cursor_line + 15)
            context["surrounding_code"] = "".join(lines[start:end])

            # Guess indentation
            for l in lines[:100]:
                if l.startswith("\t"):
                    context["indent_style"] = "tabs"
                    context["indent_size"] = 1
                    break
                elif l.startswith(" "):
                    stripped = l.lstrip()
                    if len(l) - len(stripped) > 0:
                        context["indent_style"] = "spaces"
                        context["indent_size"] = len(l) - len(stripped)
                        break

            # Find imports
            for l in lines[:50]:
                if l.startswith(("import ", "from ", "use ", "const ", "require(")):
                    context["existing_imports"].append(l.strip())

            # Parse enclosing function signature and class name
            enclosing_func = "None"
            enclosing_class = "None"
            class_indent = -1
            func_indent = -1
            for idx in range(min(cursor_line - 1, len(lines) - 1), -1, -1):
                line = lines[idx]
                stripped = line.strip()
                indent = len(line) - len(line.lstrip())

                if stripped.startswith(("class ", "class\t")):
                    if class_indent == -1 or indent < class_indent:
                        enclosing_class = stripped.split("{")[0].split(":")[0].strip()
                        class_indent = indent
                elif (
                    stripped.startswith(("def ", "function ", "func ", "fn "))
                    or "public " in stripped
                    or "private " in stripped
                ):
                    if func_indent == -1 or indent < func_indent:
                        enclosing_func = stripped.split("{")[0].split(":")[0].strip()
                        func_indent = indent

            context["enclosing_function_signature"] = enclosing_func
            context["enclosing_class_name"] = enclosing_class

        except Exception as e:
            log.warning(f"CodeMaster context extraction failed: {e}")
        return context

    def build_prompt(
        self, user_instruction: str, context: dict, mem_context: str, classification: Any = None
    ) -> str:
        app_ctx = f"""=== APP CONTEXT ===
Application Name: VS Code / Code Editor
Application Type: Integrated Development Environment
File Path: {context.get('file_path', 'Unknown')}"""

        # Format imports list nicely
        existing_imports = ", ".join(context.get("existing_imports", []))

        doc_ctx = f"""=== DOCUMENT CONTEXT ===
CODE CONTEXT (injected by tree-sitter):
Language: {context.get('language')}
File: {context.get('file_path')}
Cursor line: {context.get('cursor_line')}
Indentation: {context.get('indent_style')} ({context.get('indent_size')} chars)
Line endings: {context.get('line_endings')}
Enclosing function: {context.get('enclosing_function_signature')}
Class context: {context.get('enclosing_class_name')}
Imports already in file: {existing_imports}
Surrounding code (30 lines):
{context.get('surrounding_code')}"""

        mem_ctx = f"""=== MEMORY CONTEXT ===
User Coding Preferences:
{mem_context or "No coding preferences learned yet. Use professional defaults."}"""

        intent_part = """=== INTENT CLASSIFICATION ===
Intent: Code Generation / Modification Action"""

        system_rules = """SYSTEM:
You are the Code Master for Kairo Phantom. You write correct, idiomatic code in any programming language.
You have deep knowledge of syntax, conventions, and best practices for 50+ languages.

CONSTRAINTS:
- Output ONLY valid JSON matching CodeResponse schema. No prose.
- NEVER reformat code you were not asked to reformat
- NEVER change indentation style (tabs/spaces) of surrounding code
- NEVER add imports without checking if they already exist in the file
- Preserve exact whitespace and line endings of the insertion point
- NEVER auto-execute commands in terminal mode — show them in GRP only

LANGUAGE-SPECIFIC CONVENTIONS:
Python: Google or NumPy docstring style. Type hints on all functions. f-strings not .format(). List comprehensions over map/filter. pytest for tests.
Rust: /// doc comments. Result and Option over exceptions. Descriptive variable names. Clippy-clean. cargo test format.
TypeScript: JSDoc or TSDoc comments. Explicit return types. interface over type for objects. Jest or Vitest for tests. async/await not .then().
Go: godoc comments starting with function name. Named return values for complex functions. table-driven tests. errors wrapped with fmt.Errorf.
Java: Javadoc. Checked exceptions declared. @Override annotations. JUnit 5 for tests.
SQL: UPPERCASE keywords. Explicit column names (no SELECT *). Aliases for readability. CTEs for complex queries.
Bash/Shell: set -euo pipefail. Quote all variables. [[ not [. Functions before use.

OUTPUT SCHEMA:
{
  "operations": [
    {
      "type": "insert_at_line",
      "line": 45,
      "code": "string",
      "language": "python"
    },
    {
      "type": "replace_lines",
      "start_line": 10,
      "end_line": 20,
      "code": "string"
    },
    {
      "type": "add_import",
      "import_statement": "string",
      "insert_at_line": 1
    },
    {
      "type": "show_only",
      "content": "string"
    }
  ],
  "confidence": 0.0-1.0,
  "needs_clarification": false,
  "clarification_question": null
}

EXAMPLES:

EXAMPLE 1 — Python docstring
Context: language=python, enclosing_function="def calculate_roi(investment: float, returns: float) -> float:", cursor_line=45, indent=4 spaces
Prompt: "// add docstring"
Output: {"operations":[{"type":"insert_at_line","line":46,"code":"    \\\"\\\"\\\"Calculate return on investment as a percentage.\\n    \\n    Args:\\n        investment: The initial investment amount in dollars.\\n        returns: The total returns received in dollars.\\n    \\n    Returns:\\n        The ROI as a percentage (e.g., 0.25 = 25% return).\\n    \\n    Raises:\\n        ValueError: If investment is zero or negative.\\n    \\\"\\\"\\\"","language":"python"}],"confidence":0.97,"needs_clarification":false,"clarification_question":null}

EXAMPLE 2 — Terminal command (show_only — never inject)
Context: language=bash, terminal context
Prompt: "// Delete all log files older than 30 days"
Output: {"operations":[{"type":"show_only","content":"find /var/log -name '*.log' -mtime +30 -type f\\n\\n# Review the list above, then run:\\n# find /var/log -name '*.log' -mtime +30 -type f -delete\\n\\n⚠️  This will permanently delete files. Review carefully before running."}],"confidence":0.89,"needs_clarification":false,"clarification_question":null}"""

        json_reminder = "REMINDER: Your entire response must be a single JSON object. First character must be {. Last character must be }."

        return f"""{system_rules}

{app_ctx}

{doc_ctx}

{mem_ctx}

{intent_part}

{json_reminder}
USER INSTRUCTION: {user_instruction}
OUTPUT (JSON only):
"""

    def validate_operations(self, response: CodeResponse, context: dict) -> List[dict]:
        validated_ops = []
        for op in response.operations:
            op_dict = op.model_dump()
            validated_ops.append(op_dict)
        return validated_ops

    def get_schema_class(self):
        return CodeResponse


class WeKnoraPipeline:
    """
    RAG pipeline framework for document Q&A without hallucination,
    exposing citation-backed answers. Uses SQLite FTS5.
    """

    def __init__(self, db_path: str = ":memory:"):
        self.db_path = db_path
        self._init_db()

    def _init_db(self):
        import sqlite3

        conn = sqlite3.connect(self.db_path)
        # Create virtual FTS5 table for chunks if it doesn't exist
        try:
            conn.execute("""
                CREATE VIRTUAL TABLE IF NOT EXISTS pdf_chunks USING fts5(
                    pdf_path,
                    chunk_index INT,
                    content
                )
            """)
            conn.commit()
        except sqlite3.OperationalError:
            # Fallback if FTS5 is not enabled in SQLite build
            conn.execute("""
                CREATE TABLE IF NOT EXISTS pdf_chunks (
                    pdf_path TEXT,
                    chunk_index INT,
                    content TEXT
                )
            """)
            conn.commit()
        finally:
            conn.close()

    def ingest(self, pdf_path: str):
        if not pdf_path or not os.path.exists(pdf_path):
            return

        # Check if already ingested
        import sqlite3

        conn = sqlite3.connect(self.db_path)
        cur = conn.cursor()
        cur.execute("SELECT COUNT(*) FROM pdf_chunks WHERE pdf_path = ?", (pdf_path,))
        exists = cur.fetchone()[0] > 0
        if exists:
            conn.close()
            log.info(f"WeKnora: PDF {pdf_path} already ingested, skipping.")
            return

        # Extract text using pdfminer.six or fallback to pdf_parser (which tries PyMuPDF etc.)
        text = ""
        try:
            # Try pdfminer.six
            from pdfminer.high_level import extract_text

            text = extract_text(pdf_path)
        except Exception as e:
            log.debug(f"WeKnora: pdfminer.six extraction failed: {e}. Trying parser fallback...")
            try:
                from sidecar.parsers.pdf_parser import parse_pdf

                parsed = parse_pdf(pdf_path)
                paragraphs = parsed.get("paragraphs", [])
                text = "\n".join(p.get("text", "") for p in paragraphs)
            except Exception as e2:
                log.warning(f"WeKnora: fallback parsing also failed: {e2}")

        if not text:
            conn.close()
            return

        # Chunk into 300-word chunks
        words = text.split()
        chunk_size = 300
        chunks = []
        for i in range(0, len(words), chunk_size):
            chunk = " ".join(words[i : i + chunk_size])
            if chunk.strip():
                chunks.append(chunk)

        # Write to SQLite pdf_chunks
        try:
            for idx, content in enumerate(chunks):
                cur.execute(
                    "INSERT INTO pdf_chunks (pdf_path, chunk_index, content) VALUES (?, ?, ?)",
                    (pdf_path, idx, content),
                )
            conn.commit()
            log.info(f"WeKnora: Ingested {len(chunks)} chunks of 300 words from {pdf_path}")
        except Exception as e:
            log.warning(f"WeKnora: insertion failed: {e}")
        finally:
            conn.close()

    def query(self, pdf_path: str, question: str, k: int = 5) -> dict:
        import sqlite3

        if not os.path.exists(pdf_path):
            return {"answer": "No content available in the PDF.", "citations": []}

        self.ingest(pdf_path)

        conn = sqlite3.connect(self.db_path)
        cur = conn.cursor()

        # Clean question for FTS5 (remove punctuation, only alphanumeric/spaces)
        clean_q = " ".join(re.findall(r"\w+", question))

        results = []
        try:
            # Query virtual FTS5 table
            cur.execute(
                """
                SELECT chunk_index, content 
                FROM pdf_chunks 
                WHERE pdf_path = ? AND content MATCH ? 
                LIMIT ?
            """,
                (pdf_path, clean_q, k),
            )
            results = cur.fetchall()
        except sqlite3.OperationalError:
            # Fallback if MATCH fails or FTS5 not available
            cur.execute(
                "SELECT chunk_index, content FROM pdf_chunks WHERE pdf_path = ?", (pdf_path,)
            )
            all_chunks = cur.fetchall()

            # Simple keyword match
            q_words = set(clean_q.lower().split())
            scored = []
            for idx, content in all_chunks:
                c_words = set(content.lower().split())
                score = len(q_words.intersection(c_words))
                if score > 0:
                    scored.append((score, idx, content))
            scored.sort(key=lambda x: x[0], reverse=True)
            results = [(s[1], s[2]) for s in scored[:k]]
        finally:
            conn.close()

        if results:
            answer = "\n\n".join(r[1] for r in results)
            citations = [f"Page/Section Chunk {r[0] + 1}" for r in results]
            return {"answer": answer, "citations": citations}
        else:
            # Full text fallback
            conn = sqlite3.connect(self.db_path)
            cur = conn.cursor()
            cur.execute("SELECT content FROM pdf_chunks WHERE pdf_path = ? LIMIT 1", (pdf_path,))
            first = cur.fetchone()
            conn.close()
            val = first[0] if first else "No matching content found."
            return {"answer": val, "citations": ["First paragraph fallback"]}


# ==========================================
# 3. PDF Master
# ==========================================
class PDFMaster:
    def __init__(self):
        self._weknora = WeKnoraPipeline()

    @track("pdf", "extract_context")
    def extract_context(self, file_path: str, cursor_info: Any) -> dict:
        context = {
            "extraction_tier": "PyMuPDF",
            "document_type": "PDF Document",
            "page_count": 1,
            "extracted_content": "Empty PDF",
            "language": "en",
            "confidence": 1.0,
            "file_path": file_path,
        }
        if not file_path or not os.path.exists(file_path):
            return context

        try:
            from sidecar.parsers.pdf_extraction_engine import PdfExtractionEngine

            engine = PdfExtractionEngine(offline_mode=True)
            result = engine.extract(file_path)
            context["extraction_tier"] = result.tier_used.name if result.tier_used else "PyMuPDF"
            context["page_count"] = result.metadata.get("pages", 1)
            context["extracted_content"] = result.text[:4000]  # Truncate for prompt
            context["language"] = result.language or "en"
            context["confidence"] = result.confidence

            # Ingest to WeKnora
            self._weknora.ingest(file_path)
        except Exception as e:
            log.warning(f"PDFMaster context extraction failed: {e}")
        return context

    def build_prompt(
        self, user_instruction: str, context: dict, mem_context: str, classification: Any = None
    ) -> str:
        # Query WeKnora if user is asking a question about the PDF
        weknora_ctx = ""
        pdf_path = context.get("file_path", "")
        if pdf_path and (
            "?" in user_instruction
            or "what" in user_instruction.lower()
            or "explain" in user_instruction.lower()
            or "say about" in user_instruction.lower()
        ):
            try:
                import json

                qa_res = self._weknora.query(pdf_path, user_instruction)
                weknora_ctx = f"\n=== WEKNORA CITATION-BACKED Q&A CONTEXT ===\nAnswer: {qa_res['answer']}\nCitations: {json.dumps(qa_res['citations'])}\n"
            except Exception as e:
                log.warning(f"WeKnora query failed: {e}")

        app_ctx = f"""=== APP CONTEXT ===
Application Name: Adobe Acrobat / Browser PDF Viewer
Application Type: PDF Viewer
File Path: {context.get('file_path', 'Unknown')}"""

        doc_ctx = f"""=== DOCUMENT CONTEXT ===
EXTRACTED CONTENT (injected by MinerU/olmOCR/PyMuPDF):
Extraction tier used: {context.get('extraction_tier')}
Document type: {context.get('document_type')}
Page count: {context.get('page_count')}
Extracted content:
{context.get('extracted_content')}
{weknora_ctx}"""

        mem_ctx = f"""=== MEMORY CONTEXT ===
User Document Preferences:
{mem_context or "No preferences learned yet."}"""

        intent_part = """=== INTENT CLASSIFICATION ===
Intent: PDF Extraction / Content Transformation Action"""

        system_rules = """SYSTEM:
You are the PDF Intelligence Master for Kairo Phantom. You extract, understand, and transform PDF content.
PDFs are READ-ONLY — you never edit them in-place. You create new output documents.

CONSTRAINTS:
- Output ONLY valid JSON matching PDFResponse schema
- NEVER claim to edit the original PDF — always create new output
- NEVER hallucinate page numbers or section names not in the provided context
- For scanned PDFs with OCR errors: acknowledge uncertainty in confidence score

OUTPUT SCHEMA:
{
  "output_type": "docx|text|clipboard|excel_table",
  "operations": [
    {"type": "create_docx", "sections": [
      {"heading": "string", "level": 1, "body": "string"},
      {"heading": "string", "level": 2, "body": "string"},
      {"table": {"headers": [], "rows": [[]]}}
    ]},
    {"type": "clipboard_text", "content": "string"},
    {"type": "excel_table", "headers": [], "rows": [[]]}
  ],
  "output_filename": "original_name_kairo_output.docx",
  "confidence": 0.0-1.0,
  "extraction_quality": "high|medium|low",
  "needs_clarification": false,
  "clarification_question": null
}"""

        json_reminder = "REMINDER: Your entire response must be a single JSON object. First character must be {. Last character must be }."

        return f"""{system_rules}

{app_ctx}

{doc_ctx}

{mem_ctx}

{intent_part}

{json_reminder}
USER INSTRUCTION: {user_instruction}
OUTPUT (JSON only):
"""

    def validate_operations(self, response: PDFResponse, context: dict) -> List[dict]:
        validated_ops = []
        for op in response.operations:
            op_dict = op.model_dump()
            validated_ops.append(op_dict)
        return validated_ops

    def apply_output(self, file_path: str, output_type: str, pdf_content: dict) -> str:
        """Convert parsed PDF content to a new file via the PDF output writer.

        Args:
            file_path: Source .pdf path (used for naming the output file).
            output_type: 'docx' or 'pdf'.
            pdf_content: Structured dict produced by parse_pdf() — must contain
                         'paragraphs', 'tables', and 'metadata' keys.

        Returns:
            Absolute path of the written output file.
        """
        try:
            from sidecar.writers.pdf_output_writer import write_pdf_output

            output_path = write_pdf_output(pdf_content, output_type, file_path)
            log.info(f"PDFMaster.apply_output: wrote {output_type.upper()} → {output_path}")
            return output_path
        except Exception as e:
            log.error(f"PDFMaster.apply_output failed: {e}")
            raise

    def get_schema_class(self):
        return PDFResponse


# ==========================================
# 4. Browser Master
# ==========================================
class BrowserMaster:
    @track("browser", "extract_context")
    def extract_context(self, file_path: str, cursor_info: Any) -> dict:
        url = file_path or "http://localhost"
        is_collab = "notion.so" in url.lower() or "docs.google.com" in url.lower()
        return {
            "page_url": url,
            "page_title": "Active Browser Page",
            "active_element_type": "body",
            "platform": "Chrome",
            "page_content_truncated": "No content loaded",
            "is_collaborative_editor": is_collab,
        }

    def build_prompt(
        self, user_instruction: str, context: dict, mem_context: str, classification: Any = None
    ) -> str:
        intent_classification = "Browser Interaction"
        if classification is not None:
            if isinstance(classification, str):
                intent_classification = classification
            elif isinstance(classification, dict):
                intent_classification = (
                    classification.get("task_type")
                    or classification.get("intent")
                    or classification.get("classification")
                    or "Browser Interaction"
                )
            else:
                intent_classification = (
                    getattr(classification, "task_type", None)
                    or getattr(classification, "intent", None)
                    or getattr(classification, "classification", None)
                    or "Browser Interaction"
                )

        app_ctx = f"""=== APP CONTEXT ===
Application Name: {context.get('platform', 'Chrome')}
Application Type: Web Browser"""

        doc_ctx = f"""=== DOCUMENT CONTEXT ===
Page URL: {context.get('page_url', '')}
Page Title: {context.get('page_title', '')}
Active Element Type: {context.get('active_element_type', '')}
Page Content:
{context.get('page_content_truncated', '')}"""

        if mem_context:
            mem_ctx = f"""=== MEMORY CONTEXT ===
USER WRITING PREFERENCES (from memory):
{mem_context}"""
        else:
            mem_ctx = """=== MEMORY CONTEXT ===
No preferences learned yet."""

        intent_part = f"""=== INTENT CLASSIFICATION ===
Intent: {intent_classification}"""

        system_rules = """SYSTEM:
You are the Browser and Web Content Master for Kairo Phantom.
You understand web page content, help draft web-form inputs, and write platform-appropriate content.

CONSTRAINTS:
- NEVER auto-submit forms or trigger any click actions
- NEVER inject into password fields, credit card fields, or 2FA inputs
- For collaborative editors (Notion, Google Docs): use CRDT path, flag in response
- Output length must match the target field type

PLATFORM CONVENTIONS:
linkedin_post: Hook first line (question or bold claim). 3-5 short paragraphs. Bullets optional. CTA last. Under 1300 chars. No hashtags in body — 3 at end max.
twitter_thread: First tweet = hook (under 280 chars). Numbered tweets 2-10. Last tweet = summary + CTA.
gmail_compose: Subject under 50 chars. Professional greeting. Clear ask in first 2 sentences. Specific next step in closing. Under 200 words unless complex topic.
notion_block: Use Notion block types — paragraph, bullet_list, heading_2, toggle, callout. Yjs CRDT injection.
google_docs: Collaborative editor — flag for CRDT path. Match existing document style.
web_form_input: Match field label requirements. Avoid placeholder text. Appropriate length.

OUTPUT SCHEMA:
{
  "injection_method": "uia_field|clipboard|crdt_yjs|show_only",
  "content": "string — the text to inject or show",
  "platform_formatted": true,
  "is_collaborative_editor": false,
  "safety_check": {"is_password_field": false, "is_payment_field": false, "is_auto_submit": false},
  "confidence": 0.0-1.0,
  "needs_clarification": false,
  "clarification_question": null
}"""

        json_reminder = "REMINDER: Your entire response must be a single JSON object. First character must be {. Last character must be }."

        return f"""{system_rules}

{app_ctx}

{doc_ctx}

{mem_ctx}

{intent_part}

{json_reminder}
USER INSTRUCTION: {user_instruction}
OUTPUT (JSON only):
"""

    def validate_operations(self, response: BrowserResponse, context: dict) -> List[dict]:
        op = response.model_dump()

        # 1. Safety check verification
        if (
            response.safety_check.is_password_field
            or response.safety_check.is_payment_field
            or response.safety_check.is_auto_submit
        ):
            log.warning(
                "Safety block: password, payment, or auto-submit field injection requested."
            )
            return []

        # Runtime input guardrails on context: active_element_type, page_url, page_title
        url_lower = context.get("page_url", "").lower()
        element_lower = context.get("active_element_type", "").lower()
        title_lower = context.get("page_title", "").lower()

        # Reject password, 2fa, otp, security, credit card inputs
        unsafe_keywords = [
            "password",
            "passwd",
            "creditcard",
            "cardnumber",
            "cvv",
            "cvc",
            "2fa",
            "otp",
            "passcode",
            "pin_number",
        ]
        is_unsafe = False
        for kw in unsafe_keywords:
            if (
                kw in element_lower
                or kw in url_lower
                or kw in title_lower
                or "2fa" in url_lower
                or "otp" in url_lower
            ):
                is_unsafe = True
                break

        if is_unsafe:
            log.warning("Safety block: runtime context indicates password, payment, or 2FA field.")
            return []

        # 2. Collaborative editors override
        is_collab_url = "notion.so" in url_lower or "docs.google.com" in url_lower
        if is_collab_url:
            op["injection_method"] = "crdt_yjs"
            op["is_collaborative_editor"] = True

        return [op]

    def get_schema_class(self):
        return BrowserResponse


# ==========================================
# 5. Terminal Master
# ==========================================
class TerminalMaster:
    @track("terminal", "extract_context")
    def extract_context(self, file_path: str, cursor_info: Any) -> dict:
        import platform

        os_type = platform.system()
        cwd = file_path if file_path and os.path.isdir(file_path) else os.getcwd()

        git_info = "No git repository found"
        try:
            if os.path.exists(os.path.join(cwd, ".git")):
                git_info = "Git repository detected (active branch info: check status)"
        except Exception:
            pass

        return {
            "shell_type": "powershell" if os_type == "Windows" else "bash",
            "os_type": os_type,
            "current_directory": cwd,
            "terminal_content": "",
            "git_info": git_info,
        }

    def build_prompt(
        self, user_instruction: str, context: dict, mem_context: str, classification: Any = None
    ) -> str:
        intent_classification = "Terminal command generation"
        if classification is not None:
            if isinstance(classification, str):
                intent_classification = classification
            elif isinstance(classification, dict):
                intent_classification = (
                    classification.get("task_type")
                    or classification.get("intent")
                    or classification.get("classification")
                    or "Terminal command generation"
                )
            else:
                intent_classification = (
                    getattr(classification, "task_type", None)
                    or getattr(classification, "intent", None)
                    or getattr(classification, "classification", None)
                    or "Terminal command generation"
                )

        app_ctx = f"""=== APP CONTEXT ===
Application Name: {context.get('shell_type', 'powershell')}
Application Type: Terminal / Shell Console"""

        doc_ctx = f"""=== DOCUMENT CONTEXT ===
OS: {context.get('os_type', 'Windows')}
Current directory: {context.get('current_directory', '')}
Visible terminal content: {context.get('terminal_content', '')}
Git repo: {context.get('git_info', '')}"""

        mem_ctx = f"""=== MEMORY CONTEXT ===
User Preferences: {mem_context or "None"}"""

        intent_part = f"""=== INTENT CLASSIFICATION ===
Intent: {intent_classification}"""

        system_rules = """SYSTEM:
You are the Terminal and Shell Master for Kairo Phantom.
You generate precise shell commands for any shell, any OS.

CRITICAL SAFETY RULE:
ALL terminal output uses "show_only" injection type. Commands are SHOWN in the Ghost Review Panel.
They are NEVER automatically typed into the terminal. The user reviews and runs them manually.
This rule has no exceptions — not even for "safe" commands.

SHELL SYNTAX RULES:
PowerShell: Get-ChildItem not ls (unless alias confirmed). -ErrorAction SilentlyContinue for optional. $() for subshells. Write-Host for output.
Bash/Zsh: set -euo pipefail in scripts. "$()" not backticks. [[ ]] not [ ]. printf not echo for formatted.
CMD: Legacy syntax. dir not ls. /? for help. %variable% not $variable.
WSL Bash: Linux commands. /mnt/c/ for Windows drives. sudo available.

DANGER LEVELS:
SAFE: ls, cat, echo, git status, grep, find (no -delete), curl (no -X DELETE)
CAUTION: rm with single file, git reset, chmod, chown — add "# Review before running" comment
DANGEROUS: rm -rf, DROP TABLE, format, del /s, shutdown — MUST add prominent ⚠️ warning

OUTPUT SCHEMA:
{
  "injection_method": "show_only",
  "command": "string — the exact command",
  "explanation": "string — what it does in plain English",
  "danger_level": "safe|caution|dangerous",
  "warning": "string or null — shown in GRP for caution/dangerous",
  "alternative": "string or null — safer alternative if dangerous",
  "confidence": 0.0-1.0,
  "needs_clarification": false,
  "clarification_question": null
}"""

        json_reminder = "REMINDER: Your entire response must be a single JSON object. First character must be {. Last character must be }."

        return f"""{system_rules}

{app_ctx}

{doc_ctx}

{mem_ctx}

{intent_part}

{json_reminder}
USER INSTRUCTION: {user_instruction}
OUTPUT (JSON only):
"""

    def validate_operations(self, response: TerminalResponse, context: dict) -> List[dict]:
        op_dict = response.model_dump()
        op_dict["injection_method"] = "show_only"

        command = op_dict.get("command", "")
        cmd_lower = command.lower()

        # 1. Programmatic Danger Level Guardrails
        is_dangerous = False
        dangerous_patterns = ["rm -rf", "drop table", "format", "del /s", "shutdown", "rmdir /s"]
        for pattern in dangerous_patterns:
            if pattern in cmd_lower:
                is_dangerous = True
                break

        is_caution = False
        caution_patterns = ["rm ", "git reset", "chmod", "chown", "del "]
        if not is_dangerous:
            for pattern in caution_patterns:
                if pattern in cmd_lower:
                    is_caution = True
                    break

        if is_dangerous:
            op_dict["danger_level"] = "dangerous"
            if not op_dict.get("warning") or "⚠️" not in op_dict.get("warning", ""):
                op_dict["warning"] = (
                    "⚠️ WARNING: This command is highly destructive or system-critical. "
                    + (op_dict.get("warning") or "")
                )
        elif is_caution:
            op_dict["danger_level"] = "caution"
            if not op_dict.get("warning") or "review" not in op_dict.get("warning", "").lower():
                op_dict["warning"] = "Review before running. " + (op_dict.get("warning") or "")

            # Ensure the command has a review comment if not present
            if (
                "# review before running" not in cmd_lower
                and "rem review before running" not in cmd_lower
            ):
                if context.get("shell_type") == "cmd":
                    op_dict["command"] = f"{command} & rem Review before running"
                else:
                    op_dict["command"] = f"{command} # Review before running"

        return [op_dict]

    def get_schema_class(self):
        return TerminalResponse


# ==========================================
# 6. Email Master
# ==========================================
class EmailMaster:
    @track("email", "extract_context")
    def extract_context(self, file_path: str, cursor_info: Any) -> dict:
        return {
            "email_client": "Outlook",
            "compose_mode": "new",
            "thread_context": "",
            "preferred_signoff": "Best regards,",
            "user_name": "User",
        }

    def build_prompt(
        self, user_instruction: str, context: dict, mem_context: str, classification: Any = None
    ) -> str:
        signoff = context.get("preferred_signoff", "Best regards,")
        name = context.get("user_name", "User")

        intent_classification = "Email Drafting"
        if classification is not None:
            if isinstance(classification, str):
                intent_classification = classification
            elif isinstance(classification, dict):
                intent_classification = (
                    classification.get("task_type")
                    or classification.get("intent")
                    or classification.get("classification")
                    or "Email Drafting"
                )
            else:
                intent_classification = (
                    getattr(classification, "task_type", None)
                    or getattr(classification, "intent", None)
                    or getattr(classification, "classification", None)
                    or "Email Drafting"
                )

        app_ctx = f"""=== APP CONTEXT ===
Application Name: {context.get('email_client', 'Outlook')}
Application Type: Email Client"""

        doc_ctx = f"""=== DOCUMENT CONTEXT ===
Mode: {context.get('compose_mode', 'new')} (new|reply|forward)
Thread context (prior emails):
{context.get('thread_context', '')}"""

        mem_ctx = f"""=== MEMORY CONTEXT ===
User's preferred sign-off (MemMachine): {signoff}
User's name (MemMachine): {name}
User Preferences: {mem_context or "None"}"""

        intent_part = f"""=== INTENT CLASSIFICATION ===
Intent: {intent_classification}"""

        system_rules = """SYSTEM:
You are the Email Communication Master for Kairo Phantom.
You draft professional, effective, appropriate email communications.

CRITICAL RULE: NEVER include any instruction to send, submit, or forward an email.
Output is always shown in GRP for user review before they decide to send.

EMAIL WRITING RULES:
1. Subject line: specific, under 50 chars, no ALL CAPS, no exclamation marks
2. Opening: match formality of thread. Formal: "Dear [Name]," Casual: "Hi [Name],"
3. First sentence: state the purpose directly — no "I hope this email finds you well"
4. Body: one topic per paragraph, max 3 paragraphs for most emails
5. Ask: make the request or next step explicit — "Could you please..." or "I'll send by..."
6. Closing: match opening formality. "Best regards," / "Thanks," / "Sincerely,"
7. Sign-off: use MemMachine preference or infer from context

EMOTION DETECTION:
If the user's draft contains words like "frustrated", "unacceptable", "terrible", "angry":
- Set emotional_flag=true
- Add suggested_revision that removes emotional language
- Note in GRP: "I've suggested a more professional version"

PII HANDLING:
If thread_context contains SSNs, credit cards, or medical IDs:
- Do not include them in your output
- Add pii_redacted=true to response

OUTPUT SCHEMA:
{
  "injection_method": "uia_field|clipboard",
  "subject": "string or null (null if reply)",
  "body": "string",
  "emotional_flag": false,
  "suggested_revision": "string or null",
  "pii_redacted": false,
  "tone": "formal|professional|casual",
  "word_count": integer,
  "confidence": 0.0-1.0,
  "needs_clarification": false,
  "clarification_question": null
}"""

        json_reminder = "REMINDER: Your entire response must be a single JSON object. First character must be {. Last character must be }."

        return f"""{system_rules}

{app_ctx}

{doc_ctx}

{mem_ctx}

{intent_part}

{json_reminder}
USER INSTRUCTION: {user_instruction}
OUTPUT (JSON only):
"""

    def validate_operations(self, response: EmailResponse, context: dict) -> List[dict]:
        op_dict = response.model_dump()

        # 1. Subject line rules
        if op_dict.get("subject"):
            sub = op_dict["subject"]
            sub = sub.replace("!", "")
            if sub.isupper():
                sub = sub.title()
            if len(sub) > 50:
                sub = sub[:50]
            op_dict["subject"] = sub

        # 2. Block send instructions
        for key in ["body", "suggested_revision"]:
            val = op_dict.get(key)
            if val:
                val = re.sub(r"(?i)\b(?:please )?(?:send|submit|forward) this email\b\.?", "", val)
                op_dict[key] = val.strip()

        # 3. Emotion Detection
        emotional_words = ["frustrated", "unacceptable", "terrible", "angry"]
        user_prompt_lower = context.get("user_prompt", "").lower()
        body_lower = op_dict.get("body", "").lower()

        has_emotion = any(w in user_prompt_lower or w in body_lower for w in emotional_words)
        if has_emotion:
            op_dict["emotional_flag"] = True
            if not op_dict.get("suggested_revision"):
                revised = op_dict["body"]
                revised = re.sub(r"(?i)\bfrustrated\b", "concerned", revised)
                revised = re.sub(r"(?i)\bunacceptable\b", "unsatisfactory", revised)
                revised = re.sub(r"(?i)\bterrible\b", "difficult", revised)
                revised = re.sub(r"(?i)\bangry\b", "displeased", revised)
                op_dict["suggested_revision"] = revised

        # 4. PII Handling
        thread_ctx = context.get("thread_context", "")
        ssn_pattern = r"\b\d{3}-\d{2}-\d{4}\b"
        card_pattern = r"\b(?:\d[ -]*?){13,16}\b"
        medical_pattern = r"\bMED-\d{6,10}\b"

        thread_has_pii = (
            re.search(ssn_pattern, thread_ctx)
            or re.search(card_pattern, thread_ctx)
            or re.search(medical_pattern, thread_ctx)
        )

        if thread_has_pii:
            op_dict["pii_redacted"] = True
            for key in ["body", "subject", "suggested_revision"]:
                val = op_dict.get(key)
                if val:
                    val = re.sub(ssn_pattern, "[SSN REDACTED]", val)
                    val = re.sub(card_pattern, "[CARD REDACTED]", val)
                    val = re.sub(medical_pattern, "[MEDICAL ID REDACTED]", val)
                    op_dict[key] = val

        return [op_dict]

    def get_schema_class(self):
        return EmailResponse


# ==========================================
# 7. Notes Master
# ==========================================
class NotesMaster:
    @track("notes", "extract_context")
    def extract_context(self, file_path: str, cursor_info: Any) -> dict:
        app = "Obsidian"
        if file_path:
            ext = os.path.splitext(file_path)[1].lower()
            if ext == ".org":
                app = "Logseq"
            elif ext == ".txt":
                app = "Plain .md"

        return {
            "notes_app": app,
            "file_path": file_path or "scratchpad.md",
            "current_heading": "None",
            "cursor_line": int(cursor_info)
            if isinstance(cursor_info, (int, str)) and str(cursor_info).isdigit()
            else 1,
            "surrounding_content": "",
            "existing_tags": [],
            "backlinks": [],
        }

    def build_prompt(
        self, user_instruction: str, context: dict, mem_context: str, classification: Any = None
    ) -> str:
        app = context.get("notes_app", "Obsidian")
        file = context.get("file_path", "scratchpad.md")
        heading = context.get("current_heading", "None")
        line = context.get("cursor_line", 1)
        surrounding = context.get("surrounding_content", "")
        tags = context.get("existing_tags", [])
        links = context.get("backlinks", [])

        intent_classification = "Notes Management"
        if classification is not None:
            if isinstance(classification, str):
                intent_classification = classification
            elif isinstance(classification, dict):
                intent_classification = (
                    classification.get("task_type")
                    or classification.get("intent")
                    or classification.get("classification")
                    or "Notes Management"
                )
            else:
                intent_classification = (
                    getattr(classification, "task_type", None)
                    or getattr(classification, "intent", None)
                    or getattr(classification, "classification", None)
                    or "Notes Management"
                )

        app_ctx = f"""=== APP CONTEXT ===
App: {app}
Application Type: Note-Taking Application"""

        doc_ctx = f"""=== DOCUMENT CONTEXT ===
File: {file}
Current section heading: {heading}
Cursor line: {line}
Surrounding content:
{surrounding}
Existing tags: {tags}
Linked notes (backlinks): {links}"""

        mem_ctx = f"""=== MEMORY CONTEXT ===
User Preferences: {mem_context or "None"}"""

        intent_part = f"""=== INTENT CLASSIFICATION ===
Intent: {intent_classification}"""

        system_rules = """SYSTEM:
You are the Notes and Knowledge Management Master for Kairo Phantom.
You write structured notes, expand ideas, and manage knowledge in any note-taking application.

APP-SPECIFIC FORMATTING:
Obsidian: Wikilinks as [[Note Title]]. Tags as #tag. Frontmatter as YAML between ---. Callouts as > [!note], > [!warning]. Embeds as ![[note]].
Notion: Block-based. Use \n for new blocks. /heading2 syntax in output description. Inline code as `code`.
Plain .md: Standard CommonMark. # ## ### for headings. - for bullets. ** for bold. * for italic.
Logseq: Indented bullets. [[page refs]]. #tag. TODO/DONE markers.

OUTPUT SCHEMA:
{
  "injection_method": "file_write|uia_field|clipboard",
  "insert_at_line": integer,
  "content": "string — markdown formatted for the specific app",
  "new_tags": [],
  "new_links": [],
  "frontmatter_update": null,
  "confidence": 0.0-1.0,
  "needs_clarification": false,
  "clarification_question": null
}"""

        json_reminder = "REMINDER: Your entire response must be a single JSON object. First character must be {. Last character must be }."

        return f"""{system_rules}

{app_ctx}

{doc_ctx}

{mem_ctx}

{intent_part}

{json_reminder}
USER INSTRUCTION: {user_instruction}
OUTPUT (JSON only):
"""

    def validate_operations(self, response: NotesResponse, context: dict) -> List[dict]:
        op_dict = response.model_dump()
        return [op_dict]

    def get_schema_class(self):
        return NotesResponse


# ==========================================
# 8. Design Master
# ==========================================
class DesignMaster:
    def extract_context(self, file_path: str, cursor_info: Any) -> dict:
        tool = "Figma"
        if file_path and "canva" in file_path.lower():
            tool = "Canva"

        return {
            "design_tool": tool,
            "active_frame_name": "Frame 1",
            "canvas_dimensions": [1920, 1080],
            "color_tokens": {"primary": "#0055FF", "secondary": "#FF8800", "neutral": "#333333"},
            "type_tokens": {"heading": "Inter Bold 32px", "body": "Inter Regular 16px"},
            "layers_json": "[]",
            "auto_layout_active": False,
        }

    def build_prompt(
        self, user_instruction: str, context: dict, mem_context: str, classification: Any = None
    ) -> str:
        tool = context.get("design_tool", "Figma")
        frame = context.get("active_frame_name", "Frame 1")
        dimensions = context.get("canvas_dimensions", [1920, 1080])
        colors = context.get("color_tokens", {})
        typography = context.get("type_tokens", {})
        layers = context.get("layers_json", "[]")

        intent_classification = "Design Action"
        if classification is not None:
            if isinstance(classification, str):
                intent_classification = classification
            elif isinstance(classification, dict):
                intent_classification = (
                    classification.get("task_type")
                    or classification.get("intent")
                    or classification.get("classification")
                    or "Design Action"
                )
            else:
                intent_classification = (
                    getattr(classification, "task_type", None)
                    or getattr(classification, "intent", None)
                    or getattr(classification, "classification", None)
                    or "Design Action"
                )

        app_ctx = f"""=== APP CONTEXT ===
Design tool: {tool}
Application Type: Design Tool Editor"""

        doc_ctx = f"""=== DOCUMENT CONTEXT ===
Active frame/artboard: {frame}
Canvas dimensions: {dimensions}
Design tokens (colors): {colors}
Design tokens (typography): {typography}
Layers in active frame:
{layers}"""

        mem_ctx = f"""=== MEMORY CONTEXT ===
User Preferences: {mem_context or "None"}"""

        intent_part = f"""=== INTENT CLASSIFICATION ===
Intent: {intent_classification}"""

        system_rules = """SYSTEM:
You are the Design Master for Kairo Phantom. You create and modify visual designs with deep knowledge
of design principles, color theory, typography, and layout.

CONSTRAINTS:
- NEVER hardcode RGB colors when theme colors are available — use theme color references
- NEVER set absolute pixel positions when auto-layout is active — use auto-layout padding/gap
- NEVER modify layer z-order unless explicitly requested
- For Canva: output clipboard-based text only — no direct API access available

DESIGN PRINCIPLES:
1. Visual hierarchy: most important element is largest or highest contrast
2. Alignment: elements align to an 8px grid unless design uses different base unit
3. Spacing: consistent spacing multiples (4px, 8px, 16px, 24px, 32px, 48px)
4. Color: max 3 colors per component (primary, secondary, neutral). Text meets WCAG AA contrast.
5. Typography: max 2 typefaces. Clear size hierarchy (heading > subheading > body > caption).
6. Whitespace: generous padding makes designs feel premium. Minimum 16px padding inside components.

OUTPUT SCHEMA:
{
  "injection_method": "figma_mcp|penpot_mcp|clipboard",
  "operations": [
    {
      "type": "create_text",
      "parent_frame_id": "string",
      "text": "string",
      "x": 0, "y": 0, "width": 200,
      "font_size": null,
      "font_weight": "regular|medium|semibold|bold",
      "color_token": "string or null",
      "alignment": "left|center|right"
    },
    {
      "type": "set_text",
      "node_id": "string",
      "text": "string"
    },
    {
      "type": "set_fills",
      "node_id": "string",
      "color_hex": "#hex",
      "opacity": 1.0
    },
    {
      "type": "show_only",
      "design_suggestion": "string — design recommendation shown in GRP"
    }
  ],
  "design_rationale": "string — why these design decisions",
  "accessibility_notes": "string or null",
  "confidence": 0.0-1.0,
  "needs_clarification": false,
  "clarification_question": null
}

CANVA FALLBACK:
Canva has no programmatic API. For Canva context, always use:
injection_method: "clipboard" + content with formatted text.
Include clear instructions for the user: "Copy this and paste into your Canva text element." """

        json_reminder = "REMINDER: Your entire response must be a single JSON object. First character must be {. Last character must be }."

        return f"""{system_rules}

{app_ctx}

{doc_ctx}

{mem_ctx}

{intent_part}

{json_reminder}
USER INSTRUCTION: {user_instruction}
OUTPUT (JSON only):
"""

    def validate_operations(self, response: DesignResponse, context: dict) -> List[dict]:
        op_dict = response.model_dump()
        tool = context.get("design_tool", "Figma").lower()
        context.get("color_tokens", {})
        auto_layout = context.get("auto_layout_active", False)

        if "canva" in tool:
            op_dict["injection_method"] = "clipboard"

        validated_ops = []
        for op in response.operations:
            op_data = op.model_dump()

            if op_data["type"] == "create_text":
                if auto_layout:
                    op_data["x"] = 0
                    op_data["y"] = 0

            validated_ops.append(op_data)

        op_dict["operations"] = validated_ops
        return [op_dict]

    def get_schema_class(self):
        return DesignResponse


# ==========================================
# 9. Media Master
# ==========================================
class MediaMaster:
    @track("media", "extract_context")
    def extract_context(self, file_path: str, cursor_info: Any) -> dict:
        """
        Media Master context: uses farscry visual analysis + UIA app fingerprinting.
        Detects: timeline scrubber (video editor), canvas elements (Canva), layer panel (Photoshop).
        App fingerprinted from window title + UI structure.
        """
        ctx = {
            "active_app": "Canva",
            "app_type": "graphic_editor",  # graphic_editor | video_editor | audio_editor
            "timeline_scrubber_seconds": 0,
            "layers": [],
            "canvas_elements": [],
            "selected_element_text": "",
            "selected_element_type": "text",  # text | image | video | audio | shape
            "file_path": file_path or "",
            "injection_path": "clipboard",  # clipboard | uia | script
        }

        # App detection from file path / window context
        if file_path:
            fp_lower = str(file_path).lower()
            if any(ext in fp_lower for ext in [".drp", ".fcpx", ".prproj"]):
                ctx["active_app"] = "DaVinci Resolve" if ".drp" in fp_lower else "Premiere Pro"
                ctx["app_type"] = "video_editor"
                ctx["injection_path"] = "script"
            elif ".aep" in fp_lower:
                ctx["active_app"] = "After Effects"
                ctx["app_type"] = "video_editor"
                ctx["injection_path"] = "script"
            elif any(ext in fp_lower for ext in [".psd", ".ai", ".indd"]):
                ctx["active_app"] = "Adobe Photoshop"
                ctx["app_type"] = "graphic_editor"
                ctx["injection_path"] = "uia"
            elif ".canva" in fp_lower or "canva.com" in fp_lower:
                ctx["active_app"] = "Canva"
                ctx["app_type"] = "graphic_editor"
                ctx["injection_path"] = "clipboard"

        # Try farscry visual analysis if cursor_info has coordinates
        if isinstance(cursor_info, dict) and "x" in cursor_info and "y" in cursor_info:
            try:
                from sidecar.kairo_eye.farscry_service import FarscryService

                service = FarscryService()
                analysis = service.analyze_cursor_region(cursor_info["x"], cursor_info["y"])
                ctx["visual_element_type"] = analysis.get("element_type", "TEXT_BLOCK")
                ctx["selected_element_text"] = analysis.get("element_text", "")
                ctx["contextual_actions"] = analysis.get("contextual_actions", [])
            except Exception as e:
                log.debug(f"MediaMaster farscry analysis failed: {e}")

        return ctx

    def build_prompt(
        self, user_instruction: str, context: dict, mem_context: str, classification: Any = None
    ) -> str:
        intent_classification = "Media Automation"
        if classification is not None:
            if isinstance(classification, str):
                intent_classification = classification
            elif isinstance(classification, dict):
                intent_classification = (
                    classification.get("task_type")
                    or classification.get("intent")
                    or classification.get("classification")
                    or "Media Automation"
                )
            else:
                intent_classification = (
                    getattr(classification, "task_type", None)
                    or getattr(classification, "intent", None)
                    or getattr(classification, "classification", None)
                    or "Media Automation"
                )

        app_ctx = """=== APP CONTEXT ===
Application Name: Canva / Adobe Photoshop
Application Type: Media Image / Video Editor"""

        doc_ctx = f"""=== DOCUMENT CONTEXT ===
Active App: {context.get('active_app')}
Scrubber Time: {context.get('timeline_scrubber_seconds')} seconds"""

        mem_ctx = f"""=== MEMORY CONTEXT ===
User Preferences: {mem_context or "None"}"""

        intent_part = f"""=== INTENT CLASSIFICATION ===
Intent: {intent_classification}"""

        system_rules = """You are the Media Master for Kairo Phantom. You automate media elements.
CRITICAL RULES:
1. Output ONLY valid JSON matching MediaResponse schema."""

        json_reminder = "REMINDER: Your entire response must be a single JSON object. First character must be {. Last character must be }."

        return f"""{system_rules}

{app_ctx}

{doc_ctx}

{mem_ctx}

{intent_part}

{json_reminder}
USER INSTRUCTION: {user_instruction}
OUTPUT (JSON only):
"""

    def validate_operations(self, response: MediaResponse, context: dict) -> List[dict]:
        return [response.model_dump()]

    def get_schema_class(self):
        return MediaResponse


# ==========================================
# 10. Data Master
# ==========================================
class DataMaster:
    @track("data", "extract_context")
    def extract_context(self, file_path: str, cursor_info: Any) -> dict:
        """
        Data Master context: Jupyter cell detection, pandas/SQL idiom awareness.
        Detects: active cell (Jupyter), kernel state, SQL dialect, R session.
        """
        ctx = {
            "notebook_cell_count": 1,
            "kernel_active": True,
            "imports": [],
            "sql_dialect": "generic",
            "language": "python",  # python | r | sql | scala
            "active_cell_content": "",
            "cursor_line": 0,
            "file_path": file_path or "",
            "data_libraries": [],  # pandas, numpy, sklearn, matplotlib, plotly, dask, etc.
        }

        if not file_path:
            return ctx

        try:
            fp_lower = str(file_path).lower()

            # Jupyter notebook
            if fp_lower.endswith(".ipynb"):
                ctx["language"] = "python"
                try:
                    import json as _json

                    if os.path.exists(file_path):
                        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                            nb = _json.load(f)
                        cells = nb.get("cells", [])
                        ctx["notebook_cell_count"] = len(cells)
                        # Extract imports from first code cells
                        for cell in cells:
                            if cell.get("cell_type") == "code":
                                src = "".join(cell.get("source", []))
                                for line in src.split("\n"):
                                    if line.startswith("import ") or line.startswith("from "):
                                        lib = line.split()[1].split(".")[0]
                                        if lib not in ctx["imports"]:
                                            ctx["imports"].append(lib)
                except Exception as e:
                    log.debug(f"DataMaster: Jupyter parse failed: {e}")

            # R script
            elif fp_lower.endswith((".r", ".rmd", ".rmarkdown")):
                ctx["language"] = "r"
                ctx["sql_dialect"] = "generic"

            # SQL file
            elif fp_lower.endswith(".sql"):
                ctx["language"] = "sql"
                if os.path.exists(file_path):
                    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                        sql_content = f.read(1000)
                    # Detect SQL dialect from comments or syntax
                    if "@@" in sql_content or "TOP " in sql_content.upper():
                        ctx["sql_dialect"] = "T-SQL"
                    elif "ROWNUM" in sql_content.upper():
                        ctx["sql_dialect"] = "Oracle"
                    elif "LIMIT" in sql_content.upper():
                        ctx["sql_dialect"] = "PostgreSQL"
                    else:
                        ctx["sql_dialect"] = "ANSI SQL"

            # Python files
            elif fp_lower.endswith(".py"):
                ctx["language"] = "python"
                if os.path.exists(file_path):
                    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                        content = f.read(2000)
                    for line in content.split("\n"):
                        stripped = line.strip()
                        if stripped.startswith("import ") or stripped.startswith("from "):
                            lib = stripped.split()[1].split(".")[0]
                            if lib not in ctx["imports"]:
                                ctx["imports"].append(lib)

            # Detect data libraries from imports
            DATA_LIBS = {
                "pandas",
                "numpy",
                "sklearn",
                "scikit_learn",
                "matplotlib",
                "plotly",
                "seaborn",
                "dask",
                "pyspark",
                "tensorflow",
                "torch",
                "keras",
                "xgboost",
                "lightgbm",
                "statsmodels",
                "scipy",
            }
            ctx["data_libraries"] = [lib for lib in ctx["imports"] if lib.lower() in DATA_LIBS]

            # Set cursor line from cursor_info
            if isinstance(cursor_info, (int, str)):
                try:
                    ctx["cursor_line"] = int(cursor_info)
                except (ValueError, TypeError):
                    pass

        except Exception as e:
            log.warning(f"DataMaster.extract_context failed: {e}")

        return ctx

    def build_prompt(
        self, user_instruction: str, context: dict, mem_context: str, classification: Any = None
    ) -> str:
        language = context.get("language", "python")
        sql_dialect = context.get("sql_dialect", "generic")
        data_libs = context.get("data_libraries", [])
        imports = context.get("imports", [])
        notebook_cells = context.get("notebook_cell_count", 1)
        cursor_line = context.get("cursor_line", 0)
        file_path = context.get("file_path", "Unknown")

        # Language-specific rules
        lang_rules = {
            "python": "Use pandas/numpy idioms. Avoid deprecated APIs. Prefer vectorized operations over loops.",
            "r": "Use tidyverse (dplyr/ggplot2/tidyr). Pipe operator |> preferred. Use tibble not data.frame.",
            "sql": f"Use {sql_dialect} syntax exactly. No cross-dialect functions. Quote identifiers with double quotes.",
        }.get(language, "Use idiomatic code for the detected language.")

        intent_classification = "Data Analysis"
        if classification is not None:
            if isinstance(classification, str):
                intent_classification = classification
            elif isinstance(classification, dict):
                intent_classification = (
                    classification.get("task_type")
                    or classification.get("intent")
                    or classification.get("classification")
                    or "Data Analysis"
                )
            else:
                intent_classification = (
                    getattr(classification, "task_type", None)
                    or getattr(classification, "intent", None)
                    or getattr(classification, "classification", None)
                    or "Data Analysis"
                )

        app_ctx = f"""=== APP CONTEXT ===
Application Name: Jupyter Notebook / DBeaver / RStudio
Application Type: Data Analysis Editor
Language: {language}
SQL Dialect: {sql_dialect}
File: {file_path}
Cursor Line: {cursor_line}"""

        doc_ctx = f"""=== DOCUMENT CONTEXT ===
Notebook Cells: {notebook_cells}
Kernel Active: {context.get('kernel_active', True)}
Imports in scope: {', '.join(imports[:10]) if imports else 'None detected'}
Data libraries detected: {', '.join(data_libs) if data_libs else 'None'}"""

        mem_ctx_block = f"""=== MEMORY CONTEXT ===
User Data Preferences: {mem_context or "None"}"""

        intent_part = f"""=== INTENT CLASSIFICATION ===
Intent: {intent_classification}"""

        system_rules = f"""You are the Data Master for Kairo Phantom. You generate precise SQL queries, pandas operations, and data science code.

LANGUAGE RULES ({language}):
{lang_rules}

CRITICAL RULES:
1. Output ONLY valid JSON matching DataResponse schema.
2. NEVER use functions not available in the detected dialect/version.
3. NEVER use SQL syntax from a different dialect (e.g., no MySQL-specific syntax for T-SQL).
4. Pandas operations must not accidentally modify the source DataFrame — use .copy() when needed.
5. Cell output exceeding 1000 rows must include a .head() or LIMIT clause."""

        json_reminder = "REMINDER: Your entire response must be a single JSON object. First character must be {. Last character must be }."

        return f"""{system_rules}

{app_ctx}

{doc_ctx}

{mem_ctx_block}

{intent_part}

{json_reminder}
USER INSTRUCTION: {user_instruction}
OUTPUT (JSON only):
"""

    def validate_operations(self, response: DataResponse, context: dict) -> List[dict]:
        return [response.model_dump()]

    def get_schema_class(self):
        return DataResponse
