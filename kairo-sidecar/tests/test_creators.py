# skipif removed — tests are fully mocked, conftest provides win32com/pythoncom/os.startfile stubs
from unittest.mock import patch
import tempfile
import os
from docx import Document
import openpyxl
from pptx import Presentation

from sidecar.creators.docx_creator import DocxCreator
from sidecar.creators.xlsx_creator import XlsxCreator
from sidecar.creators.pptx_creator import PptxCreator


def test_docx_creator():
    creator = DocxCreator()
    content = {
        "title": "Test Document",
        "author": "Test Author",
        "sections": [
            {
                "heading": "Heading 1",
                "level": 1,
                "paragraphs": ["Paragraph 1", "Paragraph 2"],
                "bullets": ["Bullet 1"],
                "table": {"headers": ["Col 1", "Col 2"], "rows": [["Row 1 Col 1", "Row 1 Col 2"]]},
            }
        ],
    }

    with tempfile.TemporaryDirectory() as tmpdir:
        out_path = os.path.join(tmpdir, "test.docx")

        # Test creation
        path = creator.create(content, output_path=out_path)
        assert path == out_path
        assert os.path.exists(out_path)

        # Verify document content
        doc = Document(out_path)
        # title
        assert doc.paragraphs[0].text == "Test Document"
        # heading 1
        assert doc.paragraphs[1].text == "Heading 1"
        # paragraphs
        assert doc.paragraphs[2].text == "Paragraph 1"
        assert doc.paragraphs[3].text == "Paragraph 2"
        # bullets
        assert doc.paragraphs[4].text == "Bullet 1"
        # table
        assert len(doc.tables) == 1
        table = doc.tables[0]
        assert len(table.rows) == 2
        assert table.rows[0].cells[0].text == "Col 1"
        assert table.rows[0].cells[1].text == "Col 2"
        assert table.rows[1].cells[0].text == "Row 1 Col 1"
        assert table.rows[1].cells[1].text == "Row 1 Col 2"


def test_docx_creator_create_and_open():
    creator = DocxCreator()
    content = {"title": "Test Open"}

    with (
        patch("os.startfile") as mock_startfile,
        patch.object(creator, "create", return_value="dummy_path") as mock_create,
    ):
        path = creator.create_and_open(content)
        assert path == "dummy_path"
        mock_create.assert_called_once_with(content)
        mock_startfile.assert_called_once_with("dummy_path")


def test_xlsx_creator():
    creator = XlsxCreator()
    content = {
        "title": "Test Excel",
        "author": "Test Author",
        "sheets": [
            {
                "name": "Summary Sheet",
                "headers": ["Metric", "Value"],
                "rows": [["Revenue", 100], ["Cost", 80]],
                "totals": True,
                "cells": [{"cell": "C1", "value": "Ad-hoc"}],
            }
        ],
    }

    with tempfile.TemporaryDirectory() as tmpdir:
        out_path = os.path.join(tmpdir, "test.xlsx")
        path = creator.create(content, output_path=out_path)
        assert path == out_path
        assert os.path.exists(out_path)

        # Verify content
        wb = openpyxl.load_workbook(out_path)
        assert "Summary Sheet" in wb.sheetnames
        ws = wb["Summary Sheet"]

        assert ws["A1"].value == "Metric"
        assert ws["B1"].value == "Value"
        assert ws["A2"].value == "Revenue"
        assert ws["B2"].value == 100
        assert ws["A3"].value == "Cost"
        assert ws["B3"].value == 80
        # Totals row
        assert ws["A4"].value == "TOTAL"
        assert ws["B4"].value == '=IFERROR(SUM(B2:B3),"")'
        # Ad-hoc cell
        assert ws["C1"].value == "Ad-hoc"


def test_xlsx_creator_create_and_open():
    creator = XlsxCreator()
    content = {"title": "Test Open"}

    with (
        patch("os.startfile") as mock_startfile,
        patch.object(creator, "create", return_value="dummy_path") as mock_create,
    ):
        path = creator.create_and_open(content)
        assert path == "dummy_path"
        mock_create.assert_called_once_with(content)
        mock_startfile.assert_called_once_with("dummy_path")


def test_pptx_creator():
    creator = PptxCreator()
    content = {
        "title": "Test Presentation",
        "author": "Test Author",
        "slides": [
            {"layout": "title", "title": "Main Title", "subtitle": "Sub Title"},
            {"layout": "content", "title": "Content Title", "bullets": ["Bullet 1", "Bullet 2"]},
        ],
    }

    with tempfile.TemporaryDirectory() as tmpdir:
        out_path = os.path.join(tmpdir, "test.pptx")
        path = creator.create(content, output_path=out_path)
        assert path == out_path
        assert os.path.exists(out_path)

        # Verify content
        prs = Presentation(out_path)
        assert len(prs.slides) == 2
        # Title slide
        slide1 = prs.slides[0]
        assert slide1.shapes.title.text == "Main Title"
        assert slide1.placeholders[1].text == "Sub Title"
        # Content slide
        slide2 = prs.slides[1]
        assert slide2.shapes.title.text == "Content Title"
        # bullets text in text frame
        tf = slide2.placeholders[1].text_frame
        assert tf.paragraphs[0].text == "Bullet 1"
        assert tf.paragraphs[1].text == "Bullet 2"


def test_pptx_creator_create_and_open():
    creator = PptxCreator()
    content = {"title": "Test Open"}

    with (
        patch("os.startfile") as mock_startfile,
        patch.object(creator, "create", return_value="dummy_path") as mock_create,
    ):
        path = creator.create_and_open(content)
        assert path == "dummy_path"
        mock_create.assert_called_once_with(content)
        mock_startfile.assert_called_once_with("dummy_path")
