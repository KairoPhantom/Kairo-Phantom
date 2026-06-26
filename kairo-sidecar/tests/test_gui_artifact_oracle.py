import os
import sys
import json
import tempfile
import subprocess
import pytest
import docx
import openpyxl
import pptx

# Add parent directory of kairo-sidecar to sys.path to allow importing scripts
sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..")))

from scripts.gui_artifact_oracle import (
    verify_docx,
    verify_xlsx,
    verify_pptx,
    verify_txt,
    verify_browser,
    normalize_text,
    normalize_xlsx_value,
)


@pytest.fixture
def temp_dir():
    with tempfile.TemporaryDirectory() as td:
        yield td


def test_normalize_text():
    assert normalize_text("  hello   world  ") == "hello world"
    assert normalize_text("unicode\u00a0normalization") == "unicode normalization"


def test_normalize_xlsx_value():
    assert normalize_xlsx_value(12.34567) == "12.3457"
    assert normalize_xlsx_value(10.0) == "10"
    assert normalize_xlsx_value("  text  ") == "text"
    assert normalize_xlsx_value(True) == "True"
    assert normalize_xlsx_value(None) == ""


# ──────────────────────────────────────────────────────────────────────────────
# Word / DOCX Tests
# ──────────────────────────────────────────────────────────────────────────────


def test_docx_pass(temp_dir):
    path = os.path.join(temp_dir, "test.docx")
    doc = docx.Document()
    doc.add_paragraph("Paragraph 0")
    p1 = doc.add_paragraph("Kickoff Meeting Minutes")
    p1.style = "Heading 1"
    doc.add_paragraph("Paragraph 2")
    doc.save(path)

    # Test insert_paragraph check
    expected = {
        "action": "insert_paragraph",
        "after_paragraph_index": 0,
        "text": "Kickoff Meeting Minutes",
        "style": "Heading 1",
    }
    assert verify_docx(path, expected) is True

    # Test append_to_run check
    expected_append = {"action": "append_to_run", "paragraph_index": 1, "text": "Meeting Minutes"}
    assert verify_docx(path, expected_append) is True


def test_docx_fail_wrong_content(temp_dir):
    path = os.path.join(temp_dir, "test.docx")
    doc = docx.Document()
    doc.add_paragraph("Paragraph 0")
    doc.save(path)

    expected = {
        "action": "insert_paragraph",
        "after_paragraph_index": 0,
        "text": "Kickoff Meeting Minutes",
    }
    with pytest.raises(AssertionError) as exc_info:
        verify_docx(path, expected)
    assert "mismatch" in str(exc_info.value) or "Expected" in str(exc_info.value)


def test_docx_fail_wrong_style(temp_dir):
    path = os.path.join(temp_dir, "test.docx")
    doc = docx.Document()
    p = doc.add_paragraph("Kickoff Meeting Minutes")
    p.style = "Normal"
    doc.save(path)

    expected = {
        "action": "insert_paragraph",
        "paragraph_index": 0,
        "text": "Kickoff Meeting",
        "style": "Heading 1",
    }
    with pytest.raises(AssertionError) as exc_info:
        verify_docx(path, expected)
    assert "style mismatch" in str(exc_info.value)


def test_docx_fail_missing_content(temp_dir):
    path = os.path.join(temp_dir, "test.docx")
    doc = docx.Document()
    doc.add_paragraph("Paragraph 0")
    doc.save(path)

    # Paragraph index out of bounds
    expected = {
        "action": "insert_paragraph",
        "after_paragraph_index": 5,
        "text": "Kickoff Meeting Minutes",
    }
    with pytest.raises(AssertionError) as exc_info:
        verify_docx(path, expected)
    assert "paragraphs" in str(exc_info.value)


# ──────────────────────────────────────────────────────────────────────────────
# Excel / XLSX Tests
# ──────────────────────────────────────────────────────────────────────────────


def test_xlsx_pass(temp_dir):
    path = os.path.join(temp_dir, "test.xlsx")
    wb = openpyxl.Workbook()
    sheet = wb.active
    sheet["E1"] = "Status"
    sheet["E2"] = "=SUM(A1:D1)"
    sheet["A1"] = "Q1"
    sheet["B1"] = "Q2"
    sheet["A2"] = 100
    sheet["B2"] = 200
    wb.save(path)

    # Test write_cell value
    expected_val = {"action": "write_cell", "cell": "E1", "value": "Status"}
    assert verify_xlsx(path, expected_val) is True

    # Test write_cell formula
    expected_formula = {"action": "write_cell", "cell": "E2", "formula": "=SUM(A1:D1)"}
    assert verify_xlsx(path, expected_formula) is True

    # Test write_range matrix
    expected_range = {
        "action": "write_range",
        "range": "A1:B2",
        "values": [["Q1", "Q2"], [100, 200]],
    }
    assert verify_xlsx(path, expected_range) is True


def test_xlsx_fail_wrong_value(temp_dir):
    path = os.path.join(temp_dir, "test.xlsx")
    wb = openpyxl.Workbook()
    sheet = wb.active
    sheet["E1"] = "Status"
    wb.save(path)

    expected = {"action": "write_cell", "cell": "E1", "value": "Incorrect"}
    with pytest.raises(AssertionError) as exc_info:
        verify_xlsx(path, expected)
    assert "Value mismatch" in str(exc_info.value)


def test_xlsx_fail_wrong_formula(temp_dir):
    path = os.path.join(temp_dir, "test.xlsx")
    wb = openpyxl.Workbook()
    sheet = wb.active
    sheet["E2"] = "=SUM(A1:D1)"
    wb.save(path)

    expected = {"action": "write_cell", "cell": "E2", "formula": "=AVERAGE(A1:D1)"}
    with pytest.raises(AssertionError) as exc_info:
        verify_xlsx(path, expected)
    assert "Formula mismatch" in str(exc_info.value)


def test_xlsx_fail_missing_cell(temp_dir):
    path = os.path.join(temp_dir, "test.xlsx")
    wb = openpyxl.Workbook()
    wb.save(path)

    expected = {"action": "write_cell", "cell": "E1", "value": "Status"}
    with pytest.raises(AssertionError) as exc_info:
        verify_xlsx(path, expected)
    assert "Value mismatch" in str(exc_info.value)


# ──────────────────────────────────────────────────────────────────────────────
# PowerPoint / PPTX Tests
# ──────────────────────────────────────────────────────────────────────────────


def test_pptx_pass(temp_dir):
    path = os.path.join(temp_dir, "test.pptx")
    prs = pptx.Presentation()
    # Add title slide layout
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(path)

    expected_slides = {"action": "create_presentation", "slides": 1}
    assert verify_pptx(path, expected_slides) is True

    expected_layout = {
        "action": "slide_layout",
        "cursor_slide": 0,
        "expected_layout": prs.slide_layouts[0].name,
    }
    assert verify_pptx(path, expected_layout) is True


def test_pptx_fail_wrong_count(temp_dir):
    path = os.path.join(temp_dir, "test.pptx")
    prs = pptx.Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(path)

    expected = {"action": "create_presentation", "slides": 5}
    with pytest.raises(AssertionError) as exc_info:
        verify_pptx(path, expected)
    assert "Slide count mismatch" in str(exc_info.value)


def test_pptx_fail_wrong_layout(temp_dir):
    path = os.path.join(temp_dir, "test.pptx")
    prs = pptx.Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(path)

    expected = {"action": "slide_layout", "cursor_slide": 0, "expected_layout": "Title and Content"}
    with pytest.raises(AssertionError) as exc_info:
        verify_pptx(path, expected)
    assert "layout mismatch" in str(exc_info.value)


# ──────────────────────────────────────────────────────────────────────────────
# Command Line Interface (CLI) / Subprocess Tests
# ──────────────────────────────────────────────────────────────────────────────


def test_cli_pass(temp_dir):
    # Setup artifact docx
    artifact_path = os.path.join(temp_dir, "test.docx")
    doc = docx.Document()
    p = doc.add_paragraph("Kickoff Meeting Minutes")
    p.style = "Heading 1"
    doc.save(artifact_path)

    # Setup scenarios JSON
    scenarios_json_path = os.path.join(temp_dir, "scenarios.json")
    scenario_data = [
        {
            "id": "WORD_016",
            "category": "Word",
            "name": "Word Scenario 16: Insert End Paragraph",
            "description": "...",
            "status": "active",
            "expected_outcome": {
                "action": "insert_paragraph",
                "after_paragraph_index": -1,
                "text": "Kickoff Meeting Minutes",
                "style": "Heading 1",
            },
        }
    ]
    with open(scenarios_json_path, "w", encoding="utf-8") as f:
        json.dump(scenario_data, f)

    cmd = [
        sys.executable,
        os.path.join("scripts", "gui_artifact_oracle.py"),
        "--scenario-id",
        "WORD_016",
        "--artifact-path",
        artifact_path,
        "--scenarios-json",
        scenarios_json_path,
    ]

    # We run the command in the root repository folder
    result = subprocess.run(
        cmd,
        capture_output=True,
        text=True,
        cwd=os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..")),
    )
    assert result.returncode == 0, f"STDOUT: {result.stdout}\nSTDERR: {result.stderr}"
    assert "PASS" in result.stdout


def test_cli_fail_wrong_content(temp_dir):
    # Setup artifact docx
    artifact_path = os.path.join(temp_dir, "test.docx")
    doc = docx.Document()
    doc.add_paragraph("Different text content")
    doc.save(artifact_path)

    # Setup scenarios JSON
    scenarios_json_path = os.path.join(temp_dir, "scenarios.json")
    scenario_data = [
        {
            "id": "WORD_016",
            "category": "Word",
            "name": "...",
            "status": "active",
            "expected_outcome": {"action": "insert_paragraph", "text": "Kickoff Meeting Minutes"},
        }
    ]
    with open(scenarios_json_path, "w", encoding="utf-8") as f:
        json.dump(scenario_data, f)

    cmd = [
        sys.executable,
        os.path.join("scripts", "gui_artifact_oracle.py"),
        "--scenario-id",
        "WORD_016",
        "--artifact-path",
        artifact_path,
        "--scenarios-json",
        scenarios_json_path,
    ]

    result = subprocess.run(
        cmd,
        capture_output=True,
        text=True,
        cwd=os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..")),
    )
    assert result.returncode == 1
    assert "FAIL" in result.stdout


def test_cli_fail_file_not_found(temp_dir):
    # Setup scenarios JSON
    scenarios_json_path = os.path.join(temp_dir, "scenarios.json")
    scenario_data = [{"id": "WORD_016", "expected_outcome": {}}]
    with open(scenarios_json_path, "w", encoding="utf-8") as f:
        json.dump(scenario_data, f)

    cmd = [
        sys.executable,
        os.path.join("scripts", "gui_artifact_oracle.py"),
        "--scenario-id",
        "WORD_016",
        "--artifact-path",
        "nonexistent_file.docx",
        "--scenarios-json",
        scenarios_json_path,
    ]

    result = subprocess.run(
        cmd,
        capture_output=True,
        text=True,
        cwd=os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..")),
    )
    assert result.returncode == 1
    assert "FAIL" in result.stdout or "not found" in result.stdout.lower()


def test_txt_pass(temp_dir):
    path = os.path.join(temp_dir, "test.txt")
    with open(path, "w", encoding="utf-8") as f:
        f.write("Buy milk\nWrite code\nRun tests\r\n")

    expected_write = {"action": "write_text", "contains": "Buy milk"}
    assert verify_txt(path, expected_write) is True

    expected_format = {"action": "format_text", "crlf": True}
    assert verify_txt(path, expected_format) is True


def test_txt_fail(temp_dir):
    path = os.path.join(temp_dir, "test.txt")
    with open(path, "w", encoding="utf-8") as f:
        f.write("Only some other text\n")

    expected = {"action": "write_text", "contains": "Buy milk"}
    with pytest.raises(AssertionError):
        verify_txt(path, expected)


def test_browser_pass():
    assert verify_browser("dummy_path", {}) is True
