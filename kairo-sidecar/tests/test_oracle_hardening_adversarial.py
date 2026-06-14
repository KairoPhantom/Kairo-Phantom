import os
import pytest
import builtins
import io
from unittest.mock import patch
import pptx
import fitz

from sidecar.test_fix_loop import TestFixLoop, ProtectedPathViolation
from sidecar.oracles import verify_pptx, verify_pdf

def test_verify_patch_safety_relative_bypasses():
    """Verify that relative path bypasses are successfully blocked by verify_patch_safety."""
    loop = TestFixLoop(workspace_root=".")
    
    # List of adversarial paths that should be blocked because they refer to protected paths
    blocked_paths = [
        "foo/../kairo-sidecar/sidecar/oracles.py",
        "sidecar/../sidecar/oracles.py",
        "./kairo-sidecar/sidecar/oracles.py",
        "kairo-sidecar/sidecar/oracles.py",
        "kairo-sidecar/sidecar/../sidecar/oracles.py",
        "foo/bar/../../kairo-sidecar/sidecar/oracles.py",
        "phantom-core/src/../src/response_validator.rs",
        "scripts/ci/../ci/sbom_gate.py",
        "kairo-sidecar/sidecar/oracles.py.sig",
        "kairo-sidecar/sidecar/oracles.py.pub",
        "kairo-sidecar/sidecar/test_fix_loop.py"
    ]
    
    for path in blocked_paths:
        with pytest.raises(ProtectedPathViolation) as excinfo:
            loop.verify_patch_safety({path})
        assert "is protected" in str(excinfo.value)
        
    # Unprotected paths should not raise anything
    safe_paths = [
        "kairo-sidecar/sidecar/oracles_backup.py",
        "foo/bar/baz.py",
        "tests/test_oracles.py",
        "kairo-sidecar/sidecar/oracles.py.sig.bak"
    ]
    loop.verify_patch_safety(set(safe_paths))


def test_verify_oracles_signature_tampered_pubkey():
    """Verify that modifying the public key on disk results in verification failure."""
    loop = TestFixLoop(workspace_root=".")
    original_open = builtins.open
    
    def mock_open_modified_pub(file, mode="r", *args, **kwargs):
        filepath = os.path.abspath(str(file)).replace("\\", "/")
        if filepath.endswith("sidecar/oracles.py.pub"):
            tampered_content = b"-----BEGIN PUBLIC KEY-----\nINVALIDKEY\n-----END PUBLIC KEY-----\n"
            if "b" in mode:
                return io.BytesIO(tampered_content)
            else:
                return io.StringIO(tampered_content.decode("utf-8"))
        return original_open(file, mode, *args, **kwargs)
        
    with patch("builtins.open", side_effect=mock_open_modified_pub):
        with pytest.raises(PermissionError) as excinfo:
            loop.verify_oracles_signature()
        assert "Public key on disk does not match pinned public key!" in str(excinfo.value)


def test_verify_oracles_signature_tampered_sig():
    """Verify that modifying the signature file on disk results in signature verification failure."""
    loop = TestFixLoop(workspace_root=".")
    original_open = builtins.open
    
    def mock_open_modified_sig(file, mode="r", *args, **kwargs):
        filepath = os.path.abspath(str(file)).replace("\\", "/")
        if filepath.endswith("sidecar/oracles.py.sig"):
            # A 64-byte signature with incorrect values
            return io.BytesIO(b"\x00" * 64)
        return original_open(file, mode, *args, **kwargs)
        
    with patch("builtins.open", side_effect=mock_open_modified_sig):
        with pytest.raises(PermissionError) as excinfo:
            loop.verify_oracles_signature()
        assert "Signature verification failed" in str(excinfo.value)


def test_binned_coordinate_sorting_pptx(tmp_path):
    """Verify that binned coordinate sorting correctly aligns adjacent text blocks in PPTX."""
    pptx_path = os.path.join(tmp_path, "binned_sort.pptx")
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # blank layout
    
    # We want Shape 2 (Right Text) to have a smaller exact y-coordinate than Shape 1 (Left Text).
    # Shape 2: top=99, left=200
    # Shape 1: top=101, left=50
    #
    # With binned sorting (bin=5), both top values round to 100, and they sort left-to-right (50 < 200).
    # So "Left Text" comes before "Right Text", resulting in "Left Text Right Text".
    #
    # Add Shape 2 first to verify that creation order doesn't dictate the sort.
    tx_right = slide.shapes.add_textbox(left=200, top=99, width=100, height=50)
    tx_right.text_frame.text = "Right Text"
    
    tx_left = slide.shapes.add_textbox(left=50, top=101, width=100, height=50)
    tx_left.text_frame.text = "Left Text"
    
    prs.save(pptx_path)
    
    # verify_pptx should sort them as Left Text -> Right Text
    assert verify_pptx(pptx_path, expected_text_substrings=["Left Text Right Text"])


def test_binned_coordinate_sorting_pdf(tmp_path):
    """Verify that binned coordinate sorting correctly aligns adjacent text blocks in PDF."""
    pdf_path = os.path.join(tmp_path, "binned_sort.pdf")
    doc = fitz.open()
    page = doc.new_page()
    
    # We want two separate text blocks that have y-coordinates rounding to the same 5-unit bin.
    # From empirical testing, inserting text at (50, 57) and (400, 53) with fontsize=5 splits
    # the text into two blocks with y0_0=51.625 (Left PDF) and y0_1=47.625 (Right PDF).
    # Both y0 values round to 50, and since 50.0 < 400.0, the Left PDF block comes first in binned sorting.
    # Insert the Left PDF text block first (creation order) to ensure the layout analyzer splits them.
    page.insert_text((50, 57), "Left PDF", fontsize=5)
    page.insert_text((400, 53), "Right PDF", fontsize=5)
    
    doc.save(pdf_path)
    doc.close()
    
    # Mock pdfplumber to bypass its layout-unaware extraction and isolate the test to fitz's binned sorting logic
    from unittest.mock import MagicMock
    with patch("pdfplumber.open") as mock_open:
        mock_pdf = MagicMock()
        mock_page = MagicMock()
        mock_page.extract_text.return_value = "Left PDF Right PDF"
        mock_pdf.pages = [mock_page]
        mock_open.return_value.__enter__.return_value = mock_pdf
        
        # verify_pdf should sort them as Left PDF -> Right PDF in fitz
        assert verify_pdf(pdf_path, ["Left PDF Right PDF"])


def test_word_limits_on_slide_level0_placeholders(tmp_path):
    """Verify that word limits are correctly enforced on slide level 0 placeholders (BODY and OBJECT)."""
    # 1. Test PP_PLACEHOLDER.BODY (type 2) placeholder
    pptx_body_path = os.path.join(tmp_path, "body_limit.pptx")
    prs = pptx.Presentation()
    # Layout 2 has TITLE (0) and BODY (1) placeholders
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    body_ph = slide.placeholders[1]
    
    # Body has 13 words, which exceeds the limit of 12
    body_ph.text_frame.text = "one two three four five six seven eight nine ten eleven twelve thirteen"
    prs.save(pptx_body_path)
    
    # verify_pptx should raise AssertionError due to BODY placeholder word limit
    with pytest.raises(AssertionError) as excinfo:
        verify_pptx(pptx_body_path, bullet_word_limit=12)
    assert "bullet exceeds word limit" in str(excinfo.value)
    
    # 2. Test PP_PLACEHOLDER.OBJECT (type 7) placeholder
    pptx_obj_path = os.path.join(tmp_path, "obj_limit.pptx")
    prs = pptx.Presentation()
    # Layout 1 has TITLE (0) and OBJECT (1) placeholders
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    obj_ph = slide.placeholders[1]
    
    # Object has 13 words, which exceeds the limit of 12
    obj_ph.text_frame.text = "one two three four five six seven eight nine ten eleven twelve thirteen"
    prs.save(pptx_obj_path)
    
    # verify_pptx should raise AssertionError due to OBJECT placeholder word limit
    with pytest.raises(AssertionError) as excinfo:
        verify_pptx(pptx_obj_path, bullet_word_limit=12)
    assert "bullet exceeds word limit" in str(excinfo.value)
    
    # 3. Test TITLE (type 1) placeholder (should NOT be subject to word limit)
    pptx_title_path = os.path.join(tmp_path, "title_no_limit.pptx")
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_ph = slide.placeholders[0]
    
    title_ph.text_frame.text = "one two three four five six seven eight nine ten eleven twelve thirteen"
    # Ensure other placeholders are empty or within limit
    slide.placeholders[1].text_frame.text = "short text"
    prs.save(pptx_title_path)
    
    # This should pass without raising any AssertionError
    assert verify_pptx(pptx_title_path, bullet_word_limit=12)
    
    # 4. Test regular textbox (non-placeholder) (should NOT be subject to word limit if not starting with '-')
    pptx_box_path = os.path.join(tmp_path, "textbox_no_limit.pptx")
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # blank layout
    tx_box = slide.shapes.add_textbox(0, 0, 100, 100)
    tx_box.text_frame.text = "one two three four five six seven eight nine ten eleven twelve thirteen"
    prs.save(pptx_box_path)
    
    # This should pass without raising any AssertionError
    assert verify_pptx(pptx_box_path, bullet_word_limit=12)
    
    # 5. Test regular textbox starting with '-' (should be treated as bullet and raise AssertionError)
    pptx_dash_path = os.path.join(tmp_path, "textbox_dash_limit.pptx")
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tx_box = slide.shapes.add_textbox(0, 0, 100, 100)
    tx_box.text_frame.text = "- one two three four five six seven eight nine ten eleven twelve thirteen"
    prs.save(pptx_dash_path)
    
    with pytest.raises(AssertionError) as excinfo:
        verify_pptx(pptx_dash_path, bullet_word_limit=12)
    assert "bullet exceeds word limit" in str(excinfo.value)
