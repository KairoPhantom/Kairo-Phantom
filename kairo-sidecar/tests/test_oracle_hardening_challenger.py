# tests/test_oracle_hardening_challenger.py
import os
import io
import builtins
import pytest
import pptx
import fitz
from unittest.mock import patch, MagicMock
from sidecar.test_fix_loop import TestFixLoop, ProtectedPathViolation
from sidecar.oracles import verify_pptx, verify_pdf

@pytest.fixture
def temp_dir(tmpdir):
    return str(tmpdir)

# 1. Path Safety Verification Tests
def test_verify_patch_safety_bypasses():
    loop = TestFixLoop(workspace_root="/dummy/workspace")
    
    # These should be successfully BLOCKED
    blocked_paths = [
        "kairo-sidecar/sidecar/oracles.py",
        "foo/../kairo-sidecar/sidecar/oracles.py",
        "./kairo-sidecar/sidecar/oracles.py",
        "kairo-sidecar/sidecar/./oracles.py",
        "kairo-sidecar/sidecar/oracles.py/.",
        "foo/bar/../../kairo-sidecar/sidecar/oracles.py",
        "phantom-core/src/response_validator.rs",
        "foo/../phantom-core/src/response_validator.rs",
        "kairo-sidecar/sidecar/oracles.py.sig",
        "kairo-sidecar/sidecar/oracles.py.pub",
        "kairo-sidecar/sidecar/test_fix_loop.py",
        "scripts/ci/sbom_gate.py",
        "scripts/ci/eval_integrity_guard.py"
    ]
    
    for path in blocked_paths:
        with pytest.raises(ProtectedPathViolation):
            loop.verify_patch_safety({path})

    # These should NOT be blocked (no false positives)
    allowed_paths = [
        "kairo-sidecar/sidecar/oracles_backup.py",
        "kairo-sidecar/sidecar/not_oracles.py",
        "phantom-core/src/response_validator.rs.bak",
        "scripts/ci/sbom_gate.py.tmp",
        "kairo-sidecar/sidecar/oracles.py.bak"
    ]
    
    for path in allowed_paths:
        # Should not raise exception
        loop.verify_patch_safety({path})


# 2. Cryptographic signature and public key on disk modification tests
def test_oracles_signature_verification_failures():
    loop = TestFixLoop(workspace_root="/dummy/workspace")
    
    # We will mock builtins.open to intercept reads to oracles.py, oracles.py.pub, and oracles.py.sig
    original_open = builtins.open
    
    # Mocking standard Ed25519 public key bytes that is DIFFERENT from the pinned key
    different_pub_key = (
        b"-----BEGIN PUBLIC KEY-----\n"
        b"MCowBQYDK2VwAyEA9999999999999999999999999999999999999999990=\n"
        b"-----END PUBLIC KEY-----\n"
    )

    # Test Case 2a: modifying the public key on disk to a different one
    def mock_open_diff_pub(file, mode="r", *args, **kwargs):
        filename = os.path.basename(str(file))
        if filename == "oracles.py.pub" and "b" in mode:
            return io.BytesIO(different_pub_key)
        return original_open(file, mode, *args, **kwargs)
        
    with patch("builtins.open", side_effect=mock_open_diff_pub):
        with pytest.raises(PermissionError) as excinfo:
            loop.verify_oracles_signature()
        assert "Public key on disk does not match pinned public key" in str(excinfo.value)

    # Test Case 2b: modifying the signature file on disk
    def mock_open_bad_sig(file, mode="r", *args, **kwargs):
        filename = os.path.basename(str(file))
        if filename == "oracles.py.sig" and "b" in mode:
            return io.BytesIO(b"bad_signature_bytes_here" * 5)
        return original_open(file, mode, *args, **kwargs)

    with patch("builtins.open", side_effect=mock_open_bad_sig):
        with pytest.raises(PermissionError) as excinfo:
            loop.verify_oracles_signature()
        assert "Signature verification failed" in str(excinfo.value)

    # Test Case 2c: modifying the oracle file itself
    def mock_open_bad_oracles(file, mode="r", *args, **kwargs):
        filename = os.path.basename(str(file))
        if filename == "oracles.py" and "b" in mode:
            with original_open(file, mode, *args, **kwargs) as f:
                content = f.read()
            return io.BytesIO(content + b"\n# malicious append")
        return original_open(file, mode, *args, **kwargs)

    with patch("builtins.open", side_effect=mock_open_bad_oracles):
        with pytest.raises(PermissionError) as excinfo:
            loop.verify_oracles_signature()
        assert "Signature verification failed" in str(excinfo.value)

    # Test Case 2d: Missing oracles.py
    with patch("os.path.exists", side_effect=lambda path: False if "oracles.py" in os.path.basename(path) and not path.endswith(".py.sig") and not path.endswith(".py.pub") else True):
        with pytest.raises(PermissionError) as excinfo:
            loop.verify_oracles_signature()
        assert "oracles.py is missing" in str(excinfo.value)

    # Test Case 2e: Missing oracles.py.sig
    with patch("os.path.exists", side_effect=lambda path: False if "oracles.py.sig" in os.path.basename(path) else True):
        with pytest.raises(PermissionError) as excinfo:
            loop.verify_oracles_signature()
        assert "Signature oracles.py.sig is missing" in str(excinfo.value)


# 3. Binned Coordinate Sorting Tests
def test_pptx_binned_coordinate_sorting(temp_dir):
    pptx_path = os.path.join(temp_dir, "binned_sort.pptx")
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add textbox A visually at top=50, left=200
    tx_a = slide.shapes.add_textbox(200, 50, 100, 100)
    tx_a.text_frame.text = "A"
    
    # Add textbox B visually at top=52, left=100 (B is to the left of A, but slightly lower)
    tx_b = slide.shapes.add_textbox(100, 52, 100, 100)
    tx_b.text_frame.text = "B"
    
    prs.save(pptx_path)
    
    # Sorting by binned top:
    # A top=50 -> binned to 50
    # B top=52 -> binned to 50
    # Since they are binned to the same value, they must be sorted by left:
    # B (left=100) comes before A (left=200).
    # Thus, expected substring is "B A" (or both in that order).
    assert verify_pptx(pptx_path, expected_slide_count=1, expected_text_substrings=["B A"])


def test_pdf_binned_coordinate_sorting(temp_dir):
    pdf_path = os.path.join(temp_dir, "binned_sort.pdf")
    
    mock_page = MagicMock()
    mock_page.get_text.return_value = [
        (200, 50, 250, 60, "A", 0, 0),
        (100, 52, 150, 62, "B", 1, 0)
    ]
    
    mock_doc = MagicMock()
    mock_doc.__iter__.return_value = [mock_page]
    mock_doc.__enter__.return_value = mock_doc
    
    mock_plumber_page = MagicMock()
    mock_plumber_page.extract_text.return_value = "B A"
    
    mock_plumber_doc = MagicMock()
    mock_plumber_doc.pages = [mock_plumber_page]
    mock_plumber_doc.__enter__.return_value = mock_plumber_doc
    
    with patch("fitz.open", return_value=mock_doc), \
         patch("pdfplumber.open", return_value=mock_plumber_doc):
         
        assert verify_pdf(pdf_path, ["B A"])



# 4. Word Limit Enforcement on Slide Level 0 Placeholders
def test_pptx_word_limit_enforcement(temp_dir):
    # Test layout 1 (Title and Content)
    # Typically: placeholder index 0 is TITLE, index 1 is BODY
    pptx_path = os.path.join(temp_dir, "word_limit_test.pptx")
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    body_placeholder = slide.placeholders[1]
    title_placeholder = slide.placeholders[0]
    
    # Check that body placeholder is actually PP_PLACEHOLDER.BODY or PP_PLACEHOLDER.OBJECT
    from pptx.enum.shapes import PP_PLACEHOLDER
    assert body_placeholder.placeholder_format.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT)
    
    # 4a: Body placeholder level 0 text exceeding limit (13 words)
    tf = body_placeholder.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "one two three four five six seven eight nine ten eleven twelve thirteen"
    p.level = 0
    prs.save(pptx_path)
    
    with pytest.raises(AssertionError) as excinfo:
        verify_pptx(pptx_path, bullet_word_limit=12)
    assert "exceeds word limit" in str(excinfo.value)
    
    # 4b: Body placeholder level 0 text within limit (12 words)
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "one two three four five six seven eight nine ten eleven twelve"
    p.level = 0
    prs.save(pptx_path)
    assert verify_pptx(pptx_path, bullet_word_limit=12)
    
    # 4c: Title placeholder (not BODY/OBJECT) exceeding limit does not raise AssertionError
    tf_title = title_placeholder.text_frame
    tf_title.clear()
    p_title = tf_title.paragraphs[0]
    p_title.text = "one two three four five six seven eight nine ten eleven twelve thirteen fourteen fifteen"
    p_title.level = 0
    prs.save(pptx_path)
    assert verify_pptx(pptx_path, bullet_word_limit=12)
    
    # 4d: Title placeholder exceeding limit with level > 0 DOES raise AssertionError
    p_title.level = 1
    prs.save(pptx_path)
    with pytest.raises(AssertionError) as excinfo:
        verify_pptx(pptx_path, bullet_word_limit=12)
    assert "exceeds word limit" in str(excinfo.value)
    
    # 4e: Title placeholder exceeding limit starting with '-' DOES raise AssertionError
    p_title.level = 0
    p_title.text = "- one two three four five six seven eight nine ten eleven twelve thirteen"
    prs.save(pptx_path)
    with pytest.raises(AssertionError) as excinfo:
        verify_pptx(pptx_path, bullet_word_limit=12)
    assert "exceeds word limit" in str(excinfo.value)

    # 4f: Normal text box (not placeholder) exceeding limit does not raise AssertionError
    # Let's clean title placeholder first so it doesn't fail
    p_title.text = "Short title"
    # Clean body placeholder
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Short body"
    p.level = 0
    
    txBox = slide.shapes.add_textbox(0, 0, 100, 100)
    tf_box = txBox.text_frame
    p_box = tf_box.paragraphs[0]
    p_box.text = "one two three four five six seven eight nine ten eleven twelve thirteen fourteen fifteen"
    p_box.level = 0
    prs.save(pptx_path)
    assert verify_pptx(pptx_path, bullet_word_limit=12)


# 5. Case-insensitivity Bypass Verification (Vulnerability Detection)
def test_verify_patch_safety_case_insensitive_bypass():
    loop = TestFixLoop(workspace_root="/dummy/workspace")
    
    # On case-insensitive filesystems (like Windows), ORACLES.py refers to the same file as oracles.py.
    # However, verify_patch_safety matches case-sensitively, so "ORACLES.py" is NOT blocked.
    # We assert that the bypass succeeds (which confirms the presence of the vulnerability).
    try:
        loop.verify_patch_safety({"kairo-sidecar/sidecar/ORACLES.py"})
        bypass_succeeded = True
    except ProtectedPathViolation:
        bypass_succeeded = False
        
    assert bypass_succeeded, "Case-insensitive bypass failed (it should succeed due to the bug!)"

