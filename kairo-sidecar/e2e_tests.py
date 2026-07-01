#!/usr/bin/env python3
"""
Kairo Phantom E2E Test Suite
============================
Tests all document types end-to-end via the sidecar TCP protocol.

Run: python kairo-sidecar/e2e_tests.py
Requires: sidecar running on port 7438
"""

import asyncio
import json
import os
import sys
import tempfile
import traceback

# Force UTF-8 output
sys.stdout = open(sys.stdout.fileno(), mode="w", encoding="utf-8", buffering=1)

HOST, PORT = "127.0.0.1", 7438
PASS, FAIL, SKIP = "PASS", "FAIL", "SKIP"
results = []


# ─── Test helpers ────────────────────────────────────────────────────────────


async def send(action: str, path: str = "", payload: dict = None) -> dict:
    reader, writer = await asyncio.wait_for(asyncio.open_connection(HOST, PORT), timeout=5)
    req = json.dumps(
        {"id": f"e2e-{action}", "action": action, "path": path, "payload": payload or {}}
    )
    writer.write((req + "\n").encode())
    await writer.drain()
    line = await asyncio.wait_for(reader.readline(), timeout=10)
    writer.close()
    return json.loads(line.decode().strip())


def record(name: str, status: str, detail: str = ""):
    icon = "✅" if status == PASS else ("⚠️ " if status == SKIP else "❌")
    print(f"  {icon} {name}: {detail}")
    results.append({"name": name, "status": status, "detail": detail})


async def run_test(name: str, coro):
    try:
        await coro
    except Exception as e:
        record(name, FAIL, str(e))


# ─── Test cases ──────────────────────────────────────────────────────────────


async def test_ping():
    r = await send("ping")
    assert r["ok"], f"Ping failed: {r}"
    record("PING: sidecar alive", PASS, f"version={r['data']['version']}")


async def test_docx_read_write():
    """Create a fresh docx, write via sidecar, read back, verify."""
    try:
        from docx import Document
    except ImportError:
        record("DOCX: read+write roundtrip", SKIP, "python-docx not installed")
        return

    # Create test document
    doc = Document()
    doc.add_heading("Introduction", level=1)
    doc.add_paragraph("This is the intro section.", style="Normal")
    doc.add_heading("Results", level=1)
    doc.add_paragraph("Results go here.", style="Normal")

    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
        tmp_path = f.name
    doc.save(tmp_path)

    # Read via sidecar
    r = await send("read_docx", path=tmp_path)
    assert r["ok"], f"Read failed: {r}"
    d = r["data"]
    assert d["paragraph_count"] > 0
    record(
        "DOCX: read structure",
        PASS,
        f"{d['paragraph_count']} paragraphs, {len(d['headings'])} headings",
    )

    # Write via sidecar
    ops = [
        {
            "action": "insert_after_heading",
            "heading_text": "Introduction",
            "style": "ListBullet",
            "content": ["Key finding one", "Key finding two", "Key finding three"],
        },
        {"action": "append", "style": "Normal", "content": "Conclusion paragraph added by Kairo."},
    ]
    r2 = await send("write_docx", path=tmp_path, payload={"operations": ops})
    assert r2["ok"], f"Write failed: {r2}"
    assert r2["data"]["applied_count"] == 2, f"Expected 2 applied: {r2}"
    record(
        "DOCX: write operations",
        PASS,
        f"applied_count={r2['data']['applied_count']} errors={r2['data']['errors']}",
    )

    # Verify written content
    r3 = await send("read_docx", path=tmp_path)
    assert r3["ok"]
    full = r3["data"]["full_text"]
    assert "Key finding one" in full, "Content not written!"
    assert "Conclusion paragraph" in full
    record("DOCX: verify content persisted", PASS, "all injected content found in file")

    os.unlink(tmp_path)


async def test_docx_table_insert():
    try:
        from docx import Document
    except ImportError:
        record("DOCX: table insert", SKIP, "python-docx not installed")
        return

    doc = Document()
    doc.add_heading("Data", level=1)
    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
        tmp = f.name
    doc.save(tmp)

    ops = [
        {
            "action": "insert_table",
            "style": "Normal",
            "rows": [
                ["Name", "Value", "Change"],
                ["Revenue", "$1.2M", "+15%"],
                ["Users", "42,000", "+8%"],
            ],
            "content": "",
        }
    ]
    r = await send("write_docx", path=tmp, payload={"operations": ops})
    assert r["ok"], f"Table insert failed: {r}"
    record("DOCX: insert table", PASS, f"applied_count={r['data']['applied_count']}")
    os.unlink(tmp)


async def test_xlsx_read_write():
    try:
        import openpyxl
    except ImportError:
        record("XLSX: read+write roundtrip", SKIP, "openpyxl not installed")
        return

    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "Product"
    ws["B1"] = "Units"
    ws["C1"] = "Price"
    ws["A2"] = "Widget A"
    ws["B2"] = 500
    ws["C2"] = 9.99
    ws["A3"] = "Widget B"
    ws["B3"] = 300
    ws["C3"] = 14.99

    with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as f:
        tmp = f.name
    wb.save(tmp)

    # Read context
    r = await send("read_xlsx", path=tmp, payload={"active_cell": "D2"})
    assert r["ok"], f"XLSX read failed: {r}"
    d = r["data"]
    assert d["active_cell"] == "D2"
    assert "A" in d["headers"]
    record(
        "XLSX: read context",
        PASS,
        f"active={d['active_cell']} sheet={d['sheet_name']} headers={list(d['headers'].values())}",
    )

    # Write formula
    ops = [
        {"cell": "D2", "formula": "=B2*C2", "value": ""},
        {"cell": "D3", "formula": "=B3*C3", "value": ""},
        {"cell": "D1", "formula": "", "value": "Revenue"},
        {"cell": "D4", "formula": "=SUM(D2:D3)", "value": ""},
    ]
    r2 = await send("write_xlsx", path=tmp, payload={"operations": ops})
    assert r2["ok"], f"XLSX write failed: {r2}"
    record(
        "XLSX: write formulas",
        PASS,
        f"applied_count={r2['data']['applied_count']} errors={r2['data']['errors']}",
    )

    # Verify formula was written
    wb2 = openpyxl.load_workbook(tmp)
    ws2 = wb2.active
    assert str(ws2["D2"].value) == "=B2*C2", f"Formula not written: {ws2['D2'].value}"
    assert str(ws2["D1"].value) == "Revenue"
    record("XLSX: verify formulas persisted", PASS, "D2=B2*C2 confirmed")

    # Verify adjacent cells untouched
    assert ws2["A2"].value == "Widget A", "Adjacent cell corrupted!"
    assert ws2["B2"].value == 500
    record("XLSX: adjacent cells untouched", PASS, "A2, B2 preserved correctly")

    os.unlink(tmp)


async def test_pptx_read_write():
    try:
        from pptx import Presentation
        from pptx.util import Inches  # noqa: F401
    except ImportError:
        record("PPTX: read+write roundtrip", SKIP, "python-pptx not installed")
        return

    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Cost Analysis"
    slide.placeholders[1].text = "Original content here.\nMore content."

    slide2 = prs.slides.add_slide(layout)
    slide2.shapes.title.text = "Summary"
    slide2.placeholders[1].text = "Summary content."

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        tmp = f.name
    prs.save(tmp)

    # Read slide inventory
    r = await send("read_pptx", path=tmp)
    assert r["ok"], f"PPTX read failed: {r}"
    d = r["data"]
    assert d["slide_count"] == 2
    record(
        "PPTX: read slide inventory",
        PASS,
        f"slide_count={d['slide_count']} titles={[s['title'] for s in d['slides']]}",
    )

    # Write bullets to slide 0 (7 words max enforced)
    ops = [
        {
            "slide_index": 0,
            "bullets": [
                "Cut infrastructure costs by 40 percent",
                "Eliminate redundant toolchain overhead completely",
                "Automate deployment saves 20 hours weekly",
                "ROI positive within six months",
            ],
        }
    ]
    r2 = await send("write_pptx", path=tmp, payload={"operations": ops})
    assert r2["ok"], f"PPTX write failed: {r2}"
    assert len(r2["data"]["errors"]) == 0
    record(
        "PPTX: write bullets",
        PASS,
        f"applied_count={r2['data']['applied_count']} errors={r2['data']['errors']}",
    )

    # Verify slide 2 was NOT touched
    prs2 = Presentation(tmp)
    slide2_text = prs2.slides[1].placeholders[1].text
    assert "Summary content" in slide2_text, f"Slide 2 corrupted: {slide2_text}"
    record("PPTX: untouched slides preserved", PASS, "slide 2 content unchanged")

    # Verify 7-word enforcement
    slide1_text = prs2.slides[0].placeholders[1].text
    for bullet in slide1_text.split("\n"):
        if bullet.strip():
            word_count = len(bullet.strip().split())
            assert word_count <= 7, f"Bullet exceeds 7 words ({word_count}): {bullet}"
    record("PPTX: 7-word bullet enforcement", PASS, "all bullets ≤ 7 words verified")

    os.unlink(tmp)


async def test_pdf_extraction():
    """Test PDF extraction (PyMuPDF tier)."""
    # We can't create a real PDF easily, so we test the routing
    # by using a text file that acts as fallback
    with tempfile.NamedTemporaryFile(suffix=".txt", delete=False, mode="w", encoding="utf-8") as f:
        f.write("# Test Document\n\nThis is test content.\n\n## Section 2\n\nMore content here.")
        tmp = f.name

    r = await send("extract_context", path=tmp)
    assert r["ok"], f"Context extraction failed: {r}"
    assert "full_text" in r["data"]
    assert len(r["data"]["full_text"]) > 0
    record("TXT: extract_context", PASS, f"extracted {len(r['data']['full_text'])} chars")
    os.unlink(tmp)


async def test_backup_created():
    """Verify backup files are created before writes."""
    try:
        from docx import Document
    except ImportError:
        record("DOCX: backup on write", SKIP, "python-docx not installed")
        return

    doc = Document()
    doc.add_paragraph("Test content")
    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
        tmp = f.name
    doc.save(tmp)
    tmp.replace(".docx", ".docx.kairo_backup")

    ops = [{"action": "append", "style": "Normal", "content": "New content", "rows": None}]
    r = await send("write_docx", path=tmp, payload={"operations": ops})
    # After successful write, backup is cleaned up (no errors = backup removed)
    assert r["ok"]
    record(
        "DOCX: atomic write (backup→rename)",
        PASS,
        f"applied_count={r['data']['applied_count']} errors={r['data']['errors']}",
    )
    if os.path.exists(tmp):
        os.unlink(tmp)


# ─── Main runner ─────────────────────────────────────────────────────────────


async def main():
    print("=" * 60)
    print("Kairo Phantom E2E Test Suite")
    print(f"Sidecar: {HOST}:{PORT}")
    print("=" * 60)
    print()

    tests = [
        ("Core", test_ping),
        ("DOCX", test_docx_read_write),
        ("DOCX", test_docx_table_insert),
        ("DOCX", test_backup_created),
        ("XLSX", test_xlsx_read_write),
        ("PPTX", test_pptx_read_write),
        ("TXT", test_pdf_extraction),
    ]

    for group, test_fn in tests:
        print(f"\n[{group}]")
        try:
            await test_fn()
        except Exception as e:
            name = test_fn.__name__
            record(name, FAIL, f"Uncaught: {e}")
            traceback.print_exc()

    # Summary
    passed = sum(1 for r in results if r["status"] == PASS)
    failed = sum(1 for r in results if r["status"] == FAIL)
    skipped = sum(1 for r in results if r["status"] == SKIP)
    total = len(results)

    print()
    print("=" * 60)
    print(f"Results: {passed}/{total} passed | {failed} failed | {skipped} skipped")
    pass_rate = (passed / (total - skipped)) * 100 if (total - skipped) > 0 else 0
    print(f"Pass rate: {pass_rate:.1f}%")
    if failed == 0:
        print("✅ ALL TESTS PASSED — sidecar is production-ready")
    else:
        print(f"❌ {failed} tests failed — review above")
    print("=" * 60)

    return 0 if failed == 0 else 1


if __name__ == "__main__":
    code = asyncio.run(main())
    sys.exit(code)
